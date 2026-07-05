"""Generate the final capstone presentation from the Great Learning template.

Uses CapstoneProject_PresentationTemplate.pptx as the base (inheriting its theme,
master, and layouts), clears the example slides, and rebuilds the deck with our
content + figures. Every image is fit into a bounding box (scaled by width AND
height) so nothing overflows the slide.

Run:  uv run python scripts/generate_presentation.py
"""

from __future__ import annotations

import copy
import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "figures"
TAB = ROOT / "reports" / "tables"
TEMPLATE = ROOT / "CapstoneProject_PresentationTemplate.pptx"
METRICS = json.loads((ROOT / "models" / "metrics.json").read_text())
M = METRICS["metrics"]

GREY = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation(str(TEMPLATE))
SW, SH = prs.slide_width, prs.slide_height

# Layout indices (from the template)
L_TITLE, L_BODY, L_SECTION, L_TITLE_ONLY = 0, 1, 2, 5


def clear_slides() -> None:
    """Remove the template's example slides AND their parts (drop the rels so the
    orphaned slide parts aren't re-serialized, which would collide on save)."""
    sldIdLst = prs.slides._sldIdLst
    rId_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    for sldId in list(sldIdLst):
        prs.part.drop_rel(sldId.get(rId_attr))
        sldIdLst.remove(sldId)


def add(layout_idx: int):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_title(slide, text: str) -> None:
    slide.shapes.title.text = text


def set_body(slide, items, size: int = 14, ph_idx: int = 1) -> None:
    tf = slide.placeholders[ph_idx].text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.level = lvl
        for r in p.runs:
            r.font.size = Pt(max(size - 2 * lvl, 10))


def fit_image(slide, name, left, top, box_w, box_h, center=True):
    """Place an image scaled to fit inside (box_w, box_h), preserving aspect."""
    path = FIG / name
    if not path.exists():
        return None
    pic = slide.shapes.add_picture(str(path), left, top)
    scale = min(box_w / pic.width, box_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = left + (box_w - pic.width) // 2 if center else left
    pic.top = top + (box_h - pic.height) // 2
    return pic


def caption(slide, text, top=None):
    top = top or (SH - Inches(0.45))
    tb = slide.shapes.add_textbox(Inches(0.4), top, SW - Inches(0.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GREY


def image_slide(title, name, cap_text="", box=None):
    s = add(L_TITLE_ONLY)
    set_title(s, title)
    bl, bt, bw, bh = box or (Inches(0.5), Inches(1.05), SW - Inches(1.0), Inches(3.9))
    fit_image(s, name, bl, bt, bw, bh)
    if cap_text:
        caption(s, cap_text)
    return s


def split_slide(title, name, bullets, size=13):
    """Figure on the left, bullets on the right (TITLE_ONLY layout)."""
    s = add(L_TITLE_ONLY)
    set_title(s, title)
    fit_image(s, name, Inches(0.4), Inches(1.1), Inches(5.0), Inches(3.9))
    tb = s.shapes.add_textbox(Inches(5.6), Inches(1.15), Inches(4.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(bullets):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "- " + it if lvl == 0 else "   - " + it
        p.font.size = Pt(size - lvl)
        p.space_after = Pt(6)
    return s


def table_slide(title, df, fontsize=10):
    s = add(L_TITLE_ONLY)
    set_title(s, title)
    rows, cols = df.shape[0] + 1, df.shape[1]
    tbl = s.shapes.add_table(rows, cols, Inches(0.4), Inches(1.2),
                             SW - Inches(0.8), Inches(0.3 * rows)).table
    for j, c in enumerate(df.columns):
        cell = tbl.cell(0, j)
        cell.text = str(c)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(fontsize)
        para.font.bold = True
    for i, (_, row) in enumerate(df.iterrows(), 1):
        for j, c in enumerate(df.columns):
            v = row[c]
            cell = tbl.cell(i, j)
            cell.text = f"{v:,.0f}" if isinstance(v, float) and abs(v) > 100 else (
                f"{v:.3f}" if isinstance(v, float) else str(v))
            cell.text_frame.paragraphs[0].font.size = Pt(fontsize)
    return s


# ===================================================================== #
clear_slides()

# 1 - Title
s = add(L_TITLE)
s.shapes.title.text = "Insurance Price Prediction"
s.placeholders[1].text = ("Capstone Project - Final Presentation\n"
                          "Group << n >>  -  Batch << year >>  -  Mentor: << Name >>")

# 2 - Agenda
s = add(L_BODY); set_title(s, "Contents / Agenda")
set_body(s, [
    "Executive Summary",
    "Business Problem Overview & Solution Approach",
    "Data Overview & Preprocessing",
    "Exploratory Data Analysis (EDA) Outcomes",
    "Model Selection & Metrics of Interest",
    "Model Performance & Final Model",
    "Explainability (Permutation + SHAP)",
    "Business Insights & Recommendations",
    "Deployment",
    "Appendix",
])

# 3 - Executive Summary
s = add(L_BODY); set_title(s, "Executive Summary")
set_body(s, [
    "Problem: estimate an applicant's insurance_cost from health, lifestyle, habit, and "
    "demographic data to support fair pricing and risk triage.",
    f"Final model: LightGBM - Test R2 = {M['test_R2']:.3f}, MAE = {M['test_MAE']:,.0f}, "
    f"RMSE = {M['test_RMSE']:,.0f} (beats the reference solution).",
    "13 model families compared; tuned with RandomizedSearchCV & Optuna; XGBoost/CatBoost on GPU.",
    "Key drivers: weight, admission history, other-insurer coverage, preventive checkups.",
    "Deployed as a leak-free sklearn pipeline via FastAPI + a React UI + a Streamlit app.",
], size=14)

# 4 - Business Problem & Solution
s = add(L_BODY); set_title(s, "Business Problem Overview & Solution Approach")
set_body(s, [
    "Problem: healthcare costs are high; insurers must price policies to balance affordability "
    "against risk. Estimate the optimum insurance cost per individual (target: insurance_cost).",
    "Solution approach:",
    ("Audit & clean the data; fix schema/category quality issues.", 1),
    ("Treat missing values and outliers; engineer risk-oriented features.", 1),
    ("Compare baseline, linear, tree, and boosting models with cross-validation.", 1),
    ("Tune the best family; select on RMSE/MAE/R2 + stability + overfit gap.", 1),
    ("Explain the model and deploy it for inference.", 1),
])

# 5 - Data Overview
s = add(L_BODY); set_title(s, "Data Overview")
set_body(s, [
    "25,000 applicants x 24 columns; one row per applicant; no duplicate rows.",
    "Target insurance_cost: mean ~ 27,147, mild right skew (~ 0.33).",
    "Missing values in only two columns:",
    ("Year_last_admitted - 47.5% (meaningful: indicates no prior admission)", 1),
    ("bmi - 4.0% (imputed by age-band x gender median)", 1),
    "Mixed types: continuous, discrete counts, binary flags, nominal & ordinal categoricals.",
    "Fixed source typos in column names and the 'Salried' -> 'Salaried' category.",
])

# 6-7 - EDA
image_slide("EDA - Distributions of Continuous Attributes", "univariate_continuous.png",
            "Target is mildly right-skewed; weight, BMI, glucose, and steps are roughly symmetric.")
split_slide("EDA - Key Relationships", "boxplots_by_category.png", [
    "weight has an unusually strong positive link to cost - the dominant signal.",
    "Higher cost for applicants covered by another insurer and adventure-sports participants.",
    "Admission history adds useful risk information.",
    "Lifestyle effects act mainly through nonlinear interactions captured by boosting.",
    ("Flagged: validate weight dominance for plausibility/leakage.", 1),
])

# 8 - Preprocessing
s = add(L_BODY); set_title(s, "Data Preprocessing")
set_body(s, [
    "Duplicate check: none; dropped applicant_id (non-predictive identifier).",
    "Missing values: grouped-median BMI imputation; Year_last_admitted -> was_admitted_before + "
    "years_since_last_admitted.",
    "Outliers: capped unrealistic BMI; retained valid high-premium targets.",
    "Feature engineering: age_band, bmi_category, cholesterol_midpoint, any_major_disease_history, "
    "weight_bmi_interaction, steps_per_age.",
    "Leak-free sklearn Pipeline: feature engineering -> impute + scale + one-hot -> model "
    "(fit on training data only).",
])

# 9 - Model Selection & Metrics
s = add(L_BODY); set_title(s, "Model Selection & Metrics of Interest")
set_body(s, [
    "Baseline: DummyRegressor (mean) and LinearRegression.",
    "Advanced: Ridge/Lasso/ElasticNet, DecisionTree/RandomForest/ExtraTrees, and boosting "
    "(GradientBoosting, HistGradientBoosting, XGBoost, LightGBM, CatBoost).",
    "Metrics of interest:",
    ("MAE - average pricing error in cost units (primary, business-readable).", 1),
    ("RMSE - penalises larger errors; R2 - variance explained.", 1),
    ("MAPE/SMAPE - percentage errors; CV RMSE & train-test gap - stability/overfit.", 1),
])

# 10 - Model Performance Summary (table)
comp = pd.read_csv(TAB / "model_comparison.csv")
tbl_df = comp[["model", "train_R2", "test_R2", "test_RMSE", "test_MAE"]].head(7).round(3)
tbl_df.columns = ["Model", "Train R2", "Test R2", "Test RMSE", "Test MAE"]
table_slide("Model Performance Summary (top models by test RMSE)", tbl_df, fontsize=11)

# 11 - Final model & evaluation
split_slide("Final Model & Evaluation - LightGBM", "diagnostics.png", [
    "Selected for lowest test RMSE/MAE, strong R2, stable CV, small train/test gap.",
    f"Test MAE: {M['test_MAE']:,.0f}",
    f"Test RMSE: {M['test_RMSE']:,.0f}",
    f"Test R2: {M['test_R2']:.4f}",
    f"Test MAPE: {M['test_MAPE']:.1f}%",
    "Residuals centred & homoscedastic; predictions track actuals.",
])

# 12 - Explainability
split_slide("Explainability - Permutation + SHAP", "shap_summary.png", [
    "Two independent methods agree on the drivers:",
    ("weight (dominant)", 1),
    ("admission recency / history", 1),
    ("other-insurer coverage", 1),
    ("preventive checkups, weight change", 1),
    "SHAP shows effect direction (higher weight -> higher cost).",
])

# 13 - Business insights (lift)
split_slide("Business Insights - Lift / Gains", "decile_gains.png", [
    "Applicants sorted by predicted cost into deciles.",
    "Top deciles carry lift well above 1.0.",
    "Cumulative-gains curve sits above random.",
    "-> Effective ranking for underwriting triage and",
    ("targeting high-cost segments.", 1),
])

# 14 - Recommendations
s = add(L_BODY); set_title(s, "Business Recommendations")
set_body(s, [
    "Use as pricing decision-support, not automated underwriting (track override rate).",
    "Route extreme predictions (top decile) to manual review (reduce mispriced policies).",
    "Target wellness programmes at high-risk segments (weight, checkups, activity).",
    "Governance: review Gender/Location for fairness; validate weight dominance for leakage.",
    "Monitor drift, segment-level error, and override rates post-deployment.",
])

# 15 - Deployment showcase
s = add(L_BODY); set_title(s, "Deployment")
set_body(s, [
    "Production pipeline saved as one sklearn artifact (models/final_model.pkl).",
    "FastAPI service: /health, /schema, /predict -> cost + Low/Medium/High risk + top drivers.",
    "React + TypeScript + Tailwind UI, built dynamically from the API schema.",
    "Streamlit app for a one-click demo of the same model.",
    "GPU acceleration (XGBoost/CatBoost); reproducible via `uv run python -m insurance.train`.",
])

# 16 - Appendix divider
s = add(L_SECTION); set_title(s, "Appendix")

# 17 - Data background
s = add(L_BODY); set_title(s, "Appendix - Data Background")
set_body(s, [
    "Source: Insurance Data.csv - 25,000 rows x 24 columns; target insurance_cost.",
    "Predictors span profile (age, gender, occupation, location), health (bmi, glucose, "
    "cholesterol, disease history), habits (smoking, alcohol, exercise, steps), and policy fields.",
    "Full data dictionary: reports/tables/data_dictionary.csv; model comparison & importance "
    "tables under reports/tables/.",
    "Limitations: dataset provenance unknown; weight dominance needs plausibility/leakage review; "
    "predictors must be available at quote time; expand fairness testing before production.",
])

# 18 - Permutation importance appendix figure
image_slide("Appendix - Permutation Importance", "feature_importance.png",
            "Permutation importance of input features for the final LightGBM model.")

# 19 - Thank you
s = add(L_SECTION); set_title(s, "Thank You")

out = ROOT / "Insurance_Price_Prediction_Final_Presentation.pptx"
prs.save(out)
print(f"Saved {out.name} ({len(prs.slides._sldIdLst)} slides)")
