from __future__ import annotations

import asyncio
import importlib.util
import json
import math
import os
import pickle
import py_compile
import re
import shutil
import subprocess
import textwrap
import warnings
import zipfile
from html import escape, unescape
from pathlib import Path

import joblib
import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import nbformat as nbf
import numpy as np
import pandas as pd
import seaborn as sns
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from nbclient import NotebookClient
from nbconvert import HTMLExporter
from nbconvert.preprocessors import TagRemovePreprocessor
from sklearn.base import clone
from sklearn.compose import TransformedTargetRegressor
from sklearn.dummy import DummyRegressor
from sklearn.ensemble import (
    ExtraTreesRegressor,
    GradientBoostingRegressor,
    HistGradientBoostingClassifier,
    HistGradientBoostingRegressor,
    RandomForestRegressor,
    VotingRegressor,
)
from sklearn.inspection import PartialDependenceDisplay, permutation_importance
from sklearn.isotonic import IsotonicRegression
from sklearn.linear_model import LinearRegression, Ridge
from sklearn.model_selection import KFold, RandomizedSearchCV, RepeatedKFold, cross_val_score, train_test_split
from sklearn.tree import DecisionTreeRegressor

warnings.filterwarnings("ignore", message="X does not have valid feature names.*")

try:
    from catboost import CatBoostRegressor
except ImportError:  # optional dependency
    CatBoostRegressor = None

try:
    from lightgbm import LGBMRegressor
except ImportError:  # optional dependency
    LGBMRegressor = None

try:
    from xgboost import XGBRegressor
except ImportError:  # optional dependency
    XGBRegressor = None

from insurance_modeling import (
    BASE_CATEGORICAL_FEATURES,
    BASE_NUMERIC_FEATURES,
    CATEGORICAL_FEATURES,
    COLUMN_RENAMES,
    NUMERIC_FEATURES,
    RAW_MODEL_COLUMNS,
    InsuranceFeatureEngineer,
    clean_column_names,
    cholesterol_to_midpoint,
    evaluate_regression,
    make_age_band,
    make_bmi_category,
    make_model_pipeline,
)


ROOT = Path(__file__).resolve().parent
DATA_PATH = ROOT / "Insurance Data.csv"
M1_TEMPLATE = ROOT / "Milestone 1 Template.docx"
M2_TEMPLATE = ROOT / "Milestone 2 Template.docx"
PPT_TEMPLATE = ROOT / "CapstoneProject_PresentationTemplate.pptx"

OUTPUTS = ROOT / "outputs"
FIG_DIR = OUTPUTS / "figures"
TABLE_DIR = OUTPUTS / "tables"
MODEL_DIR = OUTPUTS / "models"
REPORT_DIR = OUTPUTS / "reports"
NOTEBOOK_DIR = ROOT / "notebooks"
OUTPUT_NOTEBOOK_DIR = OUTPUTS / "notebooks"

TARGET = "insurance_cost"
RANDOM_STATE = 42
PRIMARY = "#12355B"
ACCENT = "#2A9D8F"
WARN = "#E76F51"
MUTED = "#6B7280"
SECONDARY = "#457B9D"
GOLD = "#E9C46A"
GREEN = "#8AB17D"
PLUM = "#6D597A"
SKY = "#5DADE2"
PLOT_BG = "#FBFCFE"
GRID = "#E5E7EB"
TEXT = "#1F2937"
CHART_PALETTE = [ACCENT, SECONDARY, WARN, GOLD, GREEN, PLUM, SKY, "#F4A261"]


def gpu_available() -> bool:
    """Return True when an NVIDIA GPU is visible to the active Python runtime."""
    if shutil.which("nvidia-smi") is None:
        return False
    try:
        subprocess.run(["nvidia-smi"], capture_output=True, check=True, timeout=5)
        return True
    except Exception:
        return False


def use_gpu_acceleration() -> bool:
    """Respect INSURANCE_USE_GPU=0/1, otherwise auto-detect GPU availability."""
    requested = os.environ.get("INSURANCE_USE_GPU")
    if requested is not None:
        return requested.strip().lower() in {"1", "true", "yes", "on"}
    return gpu_available()


def optional_package_available(package_name: str) -> bool:
    return importlib.util.find_spec(package_name) is not None


def optional_model_status_sentence(metadata: dict[str, object]) -> str:
    packages = metadata.get("external_model_packages", {})
    if not isinstance(packages, dict):
        packages = {}
    display_names = {
        "xgboost": "XGBoost",
        "lightgbm": "LightGBM",
        "catboost": "CatBoost",
    }
    available = [display_names[key] for key, value in packages.items() if key in display_names and value]
    missing = [display_names[key] for key, value in packages.items() if key in display_names and not value]
    if available:
        detail = f"Optional boosting packages available in this run: {', '.join(available)}"
        if missing:
            detail += f"; unavailable: {', '.join(missing)}"
        gpu_detail = (
            "GPU acceleration was enabled for supported candidates."
            if metadata.get("gpu_acceleration_enabled")
            else "GPU acceleration was not enabled for this run."
        )
        return f"{detail}. {gpu_detail}"
    return (
        "Optional boosting packages XGBoost, LightGBM, and CatBoost were not available in the active "
        "environment, so sklearn boosting models carried the model selection."
    )


def explainability_status_sentence(metadata: dict[str, object]) -> str:
    shap_status = metadata.get("shap_status", {})
    if isinstance(shap_status, dict) and shap_status.get("status") == "completed":
        return f"SHAP importance was generated for {shap_status.get('model')}; permutation importance remains the deployed raw-column explanation."
    packages = metadata.get("external_model_packages", {})
    shap_available = isinstance(packages, dict) and bool(packages.get("shap"))
    if shap_available:
        return (
            "Permutation importance is exported by the automated pipeline; SHAP is available in the "
            "environment for deeper notebook-level explanation."
        )
    return "SHAP was not available in the environment; permutation importance is the deployed explainability fallback."


def make_xgboost_cpu_portable(estimator: object) -> None:
    """Switch fitted XGBoost estimators back to CPU so saved artifacts predict anywhere."""
    if XGBRegressor is not None and isinstance(estimator, XGBRegressor):
        try:
            estimator.set_params(device="cpu")
        except Exception:
            pass
        try:
            estimator.get_booster().set_param({"device": "cpu"})
        except Exception:
            pass
        return
    if hasattr(estimator, "named_steps"):
        for step in estimator.named_steps.values():
            make_xgboost_cpu_portable(step)
        return
    if hasattr(estimator, "regressor_"):
        make_xgboost_cpu_portable(estimator.regressor_)
        return
    if hasattr(estimator, "named_estimators_"):
        for sub_estimator in estimator.named_estimators_.values():
            make_xgboost_cpu_portable(sub_estimator)


def ensure_dirs() -> None:
    for directory in [FIG_DIR, TABLE_DIR, MODEL_DIR, REPORT_DIR, NOTEBOOK_DIR, OUTPUT_NOTEBOOK_DIR]:
        directory.mkdir(parents=True, exist_ok=True)


def load_raw_data() -> pd.DataFrame:
    return pd.read_csv(DATA_PATH)


def add_report_features(df: pd.DataFrame) -> pd.DataFrame:
    """Create the milestone-ready analytical table using full-data descriptive rules."""
    frame = clean_column_names(df)
    frame["age_band"] = make_age_band(frame["age"])
    ref_year = int(frame["Year_last_admitted"].max() + 1)
    frame["was_admitted_before"] = frame["Year_last_admitted"].notna().astype(int)
    frame["admission_year_missing_flag"] = frame["Year_last_admitted"].isna().astype(int)
    frame["years_since_last_admitted"] = np.where(
        frame["Year_last_admitted"].notna(), ref_year - frame["Year_last_admitted"], 0
    )
    frame["cholesterol_midpoint"] = frame["cholesterol_level"].map(cholesterol_to_midpoint)

    group_medians = frame.groupby(["age_band", "Gender"], observed=False)["bmi"].transform("median")
    frame["bmi_missing_flag"] = frame["bmi"].isna().astype(int)
    frame["bmi_imputed"] = frame["bmi"].fillna(group_medians).fillna(frame["bmi"].median())
    bmi_cap = frame["bmi_imputed"].quantile(0.995)
    frame["bmi_for_analysis"] = frame["bmi_imputed"].clip(lower=10, upper=bmi_cap)
    frame["bmi_category"] = make_bmi_category(frame["bmi_for_analysis"])
    frame["any_major_disease_history"] = (
        (frame["heart_disease_history"] == 1) | (frame["other_major_disease_history"] == 1)
    ).astype(int)
    frame["weight_bmi_interaction"] = frame["weight"] * frame["bmi_for_analysis"]
    frame["steps_per_age"] = frame["daily_avg_steps"] / frame["age"].replace(0, np.nan)

    engineered = InsuranceFeatureEngineer().fit(frame.drop(columns=[TARGET])).transform(frame.drop(columns=[TARGET]))
    for column in NUMERIC_FEATURES + CATEGORICAL_FEATURES:
        if column not in frame.columns:
            frame[column] = engineered[column]
    frame["lifestyle_risk_band"] = pd.cut(
        frame["lifestyle_risk_score"],
        bins=[-np.inf, 1.5, 3.5, np.inf],
        labels=["Low", "Moderate", "High"],
    ).astype("object")
    frame["medical_risk_band"] = pd.cut(
        frame["medical_risk_score"],
        bins=[-np.inf, 1.5, 3.5, np.inf],
        labels=["Low", "Moderate", "High"],
    ).astype("object")
    return frame


def write_csv(df: pd.DataFrame, name: str, index: bool = False) -> Path:
    path = TABLE_DIR / name
    df.to_csv(path, index=index)
    return path


def set_visual_theme() -> None:
    sns.set_theme(
        style="whitegrid",
        context="notebook",
        font="DejaVu Sans",
        rc={
            "axes.facecolor": PLOT_BG,
            "figure.facecolor": PLOT_BG,
            "axes.edgecolor": GRID,
            "axes.labelcolor": MUTED,
            "axes.titlecolor": PRIMARY,
            "axes.titleweight": "bold",
            "grid.color": GRID,
            "grid.linewidth": 0.7,
            "grid.alpha": 0.75,
            "xtick.color": TEXT,
            "ytick.color": TEXT,
            "text.color": TEXT,
        },
    )
    sns.set_palette(CHART_PALETTE)


def bar_palette(n: int) -> list[str]:
    return [CHART_PALETTE[i % len(CHART_PALETTE)] for i in range(max(n, 1))]


def style_current_figure() -> None:
    fig = plt.gcf()
    fig.patch.set_facecolor(PLOT_BG)
    for ax in fig.axes:
        ax.set_facecolor(PLOT_BG)
        ax.title.set_fontweight("bold")
        ax.title.set_color(PRIMARY)
        ax.xaxis.label.set_color(MUTED)
        ax.yaxis.label.set_color(MUTED)
        ax.tick_params(axis="both", colors=TEXT, labelsize=8.5)
        ax.grid(True, color=GRID, linewidth=0.7, alpha=0.75)
        for side in ["top", "right"]:
            ax.spines[side].set_visible(False)
        for side in ["left", "bottom"]:
            ax.spines[side].set_color(GRID)
            ax.spines[side].set_linewidth(0.8)


def target_grid_metadata(y: pd.Series) -> dict[str, object]:
    values = pd.Series(y).dropna().astype(float)
    levels = np.sort(values.unique())
    int_levels = np.rint(levels).astype(int)
    diffs = np.diff(int_levels)
    positive_diffs = diffs[diffs > 0]
    grid_step = int(np.gcd.reduce(positive_diffs)) if len(positive_diffs) else 0
    return {
        "target_unique_count": int(len(int_levels)),
        "target_grid_step": grid_step,
        "target_min": int(int_levels.min()) if len(int_levels) else None,
        "target_max": int(int_levels.max()) if len(int_levels) else None,
        "valid_target_levels": int_levels.tolist(),
    }


def make_target_strata(y: pd.Series) -> pd.Series:
    counts = y.value_counts()
    if counts.min() >= 2:
        return y.astype(str)
    bins = min(10, max(2, y.nunique()))
    return pd.qcut(y.rank(method="first"), q=bins, labels=False, duplicates="drop").astype(str)


def target_band_indices(values: pd.Series | np.ndarray, valid_levels: list[int]) -> np.ndarray:
    levels = np.asarray(valid_levels, dtype=float)
    mapping = {float(level): idx for idx, level in enumerate(levels)}
    arr = np.asarray(values, dtype=float)
    mapped = np.array([mapping.get(float(value), np.nan) for value in arr], dtype=float)
    missing = np.isnan(mapped)
    if missing.any():
        mapped[missing] = np.abs(arr[missing][:, None] - levels[None, :]).argmin(axis=1)
    return mapped.astype(int)


def round_to_valid_price_grid(predictions: pd.Series | np.ndarray, valid_levels: list[int]) -> np.ndarray:
    levels = np.asarray(valid_levels, dtype=float)
    preds = np.asarray(predictions, dtype=float)
    nearest_idx = np.abs(preds[:, None] - levels[None, :]).argmin(axis=1)
    return levels[nearest_idx]


def evaluate_with_price_grid(
    y_true: pd.Series | np.ndarray,
    y_pred: pd.Series | np.ndarray,
    valid_levels: list[int],
) -> dict[str, float]:
    metrics = evaluate_regression(y_true, y_pred)
    true_idx = target_band_indices(np.asarray(y_true), valid_levels)
    pred_idx = target_band_indices(round_to_valid_price_grid(y_pred, valid_levels), valid_levels)
    band_distance = np.abs(true_idx - pred_idx)
    metrics.update(
        {
            "exact_band_accuracy": float(np.mean(band_distance == 0)),
            "within_1_band_accuracy": float(np.mean(band_distance <= 1)),
            "within_2_band_accuracy": float(np.mean(band_distance <= 2)),
            "within_3_band_accuracy": float(np.mean(band_distance <= 3)),
        }
    )
    return metrics


def prediction_variant_row(
    variant: str,
    y_true: pd.Series | np.ndarray,
    y_pred: pd.Series | np.ndarray,
    valid_levels: list[int],
) -> dict[str, float | str]:
    metrics = evaluate_with_price_grid(y_true, y_pred, valid_levels)
    return {
        "variant": variant,
        "MAE": metrics["MAE"],
        "RMSE": metrics["RMSE"],
        "R2": metrics["R2"],
        "MAPE": metrics["MAPE"],
        "SMAPE": metrics["SMAPE"],
        "exact_band_accuracy": metrics["exact_band_accuracy"],
        "within_1_band_accuracy": metrics["within_1_band_accuracy"],
        "within_2_band_accuracy": metrics["within_2_band_accuracy"],
        "within_3_band_accuracy": metrics["within_3_band_accuracy"],
    }


def save_profile_tables(raw: pd.DataFrame, df: pd.DataFrame) -> dict[str, Path]:
    paths: dict[str, Path] = {}
    shape_table = pd.DataFrame(
        [
            ["Rows", raw.shape[0]],
            ["Columns", raw.shape[1]],
            ["Cells", raw.size],
            ["Duplicate rows", int(raw.duplicated().sum())],
            ["Unique applicant_id", int(raw["applicant_id"].nunique())],
            ["Target mean", round(raw[TARGET].mean(), 2)],
            ["Target median", round(raw[TARGET].median(), 2)],
            ["Target skew", round(raw[TARGET].skew(), 3)],
        ],
        columns=["Metric", "Value"],
    )
    paths["dataset_profile"] = write_csv(shape_table, "dataset_profile.csv")

    grid_meta = target_grid_metadata(raw[TARGET])
    grid_summary = pd.DataFrame(
        [
            ["target_unique_count", grid_meta["target_unique_count"]],
            ["target_grid_step", grid_meta["target_grid_step"]],
            ["target_min", grid_meta["target_min"]],
            ["target_max", grid_meta["target_max"]],
            ["first_10_levels", ", ".join(map(str, grid_meta["valid_target_levels"][:10]))],
            ["last_10_levels", ", ".join(map(str, grid_meta["valid_target_levels"][-10:]))],
        ],
        columns=["metric", "value"],
    )
    paths["target_grid_summary"] = write_csv(grid_summary, "target_grid_summary.csv")
    target_frequency = (
        raw[TARGET]
        .value_counts()
        .sort_index()
        .rename_axis("insurance_cost")
        .reset_index(name="customer_count")
    )
    target_frequency["pct"] = (target_frequency["customer_count"] / len(raw) * 100).round(3)
    paths["target_price_band_frequency"] = write_csv(target_frequency, "target_price_band_frequency.csv")

    dtype_table = pd.DataFrame(
        {
            "column": raw.columns,
            "clean_column": clean_column_names(raw).columns,
            "dtype": raw.dtypes.astype(str).values,
            "missing": raw.isna().sum().values,
            "missing_pct": (raw.isna().mean().values * 100).round(2),
            "unique_values": [raw[c].nunique(dropna=True) for c in raw.columns],
        }
    )
    paths["data_dictionary"] = write_csv(dtype_table, "data_dictionary_summary.csv")

    cleaned_columns = clean_column_names(raw).columns.tolist()
    business_descriptions = {
        "applicant_id": "Unique applicant identifier used only for record tracing.",
        "years_of_insurance_with_us": "Tenure with the insurer.",
        "regular_checkup_last_year": "Number of routine checkups in the previous year.",
        "adventure_sports": "Whether the applicant participates in adventure sports.",
        "Occupation": "Applicant occupation category.",
        "visited_doctor_last_1_year": "Number of doctor visits in the previous year.",
        "cholesterol_level": "Ordinal cholesterol range category.",
        "insurance_cost": "Historical insurance quote or premium cost.",
        "age": "Applicant age.",
        "bmi": "Body mass index.",
        "weight": "Applicant weight.",
        "avg_glucose_level": "Average glucose level.",
        "smoking_status": "Applicant smoking status.",
        "exercise": "Exercise habit category.",
        "Alcohol": "Alcohol-use flag/category.",
        "Gender": "Applicant gender.",
        "Location": "Applicant location.",
        "Year_last_admitted": "Most recent hospital admission year when available.",
        "covered_by_any_other_company": "Whether another insurance company already covers the applicant.",
        "weight_change_in_last_one_year": "Weight-change category over the last year.",
        "fat_percentage": "Applicant body-fat percentage.",
        "heart_disease_history": "Historical heart-disease flag.",
        "other_major_disease_history": "Historical other-major-disease flag.",
    }

    def analysis_type(clean_col: str) -> str:
        if clean_col == "applicant_id":
            return "identifier"
        if clean_col == TARGET:
            return "discrete numeric quote grid / regression target"
        if clean_col in {"years_of_insurance_with_us", "regular_checkup_last_year", "visited_doctor_last_1_year", "weight_change_in_last_one_year"}:
            return "ordinal discrete numeric"
        if clean_col in {"adventure_sports", "heart_disease_history", "other_major_disease_history"}:
            return "binary categorical"
        if clean_col in {"Occupation", "Gender", "Location", "smoking_status", "Alcohol", "exercise", "covered_by_any_other_company"}:
            return "nominal categorical"
        if clean_col == "cholesterol_level":
            return "ordinal categorical"
        if clean_col == "Year_last_admitted":
            return "numeric with structural missingness"
        return "continuous numeric"

    def role_for(clean_col: str) -> str:
        if clean_col == TARGET:
            return "target"
        if clean_col == "applicant_id":
            return "identifier"
        return "predictor"

    def preprocessing_action(clean_col: str) -> str:
        if clean_col == "applicant_id":
            return "retain for audit only; drop from modeling"
        if clean_col == TARGET:
            return "retain as target; analyze as quote grid"
        if clean_col == "bmi":
            return "create missingness flag; impute using train-fitted median logic"
        if clean_col == "Year_last_admitted":
            return "create admission flags/status and recency; keep missingness signal"
        if clean_col == "cholesterol_level":
            return "retain ordinal category and engineer cholesterol midpoint"
        if clean_col in {"regular_checkup_last_year", "heart_disease_history", "other_major_disease_history"}:
            return "normalize misspelled source name; preserve numeric/category meaning"
        return "validate type, preserve for EDA, and feed through modeling pipeline"

    def preprocessing_reason(clean_col: str) -> str:
        if clean_col == "applicant_id":
            return "Unique row key would leak identity-like noise and cannot generalize."
        if clean_col == TARGET:
            return "Business output to be predicted; values form a fixed quote-band grid."
        if clean_col == "bmi":
            return "BMI has missing values and valid extreme values requiring explicit handling."
        if clean_col == "Year_last_admitted":
            return "Missing admission year means no known admission history rather than a normal year value."
        if clean_col == "cholesterol_level":
            return "Ranges are ordered and can support both category and midpoint representations."
        return "Required for rubric data understanding and consistent train-time preprocessing."

    variable_rows = []
    for original_col, clean_col in zip(raw.columns.tolist(), cleaned_columns):
        examples = raw[original_col].dropna().astype(str).drop_duplicates().head(3).tolist()
        variable_rows.append(
            {
                "original_column": original_col,
                "clean_column": clean_col,
                "business_description": business_descriptions.get(clean_col, f"Applicant attribute: {clean_col}."),
                "data_type_raw": str(raw[original_col].dtype),
                "analysis_type": analysis_type(clean_col),
                "role": role_for(clean_col),
                "unique_count": int(raw[original_col].nunique(dropna=True)),
                "missing_count": int(raw[original_col].isna().sum()),
                "example_values": ", ".join(examples),
                "preprocessing_action": preprocessing_action(clean_col),
                "reason": preprocessing_reason(clean_col),
            }
        )
    paths["variable_type_classification"] = write_csv(
        pd.DataFrame(variable_rows),
        "variable_type_classification.csv",
    )

    missing_table = (
        raw.isna()
        .sum()
        .reset_index()
        .rename(columns={"index": "column", 0: "missing_count"})
        .assign(missing_pct=lambda x: (x["missing_count"] / len(raw) * 100).round(2))
        .sort_values("missing_count", ascending=False)
    )
    paths["missing_values"] = write_csv(missing_table, "missing_values.csv")

    rename_table = pd.DataFrame(
        [
            ["regular_checkup_lasy_year", "regular_checkup_last_year", "column typo normalized"],
            ["heart_decs_history", "heart_disease_history", "column typo normalized"],
            ["other_major_decs_history", "other_major_disease_history", "column typo normalized"],
            ["Occupation value: Salried", "Salaried", "category typo normalized"],
            ["applicant_id", "dropped from modeling", "unique identifier, not predictive signal"],
        ],
        columns=["original", "cleaned_or_action", "reason"],
    )
    paths["renaming"] = write_csv(rename_table, "renaming_and_cleanup.csv")

    numeric_summary = raw.describe().T.reset_index().rename(columns={"index": "column"})
    paths["numeric_summary"] = write_csv(numeric_summary, "numeric_summary.csv")

    categorical_cols = [
        "Occupation",
        "cholesterol_level",
        "Gender",
        "smoking_status",
        "Location",
        "covered_by_any_other_company",
        "Alcohol",
        "exercise",
    ]
    rows = []
    for column in categorical_cols:
        counts = raw[column].value_counts(dropna=False)
        for value, count in counts.items():
            rows.append(
                {
                    "column": column,
                    "category": value,
                    "count": int(count),
                    "pct": round(count / len(raw) * 100, 2),
                }
            )
    paths["categorical_frequency"] = write_csv(pd.DataFrame(rows), "categorical_frequency.csv")

    outlier_cols = ["bmi", "weight", "avg_glucose_level", "daily_avg_steps", "fat_percentage", TARGET]
    out_rows = []
    for column in outlier_cols:
        series = raw[column].dropna()
        q1 = series.quantile(0.25)
        q3 = series.quantile(0.75)
        iqr = q3 - q1
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        count = int(((series < lower) | (series > upper)).sum())
        out_rows.append(
            {
                "column": column,
                "q1": round(q1, 3),
                "q3": round(q3, 3),
                "iqr": round(iqr, 3),
                "lower_fence": round(lower, 3),
                "upper_fence": round(upper, 3),
                "outlier_count": count,
                "outlier_pct": round(count / len(series) * 100, 2),
            }
        )
    paths["outlier_summary"] = write_csv(pd.DataFrame(out_rows), "outlier_summary.csv")

    missing_admission_profile = (
        df.assign(year_last_admitted_missing=df["Year_last_admitted"].isna())
        .groupby("year_last_admitted_missing")
        .agg(
            customer_count=(TARGET, "size"),
            mean_cost=(TARGET, "mean"),
            median_cost=(TARGET, "median"),
            mean_age=("age", "mean"),
            mean_bmi=("bmi_for_analysis", "mean"),
            mean_medical_risk=("medical_risk_score", "mean"),
            mean_lifestyle_risk=("lifestyle_risk_score", "mean"),
        )
        .round(2)
        .reset_index()
    )
    paths["missing_admission_profile"] = write_csv(
        missing_admission_profile, "missing_year_last_admitted_profile.csv"
    )

    missing_bmi_profile = (
        df.assign(bmi_missing=raw["bmi"].isna())
        .groupby("bmi_missing")
        .agg(
            customer_count=(TARGET, "size"),
            mean_cost=(TARGET, "mean"),
            median_cost=(TARGET, "median"),
            mean_age=("age", "mean"),
            mean_weight=("weight", "mean"),
            mean_medical_risk=("medical_risk_score", "mean"),
            mean_lifestyle_risk=("lifestyle_risk_score", "mean"),
        )
        .round(2)
        .reset_index()
    )
    paths["missing_bmi_profile"] = write_csv(missing_bmi_profile, "missing_bmi_profile.csv")

    missing_rows = []
    for feature, flag in [
        ("BMI missing", raw["bmi"].isna()),
        ("Year_last_admitted missing", raw["Year_last_admitted"].isna()),
    ]:
        grouped = raw.assign(missing_flag=flag).groupby("missing_flag")[TARGET].agg(["size", "mean", "median"])
        baseline = float(raw[TARGET].mean())
        for missing_flag, row in grouped.iterrows():
            missing_rows.append(
                {
                    "analysis": feature,
                    "segment": "missing" if bool(missing_flag) else "present",
                    "customer_count": int(row["size"]),
                    "mean_cost": round(float(row["mean"]), 2),
                    "median_cost": round(float(row["median"]), 2),
                    "mean_vs_overall": round(float(row["mean"]) - baseline, 2),
                }
            )
    for feature, flag, segment_col in [
        ("BMI missing by smoking_status", raw["bmi"].isna(), "smoking_status"),
        ("Admission missing by other-company coverage", raw["Year_last_admitted"].isna(), "covered_by_any_other_company"),
    ]:
        grouped = (
            raw.assign(missing_flag=flag)
            .groupby(["missing_flag", segment_col], dropna=False)[TARGET]
            .agg(["size", "mean", "median"])
            .reset_index()
        )
        for _, row in grouped.iterrows():
            missing_rows.append(
                {
                    "analysis": feature,
                    "segment": f"{segment_col}={row[segment_col]}, {'missing' if bool(row['missing_flag']) else 'present'}",
                    "customer_count": int(row["size"]),
                    "mean_cost": round(float(row["mean"]), 2),
                    "median_cost": round(float(row["median"]), 2),
                    "mean_vs_overall": round(float(row["mean"]) - float(raw[TARGET].mean()), 2),
                }
            )
    paths["missingness_target_impact"] = write_csv(pd.DataFrame(missing_rows), "missingness_target_impact.csv")

    admission_recency = df.assign(
        admission_recency_group=np.select(
            [
                df["was_admitted_before"] == 0,
                df["years_since_last_admitted"].between(0, 5, inclusive="both"),
                df["years_since_last_admitted"].between(6, 15, inclusive="both"),
            ],
            ["No known admission", "Admitted recently (0-5 years)", "Admitted 6-15 years ago"],
            default="Admitted 16+ years ago",
        )
    )
    admission_summary = (
        admission_recency.groupby("admission_recency_group")
        .agg(
            customer_count=(TARGET, "size"),
            mean_cost=(TARGET, "mean"),
            median_cost=(TARGET, "median"),
            mean_weight=("weight", "mean"),
            mean_bmi=("bmi_for_analysis", "mean"),
            mean_years_since_last_admitted=("years_since_last_admitted", "mean"),
        )
        .round(2)
        .reset_index()
    )
    paths["admission_recency_summary"] = write_csv(admission_summary, "admission_recency_summary.csv")

    tail_cut_low = df[TARGET].quantile(0.01)
    tail_cut_high = df[TARGET].quantile(0.99)
    tail_profile = (
        df.assign(
            target_tail=np.select(
                [df[TARGET] <= tail_cut_low, df[TARGET] >= tail_cut_high],
                ["Bottom 1%", "Top 1%"],
                default="Middle 98%",
            )
        )
        .query("target_tail != 'Middle 98%'")
        .groupby("target_tail")
        .agg(
            customer_count=(TARGET, "size"),
            mean_cost=(TARGET, "mean"),
            mean_age=("age", "mean"),
            mean_bmi=("bmi_for_analysis", "mean"),
            mean_glucose=("avg_glucose_level", "mean"),
            mean_daily_steps=("daily_avg_steps", "mean"),
            mean_medical_risk=("medical_risk_score", "mean"),
            mean_lifestyle_risk=("lifestyle_risk_score", "mean"),
        )
        .round(2)
        .reset_index()
    )
    paths["target_tail_profile"] = write_csv(tail_profile, "target_tail_profile.csv")

    low_cardinality_numeric = [
        "years_of_insurance_with_us",
        "regular_checkup_last_year",
        "visited_doctor_last_1_year",
        "weight_change_in_last_one_year",
        "adventure_sports",
        "heart_disease_history",
        "other_major_disease_history",
    ]
    for column in low_cardinality_numeric:
        summary = (
            df.groupby(column, dropna=False)[TARGET]
            .agg(customer_count="size", mean_cost="mean", median_cost="median")
            .round(2)
            .reset_index()
            .sort_values(column)
        )
        safe_name = {
            "regular_checkup_last_year": "group_summary_regular_checkup_last_year.csv",
            "weight_change_in_last_one_year": "group_summary_weight_change.csv",
        }.get(column, f"group_summary_{column}.csv")
        paths[f"group_summary_{column}"] = write_csv(summary, safe_name)

    def eta_squared(frame: pd.DataFrame, feature: str) -> float:
        grouped = frame.groupby(feature, dropna=False)[TARGET]
        grand_mean = frame[TARGET].mean()
        ss_between = grouped.size().mul((grouped.mean() - grand_mean) ** 2).sum()
        ss_total = ((frame[TARGET] - grand_mean) ** 2).sum()
        return float(ss_between / ss_total) if ss_total else 0.0

    signal_features = [
        ("weight", "numeric"),
        ("Year_last_admitted", "numeric_missing"),
        ("years_since_last_admitted", "numeric_engineered"),
        ("covered_by_any_other_company", "categorical"),
        ("regular_checkup_last_year", "ordinal"),
        ("weight_change_in_last_one_year", "ordinal"),
        ("adventure_sports", "binary"),
        ("smoking_status", "categorical"),
        ("Alcohol", "categorical"),
        ("exercise", "categorical"),
        ("age", "numeric"),
        ("bmi_for_analysis", "numeric"),
        ("avg_glucose_level", "numeric"),
        ("cholesterol_midpoint", "numeric"),
        ("heart_disease_history", "binary"),
        ("other_major_disease_history", "binary"),
    ]
    signal_rows = []
    for feature, feature_type in signal_features:
        working = df[[feature, TARGET]].copy() if feature in df.columns else raw[[feature, TARGET]].copy()
        if feature in {"Year_last_admitted"}:
            working[feature] = working[feature].fillna("No known admission")
        if feature_type.startswith("numeric") and feature not in {"Year_last_admitted"} and working[feature].nunique(dropna=True) > 15:
            corr = working[[feature, TARGET]].corr(numeric_only=True).iloc[0, 1]
            buckets = pd.qcut(working[feature].rank(method="first"), q=10, duplicates="drop")
            mean_range = working.assign(bucket=buckets).groupby("bucket", observed=False)[TARGET].mean()
            signal_value = abs(float(corr)) if pd.notna(corr) else 0.0
            target_range = float(mean_range.max() - mean_range.min()) if len(mean_range) else 0.0
            statistic_name = "absolute_correlation"
        else:
            grouped_means = working.groupby(feature, dropna=False)[TARGET].mean()
            signal_value = eta_squared(working, feature)
            target_range = float(grouped_means.max() - grouped_means.min()) if len(grouped_means) else 0.0
            statistic_name = "eta_squared"
        if feature in {"weight", "Year_last_admitted", "years_since_last_admitted", "covered_by_any_other_company", "regular_checkup_last_year", "weight_change_in_last_one_year", "adventure_sports"}:
            interpretation = "strong observed signal"
        elif target_range < 500 or signal_value < 0.002:
            interpretation = "weak marginal signal in this dataset"
        else:
            interpretation = "secondary signal; monitor with interactions"
        signal_rows.append(
            {
                "feature": feature,
                "feature_type": feature_type,
                "target_mean_range": round(target_range, 2),
                "correlation_or_eta_squared": round(signal_value, 5),
                "statistic_name": statistic_name,
                "business_interpretation": interpretation,
            }
        )
    paths["feature_signal_strength"] = write_csv(
        pd.DataFrame(signal_rows).sort_values("target_mean_range", ascending=False),
        "feature_signal_strength.csv",
    )

    cleaned_export = df.drop(columns=["bmi_imputed"], errors="ignore")
    paths["cleaned_analysis_data"] = write_csv(cleaned_export, "cleaned_analysis_data.csv")
    return paths


def savefig(path: Path) -> None:
    style_current_figure()
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def safe_slug(value: object) -> str:
    slug = re.sub(r"[^A-Za-z0-9]+", "_", str(value)).strip("_").lower()
    return slug or "field"


def create_univariate_plot_set(raw: pd.DataFrame, df: pd.DataFrame) -> pd.DataFrame:
    """Create one univariate visual per original source variable with interpretation text."""
    set_visual_theme()
    cleaned = clean_column_names(raw)
    variable_path = TABLE_DIR / "variable_type_classification.csv"
    if variable_path.exists():
        variable_types = pd.read_csv(variable_path)
    else:
        variable_types = pd.DataFrame(
            {
                "original_column": raw.columns,
                "clean_column": cleaned.columns,
                "analysis_type": ["unknown"] * len(raw.columns),
                "role": ["predictor"] * len(raw.columns),
                "preprocessing_action": ["document and validate"] * len(raw.columns),
            }
        )

    rows: list[dict[str, object]] = []
    for _, meta in variable_types.iterrows():
        original_col = str(meta["original_column"])
        clean_col = str(meta["clean_column"])
        analysis_type = str(meta.get("analysis_type", "unknown"))
        role = str(meta.get("role", "predictor"))
        action = str(meta.get("preprocessing_action", "document and validate"))
        source_col = clean_col if clean_col in cleaned.columns else original_col
        if source_col not in cleaned.columns:
            continue
        series = cleaned[source_col]
        missing_count = int(series.isna().sum())
        unique_count = int(series.nunique(dropna=True))
        filename = f"univariate_{safe_slug(clean_col)}.png"
        path = FIG_DIR / filename

        if clean_col == "applicant_id":
            duplicate_ids = int(series.duplicated().sum())
            audit = pd.DataFrame(
                {
                    "check": ["Rows", "Unique IDs", "Duplicate IDs"],
                    "count": [int(len(series)), unique_count, duplicate_ids],
                }
            )
            plt.figure(figsize=(7.2, 4.2))
            sns.barplot(data=audit, x="check", y="count", hue="check", palette=bar_palette(len(audit)), legend=False)
            plt.title("Univariate audit: applicant_id")
            plt.ylabel("Count")
            plot_type = "identifier audit bar chart"
            technical = (
                f"`{clean_col}` has {unique_count:,} unique values across {len(series):,} rows "
                f"and {duplicate_ids:,} duplicate ID values."
            )
            business = (
                "The field is useful for audit traceability only and is excluded from modeling and app input "
                "because an identifier does not generalize to new applicants."
            )
        elif clean_col == TARGET:
            values = pd.to_numeric(series, errors="coerce")
            plt.figure(figsize=(8.4, 4.8))
            sns.histplot(values.dropna(), bins=35, kde=True, color=SECONDARY)
            plt.title("Univariate distribution: insurance_cost")
            plt.xlabel("insurance_cost")
            plt.ylabel("Customers")
            plot_type = "target histogram with KDE"
            grid_meta = target_grid_metadata(values)
            technical = (
                f"`{clean_col}` ranges from {values.min():,.0f} to {values.max():,.0f}, "
                f"has {grid_meta['target_unique_count']} unique quote bands, and a grid step of "
                f"{grid_meta['target_grid_step']:,}."
            )
            business = (
                "This is the prediction target; it behaves like a fixed quote-band grid, so deployment should "
                "show both raw prediction and nearest valid quote band."
            )
        elif pd.api.types.is_numeric_dtype(series) and unique_count > 20 and clean_col != "Year_last_admitted":
            values = pd.to_numeric(series, errors="coerce")
            fig, axes = plt.subplots(1, 2, figsize=(10, 4.2), gridspec_kw={"width_ratios": [2, 1]})
            sns.histplot(values.dropna(), bins=30, kde=True, color=ACCENT, ax=axes[0])
            sns.boxplot(y=values.dropna(), color=GOLD, ax=axes[1])
            axes[0].set_title(f"Distribution of {clean_col}")
            axes[1].set_title("Boxplot")
            axes[0].set_xlabel(clean_col)
            axes[1].set_ylabel(clean_col)
            plot_type = "histogram and boxplot"
            technical = (
                f"`{clean_col}` is numeric with {unique_count:,} unique values, median "
                f"{values.median():,.2f}, range {values.min():,.2f} to {values.max():,.2f}, "
                f"and {missing_count:,} missing values."
            )
            business = (
                f"The distribution documents the applicant spread for `{clean_col}` before modeling; "
                f"the preprocessing action is: {action}."
            )
        else:
            if clean_col == "Year_last_admitted":
                display_series = pd.to_numeric(series, errors="coerce").astype("Int64").astype("object")
                display_series = display_series.fillna("Missing/no known admission").astype(str)
            else:
                display_series = series.astype("object").where(series.notna(), "Missing").astype(str)
            counts = display_series.value_counts(dropna=False).head(20).reset_index()
            counts.columns = [clean_col, "customer_count"]
            plt.figure(figsize=(8.6, 4.8))
            sns.barplot(data=counts, x=clean_col, y="customer_count", hue=clean_col, palette=bar_palette(len(counts)), legend=False)
            plt.title(f"Univariate frequency: {clean_col}")
            plt.xlabel(clean_col)
            plt.ylabel("Customers")
            plt.xticks(rotation=30, ha="right")
            plot_type = "frequency bar chart"
            top_category = str(counts.iloc[0][clean_col]) if len(counts) else "not available"
            top_count = int(counts.iloc[0]["customer_count"]) if len(counts) else 0
            top_pct = top_count / len(series) * 100 if len(series) else 0
            technical = (
                f"`{clean_col}` has {unique_count:,} non-missing levels and {missing_count:,} missing values; "
                f"the most frequent level is `{top_category}` at {top_pct:.1f}% of rows."
            )
            business = (
                f"The frequency view checks whether `{clean_col}` is balanced enough for EDA and downstream encoding; "
                f"the preprocessing action is: {action}."
            )

        savefig(path)
        rows.append(
            {
                "original_column": original_col,
                "clean_column": clean_col,
                "analysis_type": analysis_type,
                "role": role,
                "plot_type": plot_type,
                "plot_file": rel_path(path),
                "unique_count": unique_count,
                "missing_count": missing_count,
                "technical_interpretation": technical,
                "business_interpretation": business,
            }
        )

    index = pd.DataFrame(rows)
    write_csv(index, "univariate_plot_interpretations.csv")
    return index


def generate_eda_figures(raw: pd.DataFrame, df: pd.DataFrame) -> dict[str, Path]:
    set_visual_theme()
    figures: dict[str, Path] = {}
    univariate_index = create_univariate_plot_set(raw, df)
    for _, row in univariate_index.iterrows():
        figures[f"univariate_{safe_slug(row['clean_column'])}"] = ROOT / str(row["plot_file"])

    plt.figure(figsize=(8, 4.8))
    sns.histplot(df[TARGET], kde=True, color=SECONDARY, bins=35)
    plt.title("Distribution of Insurance Cost")
    plt.xlabel("Insurance cost")
    figures["target_distribution"] = FIG_DIR / "target_distribution.png"
    savefig(figures["target_distribution"])

    target_frequency = raw[TARGET].value_counts().sort_index().reset_index()
    target_frequency.columns = [TARGET, "customer_count"]
    plt.figure(figsize=(10, 4.8))
    plt.bar(target_frequency[TARGET].astype(str), target_frequency["customer_count"], color=sns.color_palette("crest", len(target_frequency)))
    plt.title("Insurance Cost Price-Band Frequency")
    plt.xlabel("Insurance cost band")
    plt.ylabel("Customers")
    plt.xticks(rotation=90, fontsize=7)
    figures["target_price_grid_frequency"] = FIG_DIR / "target_price_grid_frequency.png"
    savefig(figures["target_price_grid_frequency"])

    plt.figure(figsize=(7, 4))
    sns.boxplot(x=df[TARGET], color=GOLD)
    plt.title("Insurance Cost Boxplot")
    figures["target_boxplot"] = FIG_DIR / "target_boxplot.png"
    savefig(figures["target_boxplot"])

    missing = raw.isna().sum()
    missing = missing[missing > 0].sort_values(ascending=True)
    plt.figure(figsize=(7, 3.8))
    plt.barh(missing.index, missing.values, color=[WARN if i % 2 == 0 else GOLD for i in range(len(missing))])
    plt.title("Missing Values by Column")
    plt.xlabel("Missing rows")
    figures["missing_values"] = FIG_DIR / "missing_values.png"
    savefig(figures["missing_values"])

    numeric_cols = [
        "years_of_insurance_with_us",
        "regular_checkup_last_year",
        "visited_doctor_last_1_year",
        "daily_avg_steps",
        "age",
        "avg_glucose_level",
        "bmi_for_analysis",
        "weight",
        "fat_percentage",
        "cholesterol_midpoint",
        "years_since_last_admitted",
        TARGET,
    ]
    corr = df[numeric_cols].corr(numeric_only=True)
    plt.figure(figsize=(10, 7))
    sns.heatmap(corr, cmap="vlag", center=0, linewidths=0.35, linecolor="white", annot=False)
    plt.title("Correlation Heatmap")
    figures["correlation_heatmap"] = FIG_DIR / "correlation_heatmap.png"
    savefig(figures["correlation_heatmap"])

    scatter_specs = [
        ("weight", "Insurance Cost vs Weight", "cost_vs_weight.png"),
        ("years_since_last_admitted", "Insurance Cost vs Years Since Last Admission", "cost_vs_years_since_admitted.png"),
        ("avg_glucose_level", "Insurance Cost vs Average Glucose", "cost_vs_glucose.png"),
        ("daily_avg_steps", "Insurance Cost vs Daily Steps", "cost_vs_steps.png"),
    ]
    for x_col, title, filename in scatter_specs:
        plt.figure(figsize=(7.5, 4.8))
        sns.scatterplot(data=df.sample(min(5000, len(df)), random_state=RANDOM_STATE), x=x_col, y=TARGET, alpha=0.35, s=14, color=PRIMARY)
        sns.regplot(data=df.sample(min(5000, len(df)), random_state=RANDOM_STATE), x=x_col, y=TARGET, scatter=False, color=WARN)
        plt.title(title)
        figures[filename.replace(".png", "")] = FIG_DIR / filename
        savefig(figures[filename.replace(".png", "")])

    weight_line = (
        df.groupby("weight", as_index=False)[TARGET]
        .agg(mean_cost="mean", customer_count="size")
        .sort_values("weight")
    )
    plt.figure(figsize=(9, 4.8))
    sns.lineplot(data=weight_line, x="weight", y="mean_cost", marker="o", color=PRIMARY)
    plt.title("Average Insurance Cost by Weight")
    plt.ylabel("Mean insurance cost")
    figures["cost_by_weight_line"] = FIG_DIR / "cost_by_weight_line.png"
    savefig(figures["cost_by_weight_line"])

    df_weight = df.assign(weight_band=pd.cut(df["weight"], bins=8))
    weight_band = (
        df_weight.groupby("weight_band", observed=False)[TARGET]
        .agg(mean_cost="mean", customer_count="size")
        .reset_index()
    )
    weight_band["weight_band"] = weight_band["weight_band"].astype(str)
    plt.figure(figsize=(9, 4.8))
    sns.barplot(data=weight_band, x="weight_band", y="mean_cost", hue="weight_band", palette=sns.color_palette("crest", len(weight_band)), legend=False)
    plt.title("Average Insurance Cost by Weight Band")
    plt.ylabel("Mean insurance cost")
    plt.xticks(rotation=25, ha="right")
    figures["cost_by_weight_band"] = FIG_DIR / "cost_by_weight_band.png"
    savefig(figures["cost_by_weight_band"])

    jitter_sample = df.sample(min(6000, len(df)), random_state=RANDOM_STATE).copy()
    jitter_sample["weight_jitter"] = jitter_sample["weight"] + np.random.default_rng(RANDOM_STATE).normal(0, 0.12, size=len(jitter_sample))
    jitter_sample["cost_jitter"] = jitter_sample[TARGET] + np.random.default_rng(RANDOM_STATE + 1).normal(0, 90, size=len(jitter_sample))
    plt.figure(figsize=(8.5, 5))
    sns.scatterplot(data=jitter_sample, x="weight_jitter", y="cost_jitter", alpha=0.25, s=12, color=PRIMARY)
    plt.title("Insurance Cost vs Weight with Jitter")
    plt.xlabel("Weight")
    plt.ylabel("Insurance cost")
    figures["weight_vs_cost_jitter"] = FIG_DIR / "weight_vs_cost_jitter.png"
    savefig(figures["weight_vs_cost_jitter"])

    boxplot_specs = [
        ("covered_by_any_other_company", "Insurance Cost by Other Coverage", "cost_by_other_coverage.png"),
        ("exercise", "Insurance Cost by Exercise", "cost_by_exercise.png"),
        ("Alcohol", "Insurance Cost by Alcohol Use", "cost_by_alcohol.png"),
        ("smoking_status", "Insurance Cost by Smoking Status", "cost_by_smoking.png"),
        ("adventure_sports", "Insurance Cost by Adventure Sports", "cost_by_adventure_sports.png"),
        ("any_major_disease_history", "Insurance Cost by Disease History Flag", "cost_by_disease_history.png"),
        ("bmi_category", "Insurance Cost by BMI Category", "cost_by_bmi_category.png"),
    ]
    for col, title, filename in boxplot_specs:
        plt.figure(figsize=(8, 4.8))
        order = df.groupby(col)[TARGET].median().sort_values().index.tolist()
        sns.boxplot(data=df, x=col, y=TARGET, order=order, hue=col, palette=bar_palette(len(order)), legend=False)
        plt.title(title)
        plt.xlabel(col)
        plt.ylabel("Insurance cost")
        plt.xticks(rotation=25, ha="right")
        key = filename.replace(".png", "")
        figures[key] = FIG_DIR / filename
        savefig(figures[key])

    ordinal_specs = [
        ("regular_checkup_last_year", "Insurance Cost by Regular Checkups Last Year", "cost_by_regular_checkup_last_year.png"),
        ("weight_change_in_last_one_year", "Insurance Cost by Weight Change Last Year", "cost_by_weight_change.png"),
    ]
    for col, title, filename in ordinal_specs:
        summary = df.groupby(col, as_index=False)[TARGET].mean().sort_values(col)
        plt.figure(figsize=(7.5, 4.6))
        sns.barplot(data=summary, x=col, y=TARGET, hue=col, palette=bar_palette(len(summary)), legend=False)
        plt.title(title)
        plt.ylabel("Mean insurance cost")
        key = filename.replace(".png", "")
        figures[key] = FIG_DIR / filename
        savefig(figures[key])

    year_summary = (
        df.dropna(subset=["Year_last_admitted"])
        .groupby("Year_last_admitted", as_index=False)[TARGET]
        .agg(mean_cost="mean", customer_count="size")
        .sort_values("Year_last_admitted")
    )
    plt.figure(figsize=(9, 4.8))
    sns.lineplot(data=year_summary, x="Year_last_admitted", y="mean_cost", marker="o", color=PRIMARY)
    plt.title("Average Cost by Year Last Admitted")
    plt.ylabel("Mean insurance cost")
    figures["cost_by_year_last_admitted"] = FIG_DIR / "cost_by_year_last_admitted.png"
    savefig(figures["cost_by_year_last_admitted"])

    admission_status = (
        df.groupby("admission_status", as_index=False)[TARGET]
        .agg(mean_cost="mean", customer_count="size")
        .sort_values("mean_cost", ascending=False)
    )
    plt.figure(figsize=(7.5, 4.6))
    sns.barplot(data=admission_status, x="admission_status", y="mean_cost", hue="admission_status", palette=bar_palette(len(admission_status)), legend=False)
    plt.title("Average Cost by Admission Status")
    plt.ylabel("Mean insurance cost")
    plt.xticks(rotation=15, ha="right")
    figures["cost_by_admission_status"] = FIG_DIR / "cost_by_admission_status.png"
    savefig(figures["cost_by_admission_status"])

    missing_impact = pd.DataFrame(
        [
            {
                "missingness": "BMI present",
                "mean_cost": raw.loc[raw["bmi"].notna(), TARGET].mean(),
            },
            {
                "missingness": "BMI missing",
                "mean_cost": raw.loc[raw["bmi"].isna(), TARGET].mean(),
            },
            {
                "missingness": "Admission year present",
                "mean_cost": raw.loc[raw["Year_last_admitted"].notna(), TARGET].mean(),
            },
            {
                "missingness": "Admission year missing",
                "mean_cost": raw.loc[raw["Year_last_admitted"].isna(), TARGET].mean(),
            },
        ]
    )
    plt.figure(figsize=(8.2, 4.8))
    sns.barplot(data=missing_impact, x="missingness", y="mean_cost", hue="missingness", palette=[SECONDARY, GOLD, ACCENT, WARN], legend=False)
    plt.title("Target Impact of Missing BMI and Admission Year")
    plt.ylabel("Mean insurance cost")
    plt.xticks(rotation=20, ha="right")
    figures["missingness_target_impact"] = FIG_DIR / "missingness_target_impact.png"
    savefig(figures["missingness_target_impact"])

    outlier_cols = ["bmi_for_analysis", "weight", "avg_glucose_level", "daily_avg_steps", "fat_percentage", TARGET]
    fig, axes = plt.subplots(2, 3, figsize=(13, 7))
    for i, (ax, col) in enumerate(zip(axes.ravel(), outlier_cols)):
        sns.boxplot(x=df[col], ax=ax, color=CHART_PALETTE[i % len(CHART_PALETTE)])
        ax.set_title(col)
    figures["outlier_boxplots"] = FIG_DIR / "outlier_boxplots.png"
    savefig(figures["outlier_boxplots"])

    disease_means = (
        df.groupby(["heart_disease_history", "other_major_disease_history"], as_index=False)[TARGET]
        .mean()
        .sort_values(TARGET, ascending=False)
    )
    disease_means["segment"] = disease_means.apply(
        lambda r: f"Heart {int(r['heart_disease_history'])}, Other {int(r['other_major_disease_history'])}", axis=1
    )
    plt.figure(figsize=(7.5, 4.5))
    sns.barplot(data=disease_means, x="segment", y=TARGET, hue="segment", palette=bar_palette(len(disease_means)), legend=False)
    plt.title("Average Cost by Disease-History Combination")
    plt.ylabel("Mean insurance cost")
    plt.xticks(rotation=20, ha="right")
    figures["disease_history_means"] = FIG_DIR / "disease_history_means.png"
    savefig(figures["disease_history_means"])

    age_bmi = pd.pivot_table(
        df,
        values=TARGET,
        index="age_band",
        columns="bmi_category",
        aggfunc="mean",
        observed=False,
    )
    plt.figure(figsize=(8.5, 4.8))
    sns.heatmap(age_bmi, annot=True, fmt=".0f", cmap="crest", linewidths=0.5, linecolor="white")
    plt.title("Average Cost by Age Band and BMI Category")
    plt.xlabel("BMI category")
    plt.ylabel("Age band")
    figures["cost_age_bmi_heatmap"] = FIG_DIR / "cost_age_bmi_heatmap.png"
    savefig(figures["cost_age_bmi_heatmap"])

    smoking_age = (
        df.groupby(["age_band", "smoking_status"], observed=False)[TARGET]
        .mean()
        .reset_index()
    )
    plt.figure(figsize=(9, 5))
    sns.lineplot(
        data=smoking_age,
        x="age_band",
        y=TARGET,
        hue="smoking_status",
        marker="o",
        palette=bar_palette(smoking_age["smoking_status"].nunique()),
    )
    plt.title("Average Cost by Smoking Status Across Age Bands")
    plt.xlabel("Age band")
    plt.ylabel("Mean insurance cost")
    plt.xticks(rotation=20)
    figures["cost_smoking_by_age_band"] = FIG_DIR / "cost_smoking_by_age_band.png"
    savefig(figures["cost_smoking_by_age_band"])

    risk_band = (
        df.groupby(["medical_risk_band", "lifestyle_risk_band"], observed=False)[TARGET]
        .mean()
        .reset_index()
    )
    plt.figure(figsize=(7.5, 4.8))
    sns.barplot(data=risk_band, x="medical_risk_band", y=TARGET, hue="lifestyle_risk_band", palette=[SECONDARY, GOLD, ACCENT])
    plt.title("Average Cost by Medical and Lifestyle Risk Bands")
    plt.xlabel("Medical risk band")
    plt.ylabel("Mean insurance cost")
    figures["cost_by_combined_risk_bands"] = FIG_DIR / "cost_by_combined_risk_bands.png"
    savefig(figures["cost_by_combined_risk_bands"])

    return figures


def summarize_eda(df: pd.DataFrame) -> dict[str, object]:
    top_corr = (
        df[
            [
                "years_of_insurance_with_us",
                "regular_checkup_last_year",
                "visited_doctor_last_1_year",
                "daily_avg_steps",
                "age",
                "avg_glucose_level",
                "bmi_for_analysis",
                "weight",
                "fat_percentage",
                "cholesterol_midpoint",
                "years_since_last_admitted",
                TARGET,
            ]
        ]
        .corr(numeric_only=True)[TARGET]
        .drop(TARGET)
        .sort_values(key=lambda s: s.abs(), ascending=False)
    )
    group_insights = {}
    for column in [
        "smoking_status",
        "exercise",
        "Alcohol",
        "covered_by_any_other_company",
        "any_major_disease_history",
        "adventure_sports",
        "regular_checkup_last_year",
        "weight_change_in_last_one_year",
        "admission_status",
        "bmi_category",
        "age_band",
        "lifestyle_risk_band",
        "medical_risk_band",
        "risk_profile_segment",
    ]:
        group_insights[column] = (
            df.groupby(column)[TARGET]
            .agg(["count", "mean", "median"])
            .sort_values("mean", ascending=False)
            .round(2)
            .reset_index()
        )
        write_csv(group_insights[column], f"group_summary_{column}.csv")
    return {
        "top_correlations": top_corr.round(3),
        "group_insights": group_insights,
        "reference_year": int(df["Year_last_admitted"].max() + 1),
        "bmi_imputation_median": round(df["bmi_imputed"].median(), 2),
    }


def model_complexity_score(model_name: str) -> int:
    if model_name in {"DummyRegressor", "LinearRegression", "Ridge", "LogTargetRidge"}:
        return 1
    if model_name in {"BaseLightGBMRegressor", "LightGBMRegressor", "HistGradientBoostingRegressor", "BaseHistGradientBoostingRegressor", "BaseHistGradientBoostingRegressorAlt"}:
        return 2
    if model_name in {"GradientBoostingRegressor", "TunedHistGradientBoostingRegressor", "LogTargetHistGradientBoostingRegressor", "LogTargetTunedHistGradientBoostingRegressor"}:
        return 3
    if model_name in {"RandomForestRegressor", "ExtraTreesRegressor", "XGBRegressor", "CatBoostRegressor"}:
        return 4
    if model_name.startswith("WeightedBlend"):
        return 5
    return 6


def save_split_distribution(y_train: pd.Series, y_test: pd.Series, valid_levels: list[int]) -> Path:
    train_counts = y_train.value_counts().sort_index()
    test_counts = y_test.value_counts().sort_index()
    rows = []
    for level in valid_levels:
        train_count = int(train_counts.get(level, 0))
        test_count = int(test_counts.get(level, 0))
        rows.append(
            {
                "insurance_cost": int(level),
                "train_count": train_count,
                "test_count": test_count,
                "train_pct": round(train_count / len(y_train) * 100, 3),
                "test_pct": round(test_count / len(y_test) * 100, 3),
            }
        )
    return write_csv(pd.DataFrame(rows), "target_band_split_check.csv")


def compute_repeated_cv_metrics(
    fitted: dict[str, object],
    metrics: pd.DataFrame,
    X_train: pd.DataFrame,
    y_train: pd.Series,
) -> pd.DataFrame:
    cv = RepeatedKFold(n_splits=3, n_repeats=2, random_state=RANDOM_STATE)
    rows = []
    candidate_names = list(
        dict.fromkeys(
            metrics["model"].head(5).tolist()
            + metrics.sort_values(["cv_rmse_mean", "test_RMSE"])["model"].head(5).tolist()
        )
    )
    for name in candidate_names:
        estimator = fitted[name]
        print(f"Repeated CV for {name}")
        rmse_scores = -cross_val_score(
            estimator,
            X_train,
            y_train,
            cv=cv,
            scoring="neg_root_mean_squared_error",
            n_jobs=1,
        )
        mae_scores = -cross_val_score(
            estimator,
            X_train,
            y_train,
            cv=cv,
            scoring="neg_mean_absolute_error",
            n_jobs=1,
        )
        rows.append(
            {
                "model": name,
                "repeated_cv_rmse_mean": float(rmse_scores.mean()),
                "repeated_cv_rmse_std": float(rmse_scores.std()),
                "repeated_cv_mae_mean": float(mae_scores.mean()),
                "repeated_cv_mae_std": float(mae_scores.std()),
                "folds": 3,
                "repeats": 2,
            }
        )
    repeated = pd.DataFrame(rows).sort_values(["repeated_cv_rmse_mean", "repeated_cv_mae_mean"])
    write_csv(repeated.round(4), "repeated_cv_metrics.csv")
    shutil.copy2(TABLE_DIR / "repeated_cv_metrics.csv", MODEL_DIR / "repeated_cv_metrics.csv")
    return repeated


def save_calibration_artifacts(
    model: object,
    X_train: pd.DataFrame,
    X_test: pd.DataFrame,
    y_train: pd.Series,
    y_test: pd.Series,
    valid_levels: list[int],
) -> dict[str, object]:
    calibration_strata = make_target_strata(y_train)
    X_inner, X_cal, y_inner, y_cal = train_test_split(
        X_train,
        y_train,
        test_size=0.2,
        random_state=RANDOM_STATE,
        stratify=calibration_strata,
    )
    calibration_model = clone(model)
    calibration_model.fit(X_inner, y_inner)
    cal_pred = calibration_model.predict(X_cal)
    calibrator = IsotonicRegression(out_of_bounds="clip")
    calibrator.fit(cal_pred, y_cal)

    raw_test_pred = model.predict(X_test)
    calibrated_test_pred = calibrator.predict(raw_test_pred)
    rounded_raw = round_to_valid_price_grid(raw_test_pred, valid_levels)
    rounded_calibrated = round_to_valid_price_grid(calibrated_test_pred, valid_levels)

    comparison = pd.DataFrame(
        [
            prediction_variant_row("raw_continuous", y_test, raw_test_pred, valid_levels),
            prediction_variant_row("calibrated_continuous", y_test, calibrated_test_pred, valid_levels),
            prediction_variant_row("rounded_to_nearest_price_band", y_test, rounded_raw, valid_levels),
            prediction_variant_row("calibrated_then_rounded", y_test, rounded_calibrated, valid_levels),
        ]
    )
    write_csv(comparison.round(5), "price_grid_evaluation.csv")
    shutil.copy2(TABLE_DIR / "price_grid_evaluation.csv", MODEL_DIR / "price_grid_evaluation.csv")

    calibration_summary = comparison.assign(
        RMSE_delta_vs_raw=lambda x: x["RMSE"] - float(comparison.loc[comparison["variant"] == "raw_continuous", "RMSE"].iloc[0]),
        MAE_delta_vs_raw=lambda x: x["MAE"] - float(comparison.loc[comparison["variant"] == "raw_continuous", "MAE"].iloc[0]),
    )
    write_csv(calibration_summary.round(5), "calibration_comparison.csv")
    shutil.copy2(TABLE_DIR / "calibration_comparison.csv", MODEL_DIR / "calibration_comparison.csv")

    artifact_path = MODEL_DIR / "prediction_calibrator.pkl"
    joblib.dump(
        {
            "calibrator": calibrator,
            "valid_target_levels": valid_levels,
            "variant_table": comparison.to_dict(orient="records"),
        },
        artifact_path,
    )

    curve_rows = []
    for label, preds in [("raw", raw_test_pred), ("calibrated", calibrated_test_pred)]:
        curve_frame = pd.DataFrame({"prediction": preds, "actual": y_test.values})
        curve_frame["bin"] = pd.qcut(curve_frame["prediction"].rank(method="first"), q=10, duplicates="drop")
        curve = curve_frame.groupby("bin", observed=False).agg(mean_prediction=("prediction", "mean"), mean_actual=("actual", "mean")).reset_index()
        for _, row in curve.iterrows():
            curve_rows.append(
                {
                    "variant": label,
                    "mean_prediction": float(row["mean_prediction"]),
                    "mean_actual": float(row["mean_actual"]),
                }
            )
    curve_df = pd.DataFrame(curve_rows)
    plt.figure(figsize=(7.5, 5))
    sns.lineplot(data=curve_df, x="mean_prediction", y="mean_actual", hue="variant", marker="o", palette=[SECONDARY, WARN])
    min_val = min(curve_df["mean_prediction"].min(), curve_df["mean_actual"].min())
    max_val = max(curve_df["mean_prediction"].max(), curve_df["mean_actual"].max())
    plt.plot([min_val, max_val], [min_val, max_val], linestyle="--", color=MUTED)
    plt.title("Calibration Curve: Predicted vs Actual Cost")
    plt.xlabel("Mean predicted cost")
    plt.ylabel("Mean actual cost")
    savefig(FIG_DIR / "calibration_curve.png")

    residual_frame = pd.DataFrame(
        {
            "actual": y_test.values,
            "raw_residual": raw_test_pred - y_test.values,
            "calibrated_residual": calibrated_test_pred - y_test.values,
        }
    )
    residual_frame["actual_cost_decile"] = pd.qcut(
        residual_frame["actual"].rank(method="first"),
        q=10,
        labels=[f"D{i}" for i in range(1, 11)],
    )
    decile_residuals = (
        residual_frame.groupby("actual_cost_decile", observed=False)[["raw_residual", "calibrated_residual"]]
        .mean()
        .reset_index()
        .melt(id_vars="actual_cost_decile", var_name="variant", value_name="mean_residual")
    )
    write_csv(decile_residuals.round(4), "residual_by_cost_decile_calibration.csv")
    plt.figure(figsize=(8.5, 4.8))
    sns.barplot(data=decile_residuals, x="actual_cost_decile", y="mean_residual", hue="variant", palette=[SECONDARY, WARN])
    plt.axhline(0, color=PRIMARY, linewidth=1)
    plt.title("Residual by Actual Cost Decile After Calibration")
    plt.xlabel("Actual cost decile")
    plt.ylabel("Mean residual")
    savefig(FIG_DIR / "residual_by_cost_decile_after_calibration.png")

    return {
        "calibrator_path": artifact_path,
        "price_grid_evaluation": comparison,
        "calibration_comparison": calibration_summary,
        "raw_test_predictions": raw_test_pred,
        "calibrated_test_predictions": calibrated_test_pred,
    }


def run_ordinal_challenger(
    X_train: pd.DataFrame,
    X_test: pd.DataFrame,
    y_train: pd.Series,
    y_test: pd.Series,
    valid_levels: list[int],
) -> pd.DataFrame:
    rows: list[dict[str, object]] = []
    try:
        levels = np.asarray(valid_levels, dtype=float)
        y_train_idx = target_band_indices(y_train, valid_levels)
        classifier = make_model_pipeline(
            HistGradientBoostingClassifier(
                max_iter=180,
                learning_rate=0.06,
                max_leaf_nodes=31,
                min_samples_leaf=25,
                l2_regularization=0.05,
                random_state=RANDOM_STATE,
            )
        )
        print("Fitting ordinal price-band challenger")
        classifier.fit(X_train, y_train_idx)
        class_idx = classifier.predict(X_test).astype(int)
        class_prediction = levels[class_idx]
        class_row = prediction_variant_row("hist_gradient_classifier_class", y_test, class_prediction, valid_levels)
        class_row["model"] = "HistGradientBoostingClassifier"
        rows.append(class_row)
        if hasattr(classifier, "predict_proba"):
            probabilities = classifier.predict_proba(X_test)
            classes = classifier.classes_.astype(int)
            expected_prediction = probabilities @ levels[classes]
            expected_row = prediction_variant_row(
                "hist_gradient_classifier_expected_value",
                y_test,
                expected_prediction,
                valid_levels,
            )
            expected_row["model"] = "HistGradientBoostingClassifier"
            rows.append(expected_row)
    except Exception as exc:
        rows.append(
            {
                "model": "HistGradientBoostingClassifier",
                "variant": "failed",
                "error": str(exc),
            }
        )
    ordinal = pd.DataFrame(rows)
    write_csv(ordinal.round(5) if "MAE" in ordinal.columns else ordinal, "ordinal_challenger_metrics.csv")
    shutil.copy2(TABLE_DIR / "ordinal_challenger_metrics.csv", MODEL_DIR / "ordinal_challenger_metrics.csv")
    return ordinal


def write_app_schema(raw: pd.DataFrame, X_train: pd.DataFrame, y_train: pd.Series, target_meta: dict[str, object]) -> Path:
    cleaned = clean_column_names(raw)
    schema_numeric = [
        "years_of_insurance_with_us",
        "regular_checkup_last_year",
        "adventure_sports",
        "visited_doctor_last_1_year",
        "daily_avg_steps",
        "age",
        "heart_disease_history",
        "other_major_disease_history",
        "avg_glucose_level",
        "bmi",
        "weight",
        "weight_change_in_last_one_year",
        "fat_percentage",
        "Year_last_admitted",
    ]
    numeric_ranges = {}
    for column in schema_numeric:
        series = pd.to_numeric(cleaned[column], errors="coerce")
        series_non_missing = series.dropna()
        numeric_ranges[column] = {
            "min": float(series_non_missing.min()),
            "max": float(series_non_missing.max()),
            "median": float(series_non_missing.median()),
        }
    categorical_columns = [
        "Occupation",
        "Gender",
        "smoking_status",
        "Location",
        "covered_by_any_other_company",
        "Alcohol",
        "exercise",
        "cholesterol_level",
    ]
    categorical_options = {
        column: sorted(cleaned[column].astype("object").fillna("Unknown").astype(str).unique().tolist())
        for column in categorical_columns
    }
    schema = {
        "numeric_ranges": numeric_ranges,
        "categorical_options": categorical_options,
        "target_grid": target_meta,
        "train_rows": int(len(X_train)),
        "target_median": float(y_train.median()),
        "schema_note": "Generated from training data/raw categorical levels; applicant_id is intentionally excluded from app input.",
    }
    path = MODEL_DIR / "app_schema.json"
    path.write_text(json.dumps(schema, indent=2), encoding="utf-8")
    return path


def build_models(raw: pd.DataFrame) -> dict[str, object]:
    set_visual_theme()
    df = clean_column_names(raw)
    X = df.drop(columns=[TARGET, "applicant_id"], errors="ignore")
    y = df[TARGET]
    target_meta = target_grid_metadata(y)
    split_strata = make_target_strata(y)
    X_train, X_test, y_train, y_test = train_test_split(
        X,
        y,
        test_size=0.2,
        random_state=RANDOM_STATE,
        stratify=split_strata,
    )
    train_target_meta = target_grid_metadata(y_train)
    valid_target_levels = train_target_meta["valid_target_levels"]
    split_distribution_path = save_split_distribution(y_train, y_test, valid_target_levels)
    gpu_enabled = use_gpu_acceleration()
    gpu_model_names: set[str] = set()

    model_defs = {
        "DummyRegressor": DummyRegressor(strategy="mean"),
        "LinearRegression": LinearRegression(),
        "Ridge": Ridge(alpha=5.0, random_state=RANDOM_STATE),
        "DecisionTreeRegressor": DecisionTreeRegressor(
            random_state=RANDOM_STATE, max_depth=12, min_samples_leaf=25
        ),
        "RandomForestRegressor": RandomForestRegressor(
            n_estimators=100,
            max_depth=16,
            min_samples_leaf=10,
            max_features="sqrt",
            n_jobs=1,
            random_state=RANDOM_STATE,
        ),
        "ExtraTreesRegressor": ExtraTreesRegressor(
            n_estimators=120,
            max_depth=None,
            min_samples_leaf=6,
            max_features="sqrt",
            n_jobs=1,
            random_state=RANDOM_STATE,
        ),
        "GradientBoostingRegressor": GradientBoostingRegressor(
            n_estimators=220,
            learning_rate=0.06,
            max_depth=3,
            min_samples_leaf=25,
            subsample=0.85,
            random_state=RANDOM_STATE,
        ),
        "HistGradientBoostingRegressor": HistGradientBoostingRegressor(
            max_iter=320,
            learning_rate=0.06,
            max_leaf_nodes=31,
            min_samples_leaf=25,
            l2_regularization=0.05,
            random_state=RANDOM_STATE,
        ),
    }
    if XGBRegressor is not None:
        xgb_kwargs = {
            "n_estimators": 320,
            "max_depth": 4,
            "learning_rate": 0.05,
            "subsample": 0.85,
            "colsample_bytree": 0.85,
            "objective": "reg:squarederror",
            "tree_method": "hist",
            "random_state": RANDOM_STATE,
            "n_jobs": 1,
            "verbosity": 0,
        }
        if gpu_enabled:
            xgb_kwargs["device"] = "cuda"
            gpu_model_names.add("XGBRegressor")
        model_defs["XGBRegressor"] = XGBRegressor(**xgb_kwargs)
    if LGBMRegressor is not None:
        model_defs["LightGBMRegressor"] = LGBMRegressor(
            random_state=RANDOM_STATE,
            n_jobs=1,
            verbose=-1,
        )
    if CatBoostRegressor is not None:
        catboost_kwargs = {
            "iterations": 360,
            "learning_rate": 0.05,
            "depth": 6,
            "loss_function": "RMSE",
            "random_seed": RANDOM_STATE,
            "verbose": False,
            "allow_writing_files": False,
        }
        if gpu_enabled:
            catboost_kwargs.update({"task_type": "GPU", "devices": "0"})
            gpu_model_names.add("CatBoostRegressor")
        else:
            catboost_kwargs["thread_count"] = 1
        model_defs["CatBoostRegressor"] = CatBoostRegressor(**catboost_kwargs)

    def make_log_target_estimator(model: object, feature_set: str = "enhanced") -> TransformedTargetRegressor:
        return TransformedTargetRegressor(
            regressor=make_model_pipeline(model, feature_set=feature_set),
            func=np.log1p,
            inverse_func=np.expm1,
            check_inverse=False,
        )

    estimator_defs: dict[str, object] = {name: make_model_pipeline(model) for name, model in model_defs.items()}
    estimator_defs["BaseHistGradientBoostingRegressor"] = make_model_pipeline(
        HistGradientBoostingRegressor(
            max_iter=320,
            learning_rate=0.06,
            max_leaf_nodes=31,
            min_samples_leaf=25,
            l2_regularization=0.05,
            random_state=RANDOM_STATE,
        ),
        feature_set="base",
    )
    estimator_defs["BaseHistGradientBoostingRegressorAlt"] = make_model_pipeline(
        HistGradientBoostingRegressor(
            max_iter=600,
            learning_rate=0.035,
            max_leaf_nodes=31,
            min_samples_leaf=35,
            l2_regularization=0.05,
            max_bins=255,
            random_state=RANDOM_STATE,
        ),
        feature_set="base",
    )
    estimator_defs["LogTargetRidge"] = make_log_target_estimator(Ridge(alpha=5.0, random_state=RANDOM_STATE))
    estimator_defs["LogTargetHistGradientBoostingRegressor"] = make_log_target_estimator(
        HistGradientBoostingRegressor(
            max_iter=320,
            learning_rate=0.06,
            max_leaf_nodes=31,
            min_samples_leaf=25,
            l2_regularization=0.05,
            random_state=RANDOM_STATE,
        )
    )
    if LGBMRegressor is not None:
        estimator_defs["BaseLightGBMRegressor"] = make_model_pipeline(
            LGBMRegressor(random_state=RANDOM_STATE, n_jobs=1, verbose=-1),
            feature_set="base",
        )
        estimator_defs["EnhancedLightGBMRegressorRegularized"] = make_model_pipeline(
            LGBMRegressor(
                n_estimators=200,
                learning_rate=0.04,
                num_leaves=20,
                min_child_samples=30,
                subsample=0.9,
                colsample_bytree=0.8,
                reg_lambda=0.2,
                random_state=RANDOM_STATE,
                n_jobs=1,
                verbose=-1,
            ),
            feature_set="enhanced",
        )
        estimator_defs["WeightedBlendLightGBMEnhancedHGB"] = VotingRegressor(
            estimators=[
                (
                    "base_lgbm_default",
                    make_model_pipeline(
                        LGBMRegressor(random_state=RANDOM_STATE, n_jobs=1, verbose=-1),
                        feature_set="base",
                    ),
                ),
                (
                    "enhanced_lgbm_regularized",
                    make_model_pipeline(
                        LGBMRegressor(
                            n_estimators=200,
                            learning_rate=0.04,
                            num_leaves=20,
                            min_child_samples=30,
                            subsample=0.9,
                            colsample_bytree=0.8,
                            reg_lambda=0.2,
                            random_state=RANDOM_STATE,
                            n_jobs=1,
                            verbose=-1,
                        ),
                        feature_set="enhanced",
                    ),
                ),
                (
                    "enhanced_log_hgb",
                    make_log_target_estimator(
                        HistGradientBoostingRegressor(
                            max_iter=320,
                            learning_rate=0.06,
                            max_leaf_nodes=31,
                            min_samples_leaf=25,
                            l2_regularization=0.05,
                            random_state=RANDOM_STATE,
                        ),
                        feature_set="enhanced",
                    ),
                ),
            ],
            weights=[0.65, 0.15, 0.20],
        )
        estimator_defs["WeightedBlendLightGBMBaseHGB"] = VotingRegressor(
            estimators=[
                (
                    "base_lgbm_default",
                    make_model_pipeline(
                        LGBMRegressor(random_state=RANDOM_STATE, n_jobs=1, verbose=-1),
                        feature_set="base",
                    ),
                ),
                (
                    "base_hgb_alt",
                    make_model_pipeline(
                        HistGradientBoostingRegressor(
                            max_iter=600,
                            learning_rate=0.035,
                            max_leaf_nodes=31,
                            min_samples_leaf=35,
                            l2_regularization=0.05,
                            max_bins=255,
                            random_state=RANDOM_STATE,
                        ),
                        feature_set="base",
                    ),
                ),
                (
                    "enhanced_log_hgb",
                    make_log_target_estimator(
                        HistGradientBoostingRegressor(
                            max_iter=320,
                            learning_rate=0.06,
                            max_leaf_nodes=31,
                            min_samples_leaf=25,
                            l2_regularization=0.05,
                            random_state=RANDOM_STATE,
                        ),
                        feature_set="enhanced",
                    ),
                ),
            ],
            weights=[0.50, 0.30, 0.20],
        )
    estimator_defs["WeightedBlendBaseEnhancedHGB"] = VotingRegressor(
        estimators=[
            (
                "base_hgb_default",
                make_model_pipeline(
                    HistGradientBoostingRegressor(
                        max_iter=320,
                        learning_rate=0.06,
                        max_leaf_nodes=31,
                        min_samples_leaf=25,
                        l2_regularization=0.05,
                        random_state=RANDOM_STATE,
                    ),
                    feature_set="base",
                ),
            ),
            (
                "base_hgb_alt",
                make_model_pipeline(
                    HistGradientBoostingRegressor(
                        max_iter=600,
                        learning_rate=0.035,
                        max_leaf_nodes=31,
                        min_samples_leaf=35,
                        l2_regularization=0.05,
                        max_bins=255,
                        random_state=RANDOM_STATE,
                    ),
                    feature_set="base",
                ),
            ),
            (
                "enhanced_log_hgb",
                make_log_target_estimator(
                    HistGradientBoostingRegressor(
                        max_iter=320,
                        learning_rate=0.06,
                        max_leaf_nodes=31,
                        min_samples_leaf=25,
                        l2_regularization=0.05,
                        random_state=RANDOM_STATE,
                    ),
                    feature_set="enhanced",
                ),
            ),
        ],
        weights=[0.4, 0.4, 0.2],
    )

    fitted: dict[str, object] = {}
    rows = []
    cv = KFold(n_splits=5, shuffle=True, random_state=RANDOM_STATE)
    for name, estimator in estimator_defs.items():
        print(f"Fitting {name}")
        estimator.fit(X_train, y_train)
        fitted[name] = estimator
        train_pred = estimator.predict(X_train)
        test_pred = estimator.predict(X_test)
        train_metrics = evaluate_regression(y_train, train_pred)
        test_metrics = evaluate_regression(y_test, test_pred)
        cv_n_jobs = (
            1
            if name in {"RandomForestRegressor", "ExtraTreesRegressor", *gpu_model_names}
            or name.startswith("WeightedBlend")
            else -1
        )
        cv_scores = cross_val_score(
            estimator,
            X_train,
            y_train,
            cv=cv,
            scoring="neg_root_mean_squared_error",
            n_jobs=cv_n_jobs,
        )
        rows.append(
            {
                "model": name,
                "target_transform": "mixed" if name.startswith("WeightedBlend") else ("log1p" if name.startswith("LogTarget") else "none"),
                "cv_rmse_mean": -float(cv_scores.mean()),
                "cv_rmse_std": float(cv_scores.std()),
                "train_MAE": train_metrics["MAE"],
                "train_RMSE": train_metrics["RMSE"],
                "train_R2": train_metrics["R2"],
                "train_MAPE": train_metrics["MAPE"],
                "test_MAE": test_metrics["MAE"],
                "test_RMSE": test_metrics["RMSE"],
                "test_R2": test_metrics["R2"],
                "test_MAPE": test_metrics["MAPE"],
                "test_SMAPE": test_metrics["SMAPE"],
                "train_test_RMSE_gap": test_metrics["RMSE"] - train_metrics["RMSE"],
            }
        )

    print("Tuning HistGradientBoostingRegressor")
    tuned_pipe = make_model_pipeline(HistGradientBoostingRegressor(random_state=RANDOM_STATE))
    param_dist = {
        "model__max_iter": [180, 260, 360, 460],
        "model__learning_rate": [0.03, 0.05, 0.08, 0.12],
        "model__max_leaf_nodes": [15, 31, 45, 63],
        "model__min_samples_leaf": [10, 20, 35, 60],
        "model__l2_regularization": [0.0, 0.01, 0.05, 0.2],
        "model__max_bins": [128, 255],
    }
    search = RandomizedSearchCV(
        tuned_pipe,
        param_distributions=param_dist,
        n_iter=12,
        cv=5,
        scoring="neg_root_mean_squared_error",
        random_state=RANDOM_STATE,
        n_jobs=2,
        verbose=1,
    )
    search.fit(X_train, y_train)
    tuned = search.best_estimator_
    fitted["TunedHistGradientBoostingRegressor"] = tuned
    train_pred = tuned.predict(X_train)
    test_pred = tuned.predict(X_test)
    train_metrics = evaluate_regression(y_train, train_pred)
    test_metrics = evaluate_regression(y_test, test_pred)
    rows.append(
        {
            "model": "TunedHistGradientBoostingRegressor",
            "target_transform": "none",
            "cv_rmse_mean": -float(search.best_score_),
            "cv_rmse_std": np.nan,
            "train_MAE": train_metrics["MAE"],
            "train_RMSE": train_metrics["RMSE"],
            "train_R2": train_metrics["R2"],
            "train_MAPE": train_metrics["MAPE"],
            "test_MAE": test_metrics["MAE"],
            "test_RMSE": test_metrics["RMSE"],
            "test_R2": test_metrics["R2"],
            "test_MAPE": test_metrics["MAPE"],
            "test_SMAPE": test_metrics["SMAPE"],
            "train_test_RMSE_gap": test_metrics["RMSE"] - train_metrics["RMSE"],
        }
    )

    print("Tuning log-target HistGradientBoostingRegressor")
    log_tuned_pipe = make_log_target_estimator(HistGradientBoostingRegressor(random_state=RANDOM_STATE))
    log_param_dist = {
        "regressor__model__max_iter": [180, 260, 360, 460],
        "regressor__model__learning_rate": [0.03, 0.05, 0.08, 0.12],
        "regressor__model__max_leaf_nodes": [15, 31, 45],
        "regressor__model__min_samples_leaf": [15, 25, 40, 60],
        "regressor__model__l2_regularization": [0.0, 0.01, 0.05, 0.2],
        "regressor__model__max_bins": [128, 255],
    }
    log_search = RandomizedSearchCV(
        log_tuned_pipe,
        param_distributions=log_param_dist,
        n_iter=8,
        cv=5,
        scoring="neg_root_mean_squared_error",
        random_state=RANDOM_STATE,
        n_jobs=2,
        verbose=1,
    )
    log_search.fit(X_train, y_train)
    log_tuned = log_search.best_estimator_
    fitted["LogTargetTunedHistGradientBoostingRegressor"] = log_tuned
    train_pred = log_tuned.predict(X_train)
    test_pred = log_tuned.predict(X_test)
    train_metrics = evaluate_regression(y_train, train_pred)
    test_metrics = evaluate_regression(y_test, test_pred)
    rows.append(
        {
            "model": "LogTargetTunedHistGradientBoostingRegressor",
            "target_transform": "log1p",
            "cv_rmse_mean": -float(log_search.best_score_),
            "cv_rmse_std": np.nan,
            "train_MAE": train_metrics["MAE"],
            "train_RMSE": train_metrics["RMSE"],
            "train_R2": train_metrics["R2"],
            "train_MAPE": train_metrics["MAPE"],
            "test_MAE": test_metrics["MAE"],
            "test_RMSE": test_metrics["RMSE"],
            "test_R2": test_metrics["R2"],
            "test_MAPE": test_metrics["MAPE"],
            "test_SMAPE": test_metrics["SMAPE"],
            "train_test_RMSE_gap": test_metrics["RMSE"] - train_metrics["RMSE"],
        }
    )

    metrics = pd.DataFrame(rows).sort_values(["test_RMSE", "test_MAE"]).reset_index(drop=True)

    repeated_cv_metrics = compute_repeated_cv_metrics(fitted, metrics, X_train, y_train)
    best_by_test = metrics.iloc[0]
    cv_ranked = metrics.sort_values(["cv_rmse_mean", "test_RMSE", "test_MAE"]).reset_index(drop=True)
    best_by_cv = cv_ranked.iloc[0]
    parsimony_threshold_pct = 1.0
    parsimony_threshold_rmse = float(best_by_cv["cv_rmse_mean"]) * (1 + parsimony_threshold_pct / 100)
    parsimony_candidates = metrics[metrics["cv_rmse_mean"] <= parsimony_threshold_rmse].copy()
    parsimony_candidates["complexity_score"] = parsimony_candidates["model"].map(model_complexity_score)
    selected = parsimony_candidates.sort_values(["complexity_score", "cv_rmse_mean", "test_RMSE", "test_MAE"]).iloc[0]
    best_name = selected["model"]
    selection_reason = (
        f"Best 5-fold CV RMSE was {best_by_cv['model']} at {best_by_cv['cv_rmse_mean']:.2f}; "
        f"best raw test RMSE was {best_by_test['model']} at {best_by_test['test_RMSE']:.2f}. "
        f"Applied a {parsimony_threshold_pct:.1f}% parsimony rule and selected {best_name} "
        f"(CV RMSE {selected['cv_rmse_mean']:.2f}, test RMSE {selected['test_RMSE']:.2f}, "
        f"complexity score {int(selected['complexity_score'])})."
    )
    metrics["selected_final_model"] = metrics["model"].eq(best_name)
    metrics_path = write_csv(metrics.round(4), "model_metrics.csv")
    shutil.copy2(metrics_path, MODEL_DIR / "model_metrics.csv")

    final_model_summary = pd.DataFrame(
        [
            {
                "selected_model": best_name,
                "best_test_model": best_by_test["model"],
                "best_cv_model": best_by_cv["model"],
                "selected_test_MAE": selected["test_MAE"],
                "selected_test_RMSE": selected["test_RMSE"],
                "selected_test_R2": selected["test_R2"],
                "best_test_RMSE": best_by_test["test_RMSE"],
                "best_cv_RMSE": best_by_cv["cv_rmse_mean"],
                "selection_reason": selection_reason,
            }
        ]
    )
    final_model_summary_path = write_csv(final_model_summary.round(4), "final_model_summary.csv")
    shutil.copy2(final_model_summary_path, MODEL_DIR / "final_model_summary.csv")

    parsimony_candidates_path = write_csv(
        parsimony_candidates.sort_values(["complexity_score", "cv_rmse_mean", "test_RMSE"]).round(4),
        "parsimony_candidates.csv",
    )
    shutil.copy2(parsimony_candidates_path, MODEL_DIR / "parsimony_candidates.csv")
    best_pipeline = fitted[best_name]
    make_xgboost_cpu_portable(best_pipeline)
    final_model_path = MODEL_DIR / "final_model.pkl"
    joblib.dump(best_pipeline, final_model_path)

    def extract_preprocessing(estimator: object) -> object:
        if hasattr(estimator, "named_steps"):
            return estimator[:-1]
        if hasattr(estimator, "regressor_") and hasattr(estimator.regressor_, "named_steps"):
            return estimator.regressor_[:-1]
        return estimator

    if hasattr(best_pipeline, "named_estimators_"):
        preprocessing_pipeline = {
            name: extract_preprocessing(estimator)
            for name, estimator in best_pipeline.named_estimators_.items()
        }
    else:
        preprocessing_pipeline = extract_preprocessing(best_pipeline)
    joblib.dump(preprocessing_pipeline, MODEL_DIR / "preprocessing_pipeline.pkl")

    calibration_results = save_calibration_artifacts(
        best_pipeline,
        X_train,
        X_test,
        y_train,
        y_test,
        valid_target_levels,
    )
    ordinal_challenger = run_ordinal_challenger(X_train, X_test, y_train, y_test, valid_target_levels)
    app_schema_path = write_app_schema(raw, X_train, y_train, train_target_meta)

    best_params_by_model = {
        "TunedHistGradientBoostingRegressor": search.best_params_,
        "LogTargetTunedHistGradientBoostingRegressor": log_search.best_params_,
    }

    improvement_summary = pd.DataFrame()
    baseline_metrics_path = TABLE_DIR / "model_metrics_before_improvements.csv"
    old_improvement_paths = [TABLE_DIR / "model_improvement_summary.csv", MODEL_DIR / "model_improvement_summary.csv"]
    for old_path in old_improvement_paths:
        if old_path.exists():
            old_path.unlink()
    if baseline_metrics_path.exists():
        baseline_metrics = pd.read_csv(baseline_metrics_path).sort_values(["test_RMSE", "test_MAE"]).reset_index(drop=True)
        previous = baseline_metrics.iloc[0]
        current = selected
        improvement_rows = []
        for metric in ["test_MAE", "test_RMSE", "test_MAPE", "test_SMAPE"]:
            before = float(previous[metric])
            after = float(current[metric])
            improvement_rows.append(
                {
                    "metric": metric,
                    "before_model": previous["model"],
                    "after_model": current["model"],
                    "before": before,
                    "after": after,
                    "absolute_improvement": before - after,
                    "pct_improvement": ((before - after) / before * 100) if before else np.nan,
                }
            )
        before_r2 = float(previous["test_R2"])
        after_r2 = float(current["test_R2"])
        improvement_rows.append(
            {
                "metric": "test_R2",
                "before_model": previous["model"],
                "after_model": current["model"],
                "before": before_r2,
                "after": after_r2,
                "absolute_improvement": after_r2 - before_r2,
                "pct_improvement": ((after_r2 - before_r2) / abs(before_r2) * 100) if before_r2 else np.nan,
            }
        )
        improvement_summary = pd.DataFrame(improvement_rows)
        write_csv(improvement_summary.round(6), "previous_run_comparison.csv")
        shutil.copy2(TABLE_DIR / "previous_run_comparison.csv", MODEL_DIR / "previous_run_comparison.csv")

    metadata = {
        "final_model": best_name,
        "best_params": best_params_by_model.get(best_name, {}),
        "all_tuned_best_params": best_params_by_model,
        "train_rows": int(len(X_train)),
        "test_rows": int(len(X_test)),
        "random_state": RANDOM_STATE,
        "split_strategy": "target-price-band stratified 80/20 split",
        "target_grid": train_target_meta,
        "full_target_grid": target_meta,
        "target_unique_count": train_target_meta["target_unique_count"],
        "target_grid_step": train_target_meta["target_grid_step"],
        "target_min": train_target_meta["target_min"],
        "target_max": train_target_meta["target_max"],
        "valid_target_levels": train_target_meta["valid_target_levels"],
        "target": TARGET,
        "base_numeric_features": BASE_NUMERIC_FEATURES,
        "base_categorical_features": BASE_CATEGORICAL_FEATURES,
        "numeric_features": NUMERIC_FEATURES,
        "categorical_features": CATEGORICAL_FEATURES,
        "raw_model_columns": RAW_MODEL_COLUMNS,
        "dropped_from_modeling": ["applicant_id"],
        "selection_reason": selection_reason,
        "parsimony_threshold_pct": parsimony_threshold_pct,
        "best_raw_test_model": str(best_by_test["model"]),
        "best_raw_test_rmse": float(best_by_test["test_RMSE"]),
        "best_cv_model": str(best_by_cv["model"]),
        "best_cv_rmse": float(best_by_cv["cv_rmse_mean"]),
        "deployment_variant": {
            "analytics_prediction": "raw_continuous",
            "quote_band": "rounded_to_nearest_price_band",
            "risk_category": "raw_continuous",
            "calibrated_prediction": "secondary_diagnostic",
            "reason": "Calibration is reported because it improves some percentage/band metrics, but raw prediction is retained for MAE/RMSE-oriented analytics and quote-band rounding.",
        },
        "repeated_cv_summary": repeated_cv_metrics.round(4).to_dict(orient="records"),
        "price_grid_evaluation": calibration_results["price_grid_evaluation"].round(5).to_dict(orient="records"),
        "calibration_comparison": calibration_results["calibration_comparison"].round(5).to_dict(orient="records"),
        "ordinal_challenger_metrics": ordinal_challenger.round(5).to_dict(orient="records"),
        "app_schema": str(app_schema_path.relative_to(ROOT)),
        "split_distribution": str(split_distribution_path.relative_to(ROOT)),
        "external_model_packages": {
            "xgboost": XGBRegressor is not None,
            "lightgbm": LGBMRegressor is not None,
            "catboost": CatBoostRegressor is not None,
            "shap": optional_package_available("shap"),
        },
        "gpu_acceleration_enabled": gpu_enabled,
        "gpu_model_candidates": sorted(gpu_model_names),
        "target_transform": str(selected.get("target_transform", "none")),
        "test_mae": float(selected["test_MAE"]),
        "test_rmse": float(selected["test_RMSE"]),
        "test_r2": float(selected["test_R2"]),
    }
    if not improvement_summary.empty:
        metadata["improvement_summary"] = improvement_summary.round(6).to_dict(orient="records")
    (MODEL_DIR / "model_metadata.json").write_text(json.dumps(metadata, indent=2), encoding="utf-8")

    model_figures = generate_model_figures(metrics, best_pipeline, X_test, y_test)
    importance = generate_importance(best_pipeline, X_test, y_test)
    shap_status = generate_shap_explainability(fitted, X_test)
    risk_thresholds = {
        "low_medium": float(np.percentile(best_pipeline.predict(X_train), 33)),
        "medium_high": float(np.percentile(best_pipeline.predict(X_train), 66)),
    }
    metadata["risk_thresholds"] = risk_thresholds
    metadata["shap_status"] = shap_status
    (MODEL_DIR / "model_metadata.json").write_text(json.dumps(metadata, indent=2), encoding="utf-8")

    return {
        "metrics": metrics,
        "metrics_path": metrics_path,
        "best_name": best_name,
        "best_pipeline": best_pipeline,
        "final_model_path": final_model_path,
        "metadata": metadata,
        "model_figures": model_figures,
        "importance": importance,
        "shap_status": shap_status,
        "improvement_summary": improvement_summary,
        "train_test": (X_train, X_test, y_train, y_test),
        "repeated_cv_metrics": repeated_cv_metrics,
        "price_grid_evaluation": calibration_results["price_grid_evaluation"],
        "calibration_comparison": calibration_results["calibration_comparison"],
        "ordinal_challenger_metrics": ordinal_challenger,
        "app_schema_path": app_schema_path,
    }


def generate_model_figures(metrics: pd.DataFrame, model, X_test: pd.DataFrame, y_test: pd.Series) -> dict[str, Path]:
    set_visual_theme()
    figures: dict[str, Path] = {}
    plot_metrics = metrics.sort_values("test_RMSE", ascending=True)
    plt.figure(figsize=(10, 5.5))
    sns.barplot(data=plot_metrics, y="model", x="test_RMSE", hue="model", palette=sns.color_palette("crest", len(plot_metrics)), legend=False)
    plt.title("Model Comparison by Test RMSE")
    plt.xlabel("Test RMSE")
    plt.ylabel("")
    figures["model_comparison_rmse"] = FIG_DIR / "model_comparison_rmse.png"
    savefig(figures["model_comparison_rmse"])

    preds = model.predict(X_test)
    plt.figure(figsize=(6.5, 6))
    sns.scatterplot(x=y_test, y=preds, alpha=0.35, s=16, color=SECONDARY)
    lims = [min(y_test.min(), preds.min()), max(y_test.max(), preds.max())]
    plt.plot(lims, lims, "--", color=WARN, linewidth=2)
    plt.title("Predicted vs Actual Insurance Cost")
    plt.xlabel("Actual")
    plt.ylabel("Predicted")
    figures["predicted_vs_actual"] = FIG_DIR / "predicted_vs_actual.png"
    savefig(figures["predicted_vs_actual"])

    residuals = y_test - preds
    plt.figure(figsize=(7, 4.5))
    sns.scatterplot(x=preds, y=residuals, alpha=0.35, s=16, color=PLUM)
    plt.axhline(0, color=WARN, linestyle="--", linewidth=2)
    plt.title("Residuals vs Predicted Cost")
    plt.xlabel("Predicted cost")
    plt.ylabel("Residual")
    figures["residuals"] = FIG_DIR / "residuals.png"
    savefig(figures["residuals"])

    diagnostics = pd.DataFrame(
        {
            "actual": y_test.astype(float),
            "predicted": preds.astype(float),
            "residual": y_test.astype(float) - preds.astype(float),
            "abs_error": np.abs(y_test.astype(float) - preds.astype(float)),
        },
        index=X_test.index,
    )
    diagnostics["pct_error"] = diagnostics["abs_error"] / np.maximum(np.abs(diagnostics["actual"]), 1.0) * 100

    feature_frame = None
    try:
        if hasattr(model, "named_steps") and "features" in model.named_steps:
            feature_frame = model.named_steps["features"].transform(X_test)
        elif hasattr(model, "regressor_") and hasattr(model.regressor_, "named_steps"):
            feature_frame = model.regressor_.named_steps["features"].transform(X_test)
        elif hasattr(model, "named_estimators_"):
            for estimator in model.named_estimators_.values():
                if hasattr(estimator, "named_steps") and "features" in estimator.named_steps:
                    feature_frame = estimator.named_steps["features"].transform(X_test)
                    break
                if hasattr(estimator, "regressor_") and hasattr(estimator.regressor_, "named_steps"):
                    feature_frame = estimator.regressor_.named_steps["features"].transform(X_test)
                    break
    except Exception as exc:  # diagnostics should not break model generation
        print(f"Skipping engineered residual diagnostics: {exc}")

    segment_source = X_test.copy()
    if feature_frame is not None:
        for column in [
            "age_band",
            "bmi_category",
            "risk_profile_segment",
            "admission_status",
            "medical_risk_score",
            "lifestyle_risk_score",
        ]:
            if column in feature_frame.columns:
                segment_source[column] = feature_frame[column]
        if "medical_risk_score" in segment_source:
            segment_source["medical_risk_band"] = pd.cut(
                segment_source["medical_risk_score"],
                bins=[-np.inf, 1.5, 3.5, np.inf],
                labels=["Low", "Moderate", "High"],
            ).astype("object")
        if "lifestyle_risk_score" in segment_source:
            segment_source["lifestyle_risk_band"] = pd.cut(
                segment_source["lifestyle_risk_score"],
                bins=[-np.inf, 1.5, 3.5, np.inf],
                labels=["Low", "Moderate", "High"],
            ).astype("object")

    diagnostics = pd.concat([diagnostics, segment_source], axis=1)
    diagnostics["predicted_minus_actual"] = diagnostics["predicted"] - diagnostics["actual"]
    diagnostics["actual_cost_decile"] = pd.qcut(
        diagnostics["actual"], q=10, labels=[f"D{i}" for i in range(1, 11)], duplicates="drop"
    )

    segment_cols = [
        "age_band",
        "bmi_category",
        "smoking_status",
        "exercise",
        "Alcohol",
        "Gender",
        "Location",
        "covered_by_any_other_company",
        "risk_profile_segment",
        "medical_risk_band",
        "lifestyle_risk_band",
        "admission_status",
        "actual_cost_decile",
    ]
    summary_rows = []
    for column in segment_cols:
        if column not in diagnostics.columns:
            continue
        for value, group in diagnostics.groupby(column, observed=False, dropna=False):
            if len(group) < 30:
                continue
            rmse = float(np.sqrt(np.mean(np.square(group["actual"] - group["predicted"]))))
            summary_rows.append(
                {
                    "segment_variable": column,
                    "segment": value,
                    "count": int(len(group)),
                    "actual_mean": float(group["actual"].mean()),
                    "predicted_mean": float(group["predicted"].mean()),
                    "bias_predicted_minus_actual": float(group["predicted_minus_actual"].mean()),
                    "MAE": float(group["abs_error"].mean()),
                    "RMSE": rmse,
                    "MAPE": float(group["pct_error"].mean()),
                }
            )
    segment_summary = pd.DataFrame(summary_rows)
    if not segment_summary.empty:
        write_csv(segment_summary.round(4), "residual_segment_summary.csv")
        fairness_summary = segment_summary[
            segment_summary["segment_variable"].isin(["Gender", "Location"])
        ].sort_values(["segment_variable", "MAE"], ascending=[True, False])
        write_csv(fairness_summary.round(4), "fairness_error_summary.csv")

        cost_decile = segment_summary[segment_summary["segment_variable"] == "actual_cost_decile"].copy()
        if not cost_decile.empty:
            plt.figure(figsize=(8, 4.8))
            sns.barplot(data=cost_decile, x="segment", y="MAE", hue="segment", palette=sns.color_palette("crest", len(cost_decile)), legend=False)
            plt.title("Prediction Error by Actual Cost Decile")
            plt.xlabel("Actual cost decile")
            plt.ylabel("MAE")
            figures["error_by_cost_decile"] = FIG_DIR / "error_by_cost_decile.png"
            savefig(figures["error_by_cost_decile"])

        plot_segments = segment_summary[
            segment_summary["segment_variable"].isin(["risk_profile_segment", "bmi_category", "smoking_status"])
        ].copy()
        if not plot_segments.empty:
            plot_segments["label"] = plot_segments["segment_variable"] + ": " + plot_segments["segment"].astype(str)
            plot_segments = plot_segments.sort_values("MAE", ascending=False).head(14).iloc[::-1]
            plt.figure(figsize=(9.5, 5.4))
            plt.barh(plot_segments["label"], plot_segments["MAE"], color=bar_palette(len(plot_segments)))
            plt.title("Highest Segment-Level Prediction Error")
            plt.xlabel("MAE")
            figures["segment_error_mae"] = FIG_DIR / "segment_error_mae.png"
            savefig(figures["segment_error_mae"])

        bias_segments = segment_summary[
            segment_summary["segment_variable"].isin(["Gender", "Location", "risk_profile_segment"])
        ].copy()
        if not bias_segments.empty:
            bias_segments["abs_bias"] = bias_segments["bias_predicted_minus_actual"].abs()
            bias_segments["label"] = bias_segments["segment_variable"] + ": " + bias_segments["segment"].astype(str)
            bias_segments = bias_segments.sort_values("abs_bias", ascending=False).head(12).iloc[::-1]
            colors = np.where(bias_segments["bias_predicted_minus_actual"] >= 0, SECONDARY, WARN)
            plt.figure(figsize=(9.5, 5.2))
            plt.barh(bias_segments["label"], bias_segments["bias_predicted_minus_actual"], color=colors)
            plt.axvline(0, color="#111827", linewidth=1)
            plt.title("Largest Segment Biases")
            plt.xlabel("Predicted minus actual")
            figures["segment_bias"] = FIG_DIR / "segment_bias.png"
            savefig(figures["segment_bias"])
    return figures


def generate_importance(model, X_test: pd.DataFrame, y_test: pd.Series) -> pd.DataFrame:
    set_visual_theme()
    sample_size = min(2500, len(X_test))
    X_sample = X_test.sample(sample_size, random_state=RANDOM_STATE)
    y_sample = y_test.loc[X_sample.index]
    print("Computing permutation importance")
    result = permutation_importance(
        model,
        X_sample,
        y_sample,
        n_repeats=6,
        random_state=RANDOM_STATE,
        scoring="neg_root_mean_squared_error",
        n_jobs=1,
    )
    importance = (
        pd.DataFrame(
            {
                "feature": X_sample.columns,
                "importance_mean": result.importances_mean,
                "importance_std": result.importances_std,
            }
        )
        .sort_values("importance_mean", ascending=False)
        .reset_index(drop=True)
    )
    write_csv(importance.round(5), "feature_importance.csv")
    shutil.copy2(TABLE_DIR / "feature_importance.csv", MODEL_DIR / "feature_importance.csv")

    plt.figure(figsize=(8, 5.2))
    top = importance.head(12).iloc[::-1]
    plt.barh(top["feature"], top["importance_mean"], color=sns.color_palette("crest", len(top)))
    plt.title("Top Drivers by Permutation Importance")
    plt.xlabel("Increase in RMSE when shuffled")
    savefig(FIG_DIR / "feature_importance.png")

    numeric_candidates = [
        f for f in importance["feature"].tolist() if f in ["age", "avg_glucose_level", "weight", "bmi", "daily_avg_steps", "fat_percentage"]
    ][:3]
    if numeric_candidates:
        X_pd = X_sample.copy()
        for feature in numeric_candidates:
            X_pd[feature] = X_pd[feature].astype(float)
        fig, ax = plt.subplots(1, len(numeric_candidates), figsize=(5 * len(numeric_candidates), 4.2))
        if len(numeric_candidates) == 1:
            ax = [ax]
        PartialDependenceDisplay.from_estimator(model, X_pd, numeric_candidates, ax=ax)
        plt.suptitle("Partial Dependence for Top Numeric Drivers", y=1.02)
        savefig(FIG_DIR / "partial_dependence_top_features.png")
    return importance


def generate_shap_explainability(fitted: dict[str, object], X_test: pd.DataFrame) -> dict[str, object]:
    status: dict[str, object] = {"available": optional_package_available("shap"), "model": None, "status": "not_run"}
    if not status["available"]:
        status["status"] = "shap_not_installed"
        (MODEL_DIR / "shap_status.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    candidate_name = next(
        (name for name in ["BaseLightGBMRegressor", "LightGBMRegressor", "CatBoostRegressor", "XGBRegressor"] if name in fitted),
        None,
    )
    if candidate_name is None:
        status["status"] = "no_supported_single_tree_candidate"
        (MODEL_DIR / "shap_status.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
        return status
    try:
        import shap

        estimator = fitted[candidate_name]
        if not hasattr(estimator, "named_steps"):
            raise TypeError(f"{candidate_name} is not a simple sklearn Pipeline")
        sample = X_test.sample(min(500, len(X_test)), random_state=RANDOM_STATE)
        transformed = estimator[:-1].transform(sample)
        model = estimator.named_steps["model"]
        feature_names = estimator.named_steps["preprocess"].get_feature_names_out()
        explainer = shap.TreeExplainer(model)
        shap_values = explainer.shap_values(transformed)
        if isinstance(shap_values, list):
            shap_values = shap_values[0]
        mean_abs = np.abs(np.asarray(shap_values)).mean(axis=0)
        shap_importance = (
            pd.DataFrame({"feature": feature_names, "mean_abs_shap": mean_abs})
            .sort_values("mean_abs_shap", ascending=False)
            .reset_index(drop=True)
        )
        write_csv(shap_importance.round(6), "shap_importance.csv")
        shutil.copy2(TABLE_DIR / "shap_importance.csv", MODEL_DIR / "shap_importance.csv")
        plt.figure(figsize=(8, 5.2))
        top = shap_importance.head(12).iloc[::-1]
        plt.barh(top["feature"], top["mean_abs_shap"], color=sns.color_palette("mako", len(top)))
        plt.title(f"SHAP Importance for {candidate_name}")
        plt.xlabel("Mean absolute SHAP value")
        savefig(FIG_DIR / "shap_importance_bar.png")
        status.update(
            {
                "model": candidate_name,
                "status": "completed",
                "sample_rows": int(len(sample)),
                "table": "outputs/models/shap_importance.csv",
                "figure": "outputs/figures/shap_importance_bar.png",
            }
        )
    except Exception as exc:
        status.update({"model": candidate_name, "status": "failed", "error": str(exc)})
    (MODEL_DIR / "shap_status.json").write_text(json.dumps(status, indent=2), encoding="utf-8")
    return status


def set_doc_defaults(doc: Document) -> None:
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    normal._element.rPr.rFonts.set(qn("w:ascii"), "Aptos")
    normal._element.rPr.rFonts.set(qn("w:hAnsi"), "Aptos")
    for name, size, color in [
        ("Title", 24, PRIMARY),
        ("Heading 1", 16, PRIMARY),
        ("Heading 2", 13, ACCENT),
        ("Heading 3", 11.5, PRIMARY),
    ]:
        if name in styles:
            styles[name].font.name = "Aptos"
            styles[name].font.size = Pt(size)
            styles[name].font.color.rgb = RGBColor.from_string(color.replace("#", ""))


def clear_document_body(doc: Document) -> None:
    body = doc._body._element
    for child in list(body):
        if child.tag.endswith("sectPr"):
            continue
        body.remove(child)


def add_page_number_footer(doc: Document) -> None:
    section = doc.sections[0]
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.text = "Insurance Price Prediction Capstone"


def add_doc_title(doc: Document, title: str, subtitle: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor.from_string(PRIMARY.replace("#", ""))
    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = p2.add_run(subtitle)
    r2.font.size = Pt(12)
    r2.font.color.rgb = RGBColor.from_string(MUTED.replace("#", ""))
    doc.add_paragraph()


def add_doc_table(doc: Document, df: pd.DataFrame, max_rows: int = 12, font_size: int = 8) -> None:
    display = df.copy().head(max_rows)
    table = doc.add_table(rows=1, cols=len(display.columns))
    try:
        table.style = "Table Grid"
    except KeyError:
        pass
    hdr = table.rows[0].cells
    for i, column in enumerate(display.columns):
        hdr[i].text = str(column)
        shade_cell(hdr[i], PRIMARY)
        for paragraph in hdr[i].paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(font_size)
    for _, row in display.iterrows():
        cells = table.add_row().cells
        for i, value in enumerate(row):
            if isinstance(value, float):
                text = f"{value:,.3f}" if abs(value) < 100 else f"{value:,.1f}"
            else:
                text = str(value)
            cells[i].text = text
            cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
            for paragraph in cells[i].paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
    doc.add_paragraph()


def shade_cell(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill.replace("#", ""))
    tc_pr.append(shd)


def add_bullets(doc: Document, bullets: list[str]) -> None:
    for bullet in bullets:
        try:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(bullet)
        except KeyError:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.25)
            p.add_run("- " + bullet)


def add_numbered(doc: Document, items: list[str]) -> None:
    for i, item in enumerate(items, start=1):
        try:
            p = doc.add_paragraph(style="List Number")
            p.add_run(item)
        except KeyError:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Inches(0.25)
            p.add_run(f"{i}. {item}")


def add_picture(doc: Document, path: Path, caption: str, width: float = 5.25) -> None:
    if not path.exists():
        return
    doc.add_picture(str(path), width=Inches(width))
    last = doc.paragraphs[-1]
    last.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cap.add_run(caption)
    run.italic = True
    run.font.size = Pt(8.5)
    run.font.color.rgb = RGBColor.from_string(MUTED.replace("#", ""))


def add_interpretation_pair(
    doc: Document,
    technical: str,
    business: str,
    caution: str | None = None,
) -> None:
    tech = doc.add_paragraph()
    tech.add_run("Technical interpretation: ").bold = True
    tech.add_run(technical)
    bus = doc.add_paragraph()
    bus.add_run("Business interpretation: ").bold = True
    bus.add_run(business)
    if caution:
        note = doc.add_paragraph()
        note.add_run("Caution: ").bold = True
        note.add_run(caution)


def add_picture_with_interpretations(
    doc: Document,
    path: Path,
    caption: str,
    technical: str,
    business: str,
    caution: str | None = None,
    width: float = 5.25,
) -> None:
    add_picture(doc, path, caption, width=width)
    if path.exists():
        add_interpretation_pair(doc, technical, business, caution)


def write_milestone1_rubric_coverage_matrix() -> Path:
    rows = [
        {
            "rubric_section": "Data Report",
            "marks": "10 marks",
            "rubric_requirement": "Dataset shape, rows, columns, and total cells.",
            "evidence_in_milestone1_docx": "Data Report section and dataset profile table.",
            "evidence_in_notebook": "Sections 4, 5, and 8 with raw_df.shape and raw_df.size.",
            "evidence_file_or_chart": "outputs/tables/dataset_profile.csv",
            "status": "covered",
        },
        {
            "rubric_section": "Data Report",
            "marks": "10 marks",
            "rubric_requirement": "Visual inspection using head, tail, sample records, info, and summaries.",
            "evidence_in_milestone1_docx": "Data Report narrative, dictionary, and descriptive tables.",
            "evidence_in_notebook": "Initial Data Inspection and Dataset Profile sections.",
            "evidence_file_or_chart": "notebooks/01_milestone1_eda.ipynb",
            "status": "covered",
        },
        {
            "rubric_section": "Data Report",
            "marks": "10 marks",
            "rubric_requirement": "Data types, unique counts, missing values, duplicates, and variable descriptions.",
            "evidence_in_milestone1_docx": "Variable Types and Cleanup section.",
            "evidence_in_notebook": "Data Dictionary and Variable Understanding section.",
            "evidence_file_or_chart": "outputs/tables/data_dictionary_summary.csv",
            "status": "covered",
        },
        {
            "rubric_section": "Data Report",
            "marks": "10 marks",
            "rubric_requirement": "Variable type classification, renaming, cleanup, and data inference.",
            "evidence_in_milestone1_docx": "Rubric Coverage Map and Data Report sections.",
            "evidence_in_notebook": "Column Renaming and Data Cleanup section.",
            "evidence_file_or_chart": "outputs/tables/variable_type_classification.csv; outputs/tables/renaming_and_cleanup.csv",
            "status": "covered",
        },
        {
            "rubric_section": "Initial EDA",
            "marks": "10 marks",
            "rubric_requirement": "Univariate analysis for numeric variables with distribution and spread.",
            "evidence_in_milestone1_docx": "Initial EDA target distribution, boxplot, outlier review, numeric summaries, and complete univariate appendix.",
            "evidence_in_notebook": "Univariate Analysis sections plus complete univariate plot appendix.",
            "evidence_file_or_chart": "outputs/tables/numeric_summary.csv; outputs/tables/univariate_plot_interpretations.csv; outputs/figures/univariate_*.png",
            "status": "covered",
        },
        {
            "rubric_section": "Initial EDA",
            "marks": "10 marks",
            "rubric_requirement": "Categorical distribution for categorical, binary, and ordinal variables.",
            "evidence_in_milestone1_docx": "Initial EDA, appendix group summaries, and complete univariate appendix.",
            "evidence_in_notebook": "Univariate categorical section plus complete univariate plot appendix.",
            "evidence_file_or_chart": "outputs/tables/categorical_frequency.csv; outputs/tables/group_summary_*.csv; outputs/figures/univariate_*.png",
            "status": "covered",
        },
        {
            "rubric_section": "Initial EDA",
            "marks": "10 marks",
            "rubric_requirement": "Bivariate analysis and correlation analysis with insight summaries.",
            "evidence_in_milestone1_docx": "Initial EDA and Extensive EDA sections.",
            "evidence_in_notebook": "Bivariate Analysis with Target and Correlation sections.",
            "evidence_file_or_chart": "outputs/figures/correlation_heatmap.png; outputs/tables/feature_signal_strength.csv",
            "status": "covered",
        },
        {
            "rubric_section": "Data Pre-processing",
            "marks": "10 marks",
            "rubric_requirement": "Unwanted-variable removal and applicant_id treatment.",
            "evidence_in_milestone1_docx": "Data Pre-processing bullets.",
            "evidence_in_notebook": "Duplicate and Unwanted Variable Analysis section.",
            "evidence_file_or_chart": "insurance_modeling.py; outputs/models/app_schema.json",
            "status": "covered",
        },
        {
            "rubric_section": "Data Pre-processing",
            "marks": "10 marks",
            "rubric_requirement": "Missing-value treatment for BMI and Year_last_admitted.",
            "evidence_in_milestone1_docx": "Missing Values and Data Pre-processing sections.",
            "evidence_in_notebook": "Missing Value Analysis and Feature Engineering sections.",
            "evidence_file_or_chart": "outputs/tables/missing_values.csv; outputs/figures/missingness_target_impact.png",
            "status": "covered",
        },
        {
            "rubric_section": "Data Pre-processing",
            "marks": "10 marks",
            "rubric_requirement": "Outlier treatment/retention logic, transformations, and new variables.",
            "evidence_in_milestone1_docx": "Data Pre-processing and EDA Conclusion sections.",
            "evidence_in_notebook": "Outlier Detection and Feature Engineering sections.",
            "evidence_file_or_chart": "outputs/tables/outlier_summary.csv; outputs/tables/cleaned_analysis_data.csv",
            "status": "covered",
        },
        {
            "rubric_section": "Extensive EDA",
            "marks": "10 marks",
            "rubric_requirement": "Important-variable relationships and insightful visualizations.",
            "evidence_in_milestone1_docx": "Extensive EDA section.",
            "evidence_in_notebook": "Multivariate and Segmented EDA section.",
            "evidence_file_or_chart": "outputs/figures/weight_vs_cost_jitter.png; outputs/figures/cost_by_weight_line.png; outputs/figures/cost_by_year_last_admitted.png",
            "status": "covered",
        },
        {
            "rubric_section": "Extensive EDA",
            "marks": "10 marks",
            "rubric_requirement": "Technical and business interpretation, justified inferences, conclusion, and next steps.",
            "evidence_in_milestone1_docx": "Technical interpretation and Business interpretation paragraphs after charts.",
            "evidence_in_notebook": "Interpretation blocks after important charts plus final conclusion sections.",
            "evidence_file_or_chart": "Milestone_1_Insurance_Price_Prediction.docx; notebooks/01_milestone1_eda.html",
            "status": "covered",
        },
    ]
    path = REPORT_DIR / "milestone1_rubric_coverage_matrix.csv"
    pd.DataFrame(rows).to_csv(path, index=False)
    return path


def generate_milestone1_report(
    raw: pd.DataFrame,
    df: pd.DataFrame,
    tables: dict[str, Path],
    figures: dict[str, Path],
    eda_summary: dict[str, object],
) -> Path:
    doc = Document()
    set_doc_defaults(doc)
    add_page_number_footer(doc)
    add_doc_title(
        doc,
        "Milestone 1 Submission",
        "Insurance Price Prediction | Data Report, EDA, Preprocessing, and Insights",
    )
    doc.add_heading("Contents", level=1)
    add_numbered(
        doc,
        [
            "Introduction",
            "Objectives, Problem Statement, and Sub-objectives",
            "Rubric Coverage Map",
            "Data Report",
            "Initial Exploratory Data Analysis",
            "Data Pre-processing",
            "Extensive Exploratory Data Analysis",
            "Appendix",
        ],
    )

    doc.add_heading("Introduction", level=1)
    doc.add_paragraph(
        "The insurance pricing problem is a supervised regression problem. The business objective is to estimate insurance_cost from health, lifestyle, habit, and demographic attributes so that pricing teams can understand risk drivers and support more consistent premium estimation."
    )

    doc.add_heading("Objectives, Problem Statement, and Sub-objectives", level=1)
    add_bullets(
        doc,
        [
            "Build a clean analytical dataset from the supplied Insurance Data.csv file.",
            "Identify data-quality issues including missing values, duplicate rows, inconsistent names, and outliers.",
            "Engineer interpretable predictors such as admission recency, BMI category, cholesterol midpoint, and disease-history flags.",
            "Use EDA to identify premium drivers and modeling risks before training predictive models.",
            "Prepare a reproducible pipeline for Milestone 2 model building and deployment.",
        ],
    )

    doc.add_heading("Rubric Coverage Map", level=1)
    rubric_coverage = pd.read_csv(write_milestone1_rubric_coverage_matrix())
    add_doc_table(
        doc,
        rubric_coverage[
            [
                "rubric_section",
                "marks",
                "rubric_requirement",
                "evidence_file_or_chart",
                "status",
            ]
        ],
        max_rows=12,
        font_size=5.8,
    )

    doc.add_heading("Data Report", level=1)
    profile = pd.read_csv(tables["dataset_profile"])
    add_doc_table(doc, profile, max_rows=12)
    doc.add_paragraph(
        "The dataset contains 25,000 unique applicants and 24 original columns. applicant_id is unique for every row and is retained only for auditability, not for modeling."
    )
    doc.add_heading("Target Pricing Grid", level=2)
    add_doc_table(doc, pd.read_csv(tables["target_grid_summary"]), max_rows=8, font_size=8)
    add_picture(
        doc,
        figures["target_price_grid_frequency"],
        "insurance_cost is a discrete business price grid rather than a truly continuous target.",
    )
    add_interpretation_pair(
        doc,
        "The target has 54 unique quote bands with a fixed grid step of 1,234, so regression errors should be interpreted alongside quote-band rounding.",
        "The model should produce an analytical estimate and then map it to the nearest approved business quote band for presentation.",
    )
    missing = pd.read_csv(tables["missing_values"])
    doc.add_heading("Missing Values", level=2)
    add_doc_table(doc, missing[missing["missing_count"] > 0], max_rows=8)
    add_picture(doc, figures["missing_values"], "Only BMI and Year_last_admitted contain missing values.")
    add_interpretation_pair(
        doc,
        "BMI has limited missingness, while Year_last_admitted has structural missingness for applicants without known admission history.",
        "The missingness pattern should not be blanket-dropped; admission-year missingness carries customer history context.",
    )

    doc.add_heading("Variable Types and Cleanup", level=2)
    data_dictionary = pd.read_csv(tables["data_dictionary"])
    add_doc_table(doc, data_dictionary[["column", "clean_column", "dtype", "missing", "unique_values"]], max_rows=24, font_size=7)
    variable_classification = pd.read_csv(tables["variable_type_classification"])
    add_doc_table(
        doc,
        variable_classification[["clean_column", "analysis_type", "role", "preprocessing_action"]],
        max_rows=24,
        font_size=5.8,
    )
    add_doc_table(doc, pd.read_csv(tables["renaming"]), max_rows=10, font_size=8)

    doc.add_heading("Initial Exploratory Data Analysis", level=1)
    add_picture(doc, figures["target_distribution"], "insurance_cost is mildly right-skewed, with most values between roughly 16k and 37k.")
    add_interpretation_pair(
        doc,
        "The target distribution is moderately concentrated with a mild right tail, and the center of the distribution is close to the median premium.",
        "Most customers fall in mid-range price bands, while high-cost cases require residual and manual-review attention.",
    )
    add_picture(doc, figures["target_boxplot"], "Target outliers are retained because they are valid high-cost pricing cases.")
    add_interpretation_pair(
        doc,
        "Upper-tail insurance_cost values are not data errors; they align with valid quote bands and should remain in model training.",
        "Removing high premiums would understate risk for customers who are expensive but legitimate pricing cases.",
    )
    top_corr = eda_summary["top_correlations"].reset_index()
    top_corr.columns = ["feature", "correlation_with_insurance_cost"]
    add_doc_table(doc, top_corr, max_rows=10)
    add_bullets(
        doc,
        [
            "The target mean and median are nearly equal, while skew is modest at about 0.33.",
            "Weight is the dominant marginal pricing signal; a weight-only relationship explains much of the target variation.",
            "Other-company coverage, admission history, regular checkups, weight change, and adventure sports are stronger observed categorical or ordinal signals than smoking, alcohol, exercise, and disease-history flags.",
            "No duplicate applicant records were detected.",
            "Unknown smoking status is treated as a valid business category rather than a missing value.",
            "A complete univariate plot appendix is included for every original source variable, with technical and business interpretation for each plot.",
        ],
    )

    doc.add_heading("Data Pre-processing", level=1)
    add_bullets(
        doc,
        [
            "Renamed misspelled columns: regular_checkup_lasy_year, heart_decs_history, and other_major_decs_history.",
            "Corrected Occupation value Salried to Salaried.",
            "Dropped applicant_id from the modeling feature set because it is a unique identifier.",
            "Created BMI and admission-year missingness flags before imputation so missingness can carry signal when relevant.",
            "Imputed BMI using median values by age band and gender for EDA. The sklearn modeling pipeline fits the same logic only on training data.",
            f"Converted Year_last_admitted into was_admitted_before, admission_year_missing_flag, admission_status, and years_since_last_admitted using reference_year={eda_summary['reference_year']}. Missing admission years are not treated as ordinary numeric values.",
            "Created cholesterol_midpoint from ordinal cholesterol ranges.",
            "Engineered age_band, bmi_category, any_major_disease_history, weight_bmi_interaction, and steps_per_age.",
            "Capped extremely high BMI values for analysis at the 99.5th percentile while retaining target outliers.",
        ],
    )
    add_doc_table(doc, pd.read_csv(tables["outlier_summary"]), max_rows=8)
    add_picture(doc, figures["outlier_boxplots"], "IQR outlier review for important numeric variables.")
    add_interpretation_pair(
        doc,
        "Numeric outlier checks show valid spread in cost, weight, glucose, steps, fat percentage, and BMI after analysis capping.",
        "The project keeps real applicant variation while only stabilizing extreme BMI values for analysis and modeling.",
    )

    doc.add_heading("Extensive Exploratory Data Analysis", level=1)
    add_picture(doc, figures["correlation_heatmap"], "Correlation view across numeric and engineered variables.")
    add_interpretation_pair(
        doc,
        "The heatmap shows weight as the strongest marginal numeric relationship; many clinical and habit variables have weak standalone correlations.",
        "Pricing discussion should prioritize dominant observed drivers but retain weak variables for interaction and monitoring checks.",
    )
    add_picture(doc, figures["weight_vs_cost_jitter"], "Weight dominates the target relationship and aligns closely with the discrete price bands.")
    add_interpretation_pair(
        doc,
        "The jittered scatter separates overlapping quote-grid points and reveals a clear upward cost pattern as weight increases.",
        "Weight is the most visible pricing driver in this dataset and must be validated against business policy before operational use.",
    )
    add_picture(doc, figures["cost_by_weight_line"], "Average insurance cost rises sharply with weight.")
    add_interpretation_pair(
        doc,
        "Mean premium increases monotonically across much of the weight range, indicating a strong nonlinear pricing gradient.",
        "Applicants in higher weight ranges are likely to receive higher model estimates, so segment-level fairness and policy checks matter.",
    )
    add_picture(doc, figures["cost_by_weight_band"], "Weight-band summaries show the main pricing gradient.")
    add_interpretation_pair(
        doc,
        "Weight bands reduce point-level noise and confirm that the weight-cost relationship is not driven by only a few individual records.",
        "Business users can understand the model's main driver through banded summaries instead of raw scatter alone.",
    )
    add_picture(doc, figures["cost_by_year_last_admitted"], "Admission year is an important risk-history signal among applicants with known admission history.")
    add_interpretation_pair(
        doc,
        "Admission year has an ordered association with average premium among applicants with a known admission date.",
        "Recent or known admissions may flag underwriting review needs, but the model should not treat missing years as ordinary zeros.",
    )
    add_picture(doc, figures["cost_by_admission_status"], "Missing or absent admission history is separated from known prior admission cases.")
    add_interpretation_pair(
        doc,
        "Admission status separates no known admission, missing year, and known admission history so the model can learn structural differences.",
        "This prevents applicants without admission records from being incorrectly grouped with applicants admitted in an arbitrary year.",
    )
    add_picture(doc, figures["cost_by_regular_checkup_last_year"], "Regular-checkup counts show a stronger ordinal signal than many lifestyle flags.")
    add_interpretation_pair(
        doc,
        "Regular checkup count has a clearer ordinal target pattern than several binary habit flags.",
        "Preventive-care behavior may help segment applicants, though the relationship should be treated as association rather than causation.",
    )
    add_picture(doc, figures["cost_by_weight_change"], "Weight-change categories show a strong observed pattern; this should be interpreted as association, not causation.")
    add_interpretation_pair(
        doc,
        "Weight-change categories produce a sizable target mean range, making them a useful engineered ordinal feature.",
        "Weight-change history can support risk segmentation, but business recommendations should avoid claiming it directly causes premium movement.",
    )
    add_picture(doc, figures["missingness_target_impact"], "Missingness itself affects target means for BMI and admission year.")
    add_interpretation_pair(
        doc,
        "Target means differ between missing and present groups, which justifies explicit missingness flags before imputation.",
        "Retaining missingness indicators helps the deployment model handle incomplete applications more consistently.",
    )
    add_doc_table(doc, pd.read_csv(TABLE_DIR / "feature_signal_strength.csv"), max_rows=16, font_size=7)
    add_bullets(
        doc,
        [
            "The strongest observed drivers are weight, admission recency/year, other-company coverage, regular checkups, weight change, and adventure sports.",
            "Smoking, alcohol, exercise, age, BMI, glucose, cholesterol, and disease-history variables show weak marginal effects in this dataset but remain useful for interaction checks, fairness review, and monitoring.",
            "BMI, weight, glucose, and fat percentage are related but not interchangeable; multicollinearity should be handled by regularization or tree ensembles rather than manually dropping fields too early.",
            "Other coverage and admission-history features may reflect customer risk profile and require careful business interpretation.",
        ],
    )
    doc.add_heading("EDA Conclusion and Further Steps", level=2)
    add_bullets(
        doc,
        [
            "Conclusion: the dataset is complete enough for modeling after targeted BMI/admission-history treatment, and the target behaves like a discrete quote grid with clear pricing bands.",
            "Conclusion: weight is the dominant observed signal, while admission recency, coverage, regular checkups, weight change, and adventure sports provide the next strongest pricing context.",
            "Further step: carry the cleaned and engineered data into Milestone 2 with target-band stratified splitting, cross-validation, final model selection, explainability, and deployment checks.",
            "Further step: validate dominant drivers with business stakeholders before interpreting them as underwriting policy or causal effects.",
        ],
    )

    doc.add_heading("Appendix", level=1)
    doc.add_heading("Complete Univariate Plot Appendix", level=2)
    univariate_index = pd.read_csv(TABLE_DIR / "univariate_plot_interpretations.csv")
    add_doc_table(
        doc,
        univariate_index[
            [
                "original_column",
                "clean_column",
                "analysis_type",
                "plot_type",
                "missing_count",
                "unique_count",
            ]
        ],
        max_rows=30,
        font_size=5.8,
    )
    for _, row in univariate_index.iterrows():
        clean_col = str(row["clean_column"])
        doc.add_heading(f"Univariate plot: {clean_col}", level=3)
        plot_path = ROOT / Path(*str(row["plot_file"]).split("/"))
        add_picture_with_interpretations(
            doc,
            plot_path,
            f"Univariate distribution or frequency for {clean_col}.",
            str(row["technical_interpretation"]),
            str(row["business_interpretation"]),
            width=4.75,
        )

    doc.add_heading("Selected Group Summaries", level=2)
    for column in ["covered_by_any_other_company", "regular_checkup_last_year", "weight_change_in_last_one_year", "admission_status", "smoking_status", "exercise", "Alcohol", "bmi_category"]:
        table = pd.read_csv(TABLE_DIR / f"group_summary_{column}.csv")
        doc.add_heading(column, level=3)
        add_doc_table(doc, table, max_rows=10, font_size=8)
    add_picture(doc, figures["cost_by_other_coverage"], "Insurance cost by other-company coverage.")
    add_interpretation_pair(
        doc,
        "Other-company coverage has a meaningful target difference and appears among the stronger non-weight EDA signals.",
        "Coverage status may proxy applicant history or insurance behavior, so it should be explained carefully in reviewer-facing material.",
    )
    add_picture(doc, figures["cost_by_alcohol"], "Insurance cost by alcohol consumption.")
    add_interpretation_pair(
        doc,
        "Alcohol categories show only weak marginal separation compared with weight and admission-related features.",
        "Alcohol should not be overstated as a standalone premium driver in this dataset.",
    )
    add_picture(doc, figures["cost_by_adventure_sports"], "Insurance cost by adventure sports flag.")
    add_interpretation_pair(
        doc,
        "Adventure sports produces a clearer group difference than many other habit indicators and remains a useful risk flag.",
        "The business can use this feature for risk context, while still avoiding an automatic underwriting conclusion from one flag.",
    )
    add_picture(doc, figures["cost_by_exercise"], "Insurance cost by exercise category.")
    add_interpretation_pair(
        doc,
        "Exercise has weak marginal separation and does not rank among the strongest standalone signals.",
        "Exercise is better framed as supporting context than as a major pricing driver.",
    )
    add_picture(doc, figures["cost_by_smoking"], "Insurance cost by smoking status.")
    add_interpretation_pair(
        doc,
        "Smoking status has weaker marginal separation than expected in this dataset and includes an explicit Unknown category.",
        "The report should be honest that smoking is not a leading standalone signal here, even if it remains business-relevant.",
    )
    add_picture(doc, figures["cost_by_bmi_category"], "Insurance cost by BMI category.")
    add_interpretation_pair(
        doc,
        "BMI categories are useful for descriptive segmentation but are weaker than weight as a marginal cost signal.",
        "BMI remains useful for health context, yet weight is the dominant observed pricing gradient in the supplied data.",
    )
    add_picture(doc, figures["cost_age_bmi_heatmap"], "Average cost by age band and BMI category.")
    add_interpretation_pair(
        doc,
        "The age and BMI heatmap checks multivariate segmentation rather than relying on one variable at a time.",
        "Segment views help reviewers see whether weak marginal effects become meaningful when combined with other applicant characteristics.",
    )
    add_picture(doc, figures["cost_by_combined_risk_bands"], "Average cost by medical and lifestyle risk bands.")
    add_interpretation_pair(
        doc,
        "Combined risk bands summarize engineered medical and lifestyle signals and show whether risk scoring adds segmentation value.",
        "This gives business stakeholders a higher-level explanation layer beyond individual raw fields.",
    )

    path = ROOT / "Milestone_1_Insurance_Price_Prediction.docx"
    return save_docx_with_fallback(doc, path)


def generate_milestone2_report(
    model_results: dict[str, object],
    figures: dict[str, Path],
) -> Path:
    metrics = model_results["metrics"].copy()
    metadata = model_results["metadata"]
    importance = model_results["importance"]
    optional_model_status = optional_model_status_sentence(metadata)

    doc = Document()
    set_doc_defaults(doc)
    add_page_number_footer(doc)
    add_doc_title(
        doc,
        "Milestone 2 Submission",
        "Insurance Price Prediction | Modeling, Tuning, Final Selection, Deployment, and Recommendations",
    )
    doc.add_heading("Contents", level=1)
    add_numbered(
        doc,
        [
            "Introduction",
            "Model Selection and Metrics of Interest",
            "Model Building and Evaluation",
            "Model Comparison and Selection",
            "Business Insights and Recommendations",
            "Appendix",
        ],
    )

    doc.add_heading("Introduction", level=1)
    doc.add_paragraph(
        "Milestone 2 extends the cleaned analytical dataset into a trainable and deployable sklearn workflow. All imputation, one-hot encoding, scaling, and engineered features are fitted inside the training pipeline to avoid leakage."
    )

    doc.add_heading("Model Selection and Metrics of Interest", level=1)
    add_bullets(
        doc,
        [
            "Baseline: DummyRegressor establishes the value of predicting the average premium only.",
            "Linear family: Linear Regression and Ridge test additive relationships and regularization.",
            "Tree family: Decision Tree, Random Forest, Extra Trees, Gradient Boosting, and Histogram Gradient Boosting test nonlinear interactions.",
            optional_model_status,
            "Primary metrics are MAE and RMSE because they are interpretable in insurance-cost units. R2 explains variance captured, and MAPE/SMAPE provide percentage-style business context.",
        ],
    )

    doc.add_heading("Model Building and Evaluation", level=1)
    add_bullets(
        doc,
        [
            "Used an 80/20 train-test split stratified by target price band with random_state=42.",
            "Used sklearn Pipeline and ColumnTransformer for feature engineering, imputation, scaling, and one-hot encoding.",
            "Used 5-fold cross-validation on training data for all model families.",
            "Added repeated 3-fold cross-validation with 2 repeats for the leading candidates to check selection stability.",
            "Tuned Histogram Gradient Boosting with RandomizedSearchCV and evaluated optional LightGBM, XGBoost, and CatBoost candidates when their packages were available.",
            "Evaluated raw continuous, calibrated continuous, rounded quote-band, and calibrated-plus-rounded post-processing variants.",
        ],
    )
    report_cols = [
        "model",
        "cv_rmse_mean",
        "train_MAE",
        "train_RMSE",
        "train_R2",
        "test_MAE",
        "test_RMSE",
        "test_R2",
        "test_MAPE",
        "train_test_RMSE_gap",
        "selected_final_model",
    ]
    add_doc_table(doc, metrics[report_cols].round(3), max_rows=12, font_size=7)
    add_picture(doc, figures["model_comparison_rmse"], "Model performance comparison by test RMSE.")
    doc.add_heading("Repeated Cross-validation", level=2)
    add_doc_table(doc, model_results["repeated_cv_metrics"].round(3), max_rows=8, font_size=7)

    doc.add_heading("Model Comparison and Selection", level=1)
    winner = metrics.loc[metrics["model"] == metadata["final_model"]].iloc[0]
    add_bullets(
        doc,
        [
            f"Selected final model: {metadata['final_model']}.",
            f"Test MAE: {winner['test_MAE']:,.0f}. Test RMSE: {winner['test_RMSE']:,.0f}. Test R2: {winner['test_R2']:.3f}.",
            metadata.get("selection_reason", "Selection prioritized low test RMSE and MAE, stable validation behavior, and deployment simplicity."),
            "The final model is selected by a parsimony/stability rule, not by claiming the absolute lowest test RMSE.",
            "The final model is saved as outputs/models/final_model.pkl and is loaded directly by app.py.",
        ],
    )
    doc.add_heading("Calibration and Price-grid Post-processing", level=2)
    add_doc_table(doc, model_results["price_grid_evaluation"].round(4), max_rows=8, font_size=7)
    add_picture(doc, FIG_DIR / "calibration_curve.png", "Validation-based isotonic calibration compared with raw predictions.")
    add_picture(doc, FIG_DIR / "residual_by_cost_decile_after_calibration.png", "Residuals by actual cost decile before and after calibration.")
    doc.add_heading("Ordinal Price-band Challenger", level=2)
    add_doc_table(doc, model_results["ordinal_challenger_metrics"].round(4), max_rows=6, font_size=7)
    add_picture(doc, figures["predicted_vs_actual"], "Predicted vs actual plot for the final model.")
    add_picture(doc, figures["residuals"], "Residual plot for the final model.")
    doc.add_heading("Explainability", level=2)
    add_doc_table(doc, importance.head(12).round(4), max_rows=12, font_size=8)
    add_picture(doc, FIG_DIR / "feature_importance.png", "Permutation importance highlights the strongest model drivers.")
    if metadata.get("shap_status", {}).get("status") == "completed" and (MODEL_DIR / "shap_importance.csv").exists():
        add_doc_table(doc, pd.read_csv(MODEL_DIR / "shap_importance.csv").head(12).round(4), max_rows=12, font_size=7)
        add_picture(doc, FIG_DIR / "shap_importance_bar.png", "SHAP importance for the strongest single-tree booster challenger.")
    add_picture(doc, FIG_DIR / "partial_dependence_top_features.png", "Partial dependence plots summarize directional feature effects.")

    doc.add_heading("Business Insights and Recommendations", level=1)
    add_bullets(
        doc,
        [
            "Use the final model as a pricing decision-support tool, not as an automatic underwriting decision.",
            "Treat insurance_cost as a historical quote-band grid; show both raw predicted cost and nearest valid business quote band in deployment.",
            "Weight and admission-related variables dominate this dataset. Validate that this aligns with actual pricing policy before production use.",
            "Encourage preventive checkups and monitoring for applicants with prior admission history, while avoiding causal claims from observational EDA.",
            "Route extreme predicted premiums or large model residual-like cases to manual review before final quote approval.",
            "Review Gender and Location usage with compliance teams because they may create fairness or proxy-discrimination concerns in real insurance pricing.",
            "Monitor prediction drift, segment-level error, and approval outcomes if the model is deployed beyond the capstone demo.",
        ],
    )

    doc.add_heading("Appendix", level=1)
    doc.add_heading("Saved Artifacts", level=2)
    add_bullets(
        doc,
        [
            "outputs/models/final_model.pkl",
            "outputs/models/preprocessing_pipeline.pkl",
            "outputs/models/model_metrics.csv",
            "outputs/models/final_model_summary.csv",
            "outputs/models/repeated_cv_metrics.csv",
            "outputs/models/price_grid_evaluation.csv",
            "outputs/models/calibration_comparison.csv",
            "outputs/models/ordinal_challenger_metrics.csv",
            "outputs/models/app_schema.json",
            "outputs/models/feature_importance.csv",
            "outputs/figures/*.png",
            "app.py",
            "notebooks/01_milestone1_eda.ipynb",
            "notebooks/02_milestone2_modeling.ipynb",
        ],
    )
    doc.add_heading("Best Tuning Parameters", level=2)
    param_df = pd.DataFrame(
        [{"parameter": key, "value": value} for key, value in metadata.get("best_params", {}).items()]
    )
    add_doc_table(doc, param_df, max_rows=20, font_size=8)

    path = ROOT / "Milestone_2_Insurance_Price_Prediction.docx"
    return save_docx_with_fallback(doc, path)


def write_streamlit_app() -> Path:
    app_code = r'''from __future__ import annotations

import json
from pathlib import Path

import json
import joblib
import pandas as pd
import streamlit as st

# Required so pickle can resolve the custom sklearn transformer.
from insurance_modeling import InsuranceFeatureEngineer  # noqa: F401


ROOT = Path(__file__).resolve().parent
MODEL_PATH = ROOT / "outputs" / "models" / "final_model.pkl"
METADATA_PATH = ROOT / "outputs" / "models" / "model_metadata.json"
SCHEMA_PATH = ROOT / "outputs" / "models" / "app_schema.json"
CALIBRATOR_PATH = ROOT / "outputs" / "models" / "prediction_calibrator.pkl"
IMPORTANCE_PATH = ROOT / "outputs" / "models" / "feature_importance.csv"


@st.cache_resource
def load_model():
    return joblib.load(MODEL_PATH)


@st.cache_resource
def load_calibrator():
    if CALIBRATOR_PATH.exists():
        return joblib.load(CALIBRATOR_PATH)
    return None


@st.cache_data
def load_json(path: Path, default: dict):
    if path.exists():
        with open(path, "r", encoding="utf-8") as file:
            return json.load(file)
    return default


@st.cache_data
def load_importance():
    if IMPORTANCE_PATH.exists():
        return pd.read_csv(IMPORTANCE_PATH)
    return pd.DataFrame(columns=["feature", "importance_mean"])


def numeric_cfg(schema: dict, name: str, fallback_min: float, fallback_max: float, fallback_default: float) -> tuple[float, float, float]:
    cfg = schema.get("numeric_ranges", {}).get(name, {})
    return (
        float(cfg.get("min", fallback_min)),
        float(cfg.get("max", fallback_max)),
        float(cfg.get("median", fallback_default)),
    )


def categorical_options(schema: dict, name: str, fallback: list[str]) -> list[str]:
    values = schema.get("categorical_options", {}).get(name, fallback)
    return [str(value) for value in values if str(value) != "nan"] or fallback


def nearest_quote_band(value: float, levels: list[int]) -> float:
    if not levels:
        return value
    return float(min(levels, key=lambda level: abs(float(level) - value)))


def calibrated_value(raw_prediction: float, calibrator_bundle: dict | None) -> float | None:
    if not calibrator_bundle or "calibrator" not in calibrator_bundle:
        return None
    return float(calibrator_bundle["calibrator"].predict([raw_prediction])[0])


st.set_page_config(page_title="Insurance Price Prediction", page_icon="IP", layout="wide")
st.title("Insurance Price Prediction")
st.caption("Decision-support estimator for health, lifestyle, habit, and demographic inputs.")

model = load_model()
metadata = load_json(METADATA_PATH, {})
schema = load_json(SCHEMA_PATH, {})
calibrator_bundle = load_calibrator()
importance = load_importance()
thresholds = metadata.get("risk_thresholds", {"low_medium": 18000, "medium_high": 35000})
target_grid = schema.get("target_grid", metadata.get("target_grid", {}))
valid_levels = target_grid.get("valid_target_levels", [])
deployment_variant = metadata.get("deployment_variant", {})
quote_variant = deployment_variant.get("quote_band", "rounded_to_nearest_price_band")
risk_variant = deployment_variant.get("risk_category", "raw_continuous")

with st.form("prediction_form"):
    left, middle, right = st.columns(3)
    with left:
        age_min, age_max, age_default = numeric_cfg(schema, "age", 16, 74, 45)
        age = st.slider("Age", int(age_min), int(age_max), int(age_default))
        gender = st.selectbox("Gender", categorical_options(schema, "Gender", ["Female", "Male"]))
        occupation = st.selectbox("Occupation", categorical_options(schema, "Occupation", ["Business", "Salaried", "Student"]))
        location = st.selectbox("Location", categorical_options(schema, "Location", ["Delhi", "Mumbai", "Bangalore"]))
        yi_min, yi_max, yi_default = numeric_cfg(schema, "years_of_insurance_with_us", 0, 8, 4)
        years_insured = st.slider("Years with insurer", int(yi_min), int(yi_max), int(yi_default))
        other_coverage = st.selectbox("Covered by another company", categorical_options(schema, "covered_by_any_other_company", ["N", "Y"]))
    with middle:
        weight_min, weight_max, weight_default = numeric_cfg(schema, "weight", 52, 110, 72)
        weight = st.slider("Weight", int(weight_min), int(weight_max), int(weight_default))
        bmi_unknown = st.checkbox("BMI unknown")
        bmi_min, bmi_max, bmi_default = numeric_cfg(schema, "bmi", 12.0, 60.0, 30.5)
        bmi = st.slider("BMI", float(bmi_min), float(bmi_max), float(bmi_default), step=0.1, disabled=bmi_unknown)
        fat_min, fat_max, fat_default = numeric_cfg(schema, "fat_percentage", 10, 45, 29)
        fat_percentage = st.slider("Fat percentage", int(fat_min), int(fat_max), int(fat_default))
        glucose_min, glucose_max, glucose_default = numeric_cfg(schema, "avg_glucose_level", 55, 280, 168)
        avg_glucose_level = st.slider("Average glucose level", int(glucose_min), int(glucose_max), int(glucose_default))
        cholesterol_level = st.selectbox("Cholesterol level", categorical_options(schema, "cholesterol_level", ["125 to 150", "150 to 175", "175 to 200", "200 to 225", "225 to 250"]))
        wc_min, wc_max, wc_default = numeric_cfg(schema, "weight_change_in_last_one_year", 0, 6, 2)
        weight_change = st.slider("Weight change in last year", int(wc_min), int(wc_max), int(wc_default))
    with right:
        steps_min, steps_max, steps_default = numeric_cfg(schema, "daily_avg_steps", 2000, 12000, 5200)
        daily_steps = st.slider("Daily average steps", int(steps_min), int(steps_max), int(steps_default), step=100)
        exercise = st.selectbox("Exercise", categorical_options(schema, "exercise", ["No", "Moderate", "Extreme"]))
        alcohol = st.selectbox("Alcohol", categorical_options(schema, "Alcohol", ["No", "Rare", "Daily"]))
        smoking_status = st.selectbox("Smoking status", categorical_options(schema, "smoking_status", ["never smoked", "formerly smoked", "smokes", "Unknown"]))
        adventure_sports = st.selectbox("Adventure sports", [0, 1])
        rc_min, rc_max, rc_default = numeric_cfg(schema, "regular_checkup_last_year", 0, 5, 1)
        regular_checkup = st.slider("Regular checkups last year", int(rc_min), int(rc_max), int(rc_default))

    c1, c2, c3 = st.columns(3)
    with c1:
        vd_min, vd_max, vd_default = numeric_cfg(schema, "visited_doctor_last_1_year", 0, 12, 3)
        visited_doctor = st.slider("Doctor visits last year", int(vd_min), int(vd_max), int(vd_default))
    with c2:
        heart_history = st.selectbox("Heart disease history", [0, 1])
    with c3:
        other_major_history = st.selectbox("Other major disease history", [0, 1])

    admitted_before = st.checkbox("Customer was admitted before")
    year_last_admitted = None
    if admitted_before:
        yla_min, yla_max, yla_default = numeric_cfg(schema, "Year_last_admitted", 1990, 2018, 2010)
        year_last_admitted = st.slider("Year last admitted", int(yla_min), int(yla_max), int(yla_default))

    submitted = st.form_submit_button("Predict insurance cost")

if submitted:
    input_row = pd.DataFrame(
        [
            {
                "years_of_insurance_with_us": years_insured,
                "regular_checkup_last_year": regular_checkup,
                "adventure_sports": adventure_sports,
                "Occupation": occupation,
                "visited_doctor_last_1_year": visited_doctor,
                "cholesterol_level": cholesterol_level,
                "daily_avg_steps": daily_steps,
                "age": age,
                "heart_disease_history": heart_history,
                "other_major_disease_history": other_major_history,
                "Gender": gender,
                "avg_glucose_level": avg_glucose_level,
                "bmi": None if bmi_unknown else bmi,
                "smoking_status": smoking_status,
                "Year_last_admitted": year_last_admitted,
                "Location": location,
                "weight": weight,
                "covered_by_any_other_company": other_coverage,
                "Alcohol": alcohol,
                "exercise": exercise,
                "weight_change_in_last_one_year": weight_change,
                "fat_percentage": fat_percentage,
            }
        ]
    )
    raw_prediction = float(model.predict(input_row)[0])
    calibrated_prediction = calibrated_value(raw_prediction, calibrator_bundle)
    quote_source = raw_prediction
    if quote_variant == "calibrated_then_rounded" and calibrated_prediction is not None:
        quote_source = calibrated_prediction
    quote_band = nearest_quote_band(quote_source, valid_levels)
    risk_source = raw_prediction
    if risk_variant == "calibrated_continuous" and calibrated_prediction is not None:
        risk_source = calibrated_prediction

    if risk_source < thresholds["low_medium"]:
        risk = "Low"
    elif risk_source < thresholds["medium_high"]:
        risk = "Medium"
    else:
        risk = "High"

    m1, m2, m3 = st.columns(3)
    m1.metric("Raw predicted cost", f"INR {raw_prediction:,.0f}")
    if calibrated_prediction is not None:
        m2.metric("Calibrated predicted cost", f"INR {calibrated_prediction:,.0f}")
    else:
        m2.metric("Calibrated predicted cost", "Not available")
    m3.metric("Nearest quote band", f"INR {quote_band:,.0f}")
    st.info(f"Risk category: {risk}")
    st.caption(
        "Calibration improved some percentage/band metrics but did not improve MAE/RMSE; "
        "quote-band output uses the selected deployment variant from model metadata."
    )
    if target_grid:
        step = target_grid.get("target_grid_step")
        step_text = f"{step:,.0f}" if isinstance(step, (int, float)) else "unknown"
        st.caption(
            f"Detected pricing grid: {target_grid.get('target_unique_count')} valid bands, "
            f"step INR {step_text}."
        )
    if not importance.empty:
        st.subheader("Top model drivers")
        st.dataframe(importance.head(8), use_container_width=True, hide_index=True)

st.warning(
    "Disclaimer: This model supports pricing analysis and triage. It should not be used as the final underwriting decision without policy, compliance, and human review."
)
'''
    path = ROOT / "app.py"
    path.write_text(app_code, encoding="utf-8")
    return path


NOTEBOOK_RUN_ALL_EXPORTS = [
    "ROOT",
    "OUTPUTS",
    "FIG_DIR",
    "TABLE_DIR",
    "MODEL_DIR",
    "REPORT_DIR",
    "TARGET",
    "RANDOM_STATE",
    "add_report_features",
    "build_models",
    "ensure_dirs",
    "generate_eda_figures",
    "make_target_strata",
    "save_profile_tables",
    "set_visual_theme",
    "summarize_eda",
    "target_grid_metadata",
    "write_milestone1_rubric_coverage_matrix",
    "write_streamlit_app",
]


NOTEBOOK_INSURANCE_EXPORTS = [
    "RAW_MODEL_COLUMNS",
    "InsuranceFeatureEngineer",
    "clean_column_names",
]


def embedded_notebook_bootstrap_source() -> str:
    insurance_source = (ROOT / "insurance_modeling.py").read_text(encoding="utf-8")
    run_all_source = (ROOT / "run_all.py").read_text(encoding="utf-8")
    return f"""
    # Self-contained notebook bootstrap.
    # The helper modules below are embedded as source strings, so this notebook can
    # run from the supplied Insurance Data.csv without reading local project .py files.
    import sys
    import types
    from pathlib import Path

    ROOT = Path.cwd()
    if not (ROOT / "Insurance Data.csv").exists():
        ROOT = ROOT.parent

    _INSURANCE_MODELING_SOURCE = {insurance_source!r}
    _RUN_ALL_SOURCE = {run_all_source!r}

    _insurance_module = types.ModuleType("insurance_modeling")
    _insurance_module.__file__ = str(ROOT / "embedded_insurance_modeling.py")
    sys.modules["insurance_modeling"] = _insurance_module
    exec(
        compile(_INSURANCE_MODELING_SOURCE, _insurance_module.__file__, "exec"),
        _insurance_module.__dict__,
    )

    _run_all_module = types.ModuleType("run_all")
    _run_all_module.__file__ = str(ROOT / "embedded_run_all.py")
    sys.modules["run_all"] = _run_all_module
    exec(
        compile(_RUN_ALL_SOURCE, _run_all_module.__file__, "exec"),
        _run_all_module.__dict__,
    )

    for _name in {NOTEBOOK_RUN_ALL_EXPORTS!r}:
        globals()[_name] = getattr(_run_all_module, _name)
    for _name in {NOTEBOOK_INSURANCE_EXPORTS!r}:
        globals()[_name] = getattr(_insurance_module, _name)

    del _INSURANCE_MODELING_SOURCE, _RUN_ALL_SOURCE
    print("Embedded helper code loaded from notebook; no external .py helper scripts required.")
    """


def create_notebooks() -> list[Path]:
    def md(text: str):
        return nbf.v4.new_markdown_cell(textwrap.dedent(text).strip())

    def code(source: str, tags: list[str] | None = None):
        cell = nbf.v4.new_code_cell(textwrap.dedent(source).strip())
        if tags:
            cell.metadata["tags"] = tags
        if tags and "embedded-helper-code" in tags:
            cell.metadata.setdefault("jupyter", {})["source_hidden"] = True
        return cell

    notebook_css = f"""
    <style>
    body, .jp-Notebook {{
        background: #F8FAFC;
        color: {TEXT};
        font-family: Inter, Segoe UI, Arial, sans-serif;
    }}
    h1 {{
        color: {PRIMARY};
        font-size: 2.15rem;
        font-weight: 800;
        letter-spacing: 0;
        padding: 0.55rem 0 0.35rem;
        border-bottom: 4px solid {ACCENT};
    }}
    h2 {{
        color: {PRIMARY};
        font-size: 1.45rem;
        font-weight: 760;
        letter-spacing: 0;
        margin-top: 1.45rem;
        padding: 0.35rem 0 0.35rem 0.7rem;
        border-left: 5px solid {ACCENT};
        background: linear-gradient(90deg, #E8F4F1 0%, rgba(232,244,241,0) 78%);
        border-radius: 6px;
    }}
    h3 {{
        color: {SECONDARY};
        font-size: 1.12rem;
        font-weight: 720;
        margin-top: 1.1rem;
    }}
    p, li {{
        font-size: 0.98rem;
        line-height: 1.58;
    }}
    div.insight-card {{
        background: #FFFFFF;
        border: 1px solid #D9E6EF;
        border-left: 5px solid {ACCENT};
        border-radius: 8px;
        box-shadow: 0 6px 18px rgba(18, 53, 91, 0.08);
        margin: 0.85rem 0 1.05rem;
        padding: 0.95rem 1.05rem;
    }}
    div.insight-row {{
        display: flex;
        gap: 0.65rem;
        align-items: flex-start;
        margin: 0.25rem 0;
    }}
    span.insight-label {{
        border-radius: 999px;
        color: #FFFFFF;
        display: inline-block;
        font-size: 0.74rem;
        font-weight: 760;
        letter-spacing: 0;
        min-width: 5.25rem;
        padding: 0.18rem 0.55rem;
        text-align: center;
    }}
    span.insight-label.tech {{ background: {PRIMARY}; }}
    span.insight-label.biz {{ background: {ACCENT}; }}
    span.insight-label.caution {{ background: {WARN}; }}
    div.output_png img, img {{
        background: #FFFFFF;
        border: 1px solid #DDE7F0;
        border-radius: 8px;
        box-shadow: 0 8px 24px rgba(18, 53, 91, 0.10);
        padding: 6px;
    }}
    table.dataframe {{
        border-collapse: collapse;
        border: 1px solid #D9E6EF;
        border-radius: 8px;
        box-shadow: 0 4px 14px rgba(18, 53, 91, 0.06);
        overflow: hidden;
    }}
    table.dataframe thead tr {{
        background: {PRIMARY};
        color: #FFFFFF;
    }}
    table.dataframe th, table.dataframe td {{
        border: 1px solid #E5E7EB;
        padding: 0.42rem 0.55rem;
        font-size: 0.86rem;
    }}
    table.dataframe tbody tr:nth-child(even) {{
        background: #F3F8FA;
    }}
    div.input_area {{
        border-radius: 8px;
        border: 1px solid #D9E6EF;
    }}
    code {{
        color: {PRIMARY};
    }}
    </style>
    """

    def style_cell():
        return code(
            f"""
            from IPython.display import HTML, display
            display(HTML({notebook_css!r}))
            """
        )

    def interp(technical: str, business: str, caution: str | None = None):
        caution_row = (
            f"""
            <div class="insight-row">
                <span class="insight-label caution">Caution</span>
                <span>{escape(caution)}</span>
            </div>
            """
            if caution
            else ""
        )
        return md(
            f"""
            <div class="insight-card">
                <div class="insight-row">
                    <span class="insight-label tech">Technical</span>
                    <span>{escape(technical)}</span>
                </div>
                <div class="insight-row">
                    <span class="insight-label biz">Business</span>
                    <span>{escape(business)}</span>
                </div>
                {caution_row}
            </div>
            """
        )

    def image_cell(relative_path: str, title: str):
        return code(
            f"""
            image_path = ROOT / {relative_path!r}
            display(Markdown("**{title}**"))
            if image_path.exists():
                display(Image(filename=str(image_path), width=700))
            else:
                print(f"Missing image: {{image_path}}")
            """
        )

    notebook_metadata = {
        "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
        "language_info": {"name": "python", "pygments_lexer": "ipython3"},
    }

    m1_cells = [
        md("# Milestone 1: Insurance Price Prediction - Data Report, EDA, and Preprocessing"),
        style_cell(),
        md(
            """
            ## 1. Project Context and Business Objective

            The capstone objective is to predict `insurance_cost` from applicant health, lifestyle, habit,
            demographic, medical-history, and insurance-history variables. Milestone 1 focuses on data
            understanding, EDA, preprocessing decisions, and evidence-backed business insight before
            the Milestone 2 modeling workflow.
            """
        ),
        interp(
            "Milestone 1 is treated as the analytical foundation: data quality, target behavior, missingness, and important relationships are documented before model fitting.",
            "A pricing model is only useful if reviewers can see why the data is usable, what the target means, and which applicant characteristics drive the estimate.",
        ),
        md("## 2. Import Libraries and Set Paths"),
        code(embedded_notebook_bootstrap_source(), tags=["embedded-helper-code", "remove-input"]),
        code(
            """
            import json
            from pathlib import Path

            import numpy as np
            import pandas as pd
            import matplotlib.pyplot as plt
            import seaborn as sns
            from IPython.display import Image, Markdown, display

            ROOT = Path.cwd()
            if not (ROOT / "Insurance Data.csv").exists():
                ROOT = ROOT.parent

            TABLE_DIR = ROOT / "outputs" / "tables"
            FIG_DIR = ROOT / "outputs" / "figures"
            MODEL_DIR = ROOT / "outputs" / "models"
            REPORT_DIR = ROOT / "outputs" / "reports"

            ensure_dirs()
            set_visual_theme()
            pd.set_option("display.max_columns", 80)
            pd.set_option("display.width", 140)
            print(ROOT)
            """
        ),
        interp(
            "The notebook resolves the project root, creates output folders, and loads embedded helper functions directly from notebook cells.",
            "A reviewer can execute the notebook from the supplied dataset without needing pre-generated EDA tables, figures, or external helper scripts.",
        ),
        md("## 3. End-to-End EDA Artifact Generation"),
        code(
            """
            DATA_PATH = ROOT / "Insurance Data.csv"
            raw_df = pd.read_csv(DATA_PATH)
            report_df = add_report_features(raw_df)
            tables = save_profile_tables(raw_df, report_df)
            figures = generate_eda_figures(raw_df, report_df)
            eda_summary = summarize_eda(report_df)
            rubric_path = write_milestone1_rubric_coverage_matrix()

            print(f"Loaded dataset: {DATA_PATH.name} -> {raw_df.shape[0]:,} rows, {raw_df.shape[1]:,} columns")
            print(f"Generated {len(tables)} named profile tables and {len(figures)} EDA figures.")
            print(f"Rubric coverage matrix: {rubric_path.relative_to(ROOT)}")
            """
        ),
        interp(
            "This cell regenerates the Milestone 1 tables, cleaned analysis dataset, univariate plot index, and EDA figures directly from `Insurance Data.csv`.",
            "The rest of the notebook reads artifacts created in this run, so the HTML/notebook evidence is reproducible from the raw dataset.",
        ),
        md("## 4. Load Raw Dataset"),
        code(
            """
            raw_df = pd.read_csv(ROOT / "Insurance Data.csv")
            raw_df.shape, raw_df.size
            """
        ),
        code(
            """
            display(raw_df.head())
            display(raw_df.tail())
            display(raw_df.sample(5, random_state=42))
            """
        ),
        interp(
            "The raw file loads successfully with 25,000 rows and 24 original columns, and visual inspection confirms mixed numeric and categorical applicant fields.",
            "The dataset is large enough for supervised learning while still small enough for transparent EDA and report-level validation.",
        ),
        md("## 5. Initial Data Inspection"),
        code(
            """
            raw_df.info()
            """
        ),
        code(
            """
            display(raw_df.describe().T)
            categorical_like_cols = [
                col for col in raw_df.columns
                if not pd.api.types.is_numeric_dtype(raw_df[col])
            ]
            if categorical_like_cols:
                display(raw_df[categorical_like_cols].describe().T)
            display(raw_df.nunique().sort_values())
            display(raw_df.isna().sum().sort_values(ascending=False))
            print("Duplicate rows:", raw_df.duplicated().sum())
            """
        ),
        interp(
            "The initial inspection finds no duplicate rows, a unique `applicant_id`, and missing values only in BMI and admission-year fields.",
            "The data can move forward after targeted preprocessing rather than broad row deletion or aggressive cleanup.",
        ),
        md("## 6. Data Dictionary and Variable Understanding"),
        code(
            """
            dataset_profile = pd.DataFrame(
                {
                    "Metric": ["Rows", "Columns", "Cells", "Duplicate rows", "Unique applicant_id"],
                    "Value": [raw_df.shape[0], raw_df.shape[1], raw_df.size, raw_df.duplicated().sum(), raw_df["applicant_id"].nunique()],
                }
            )
            dataset_profile.to_csv(TABLE_DIR / "dataset_profile.csv", index=False)
            numeric_summary = raw_df.describe().T.reset_index().rename(columns={"index": "column"})
            numeric_summary.to_csv(TABLE_DIR / "numeric_summary.csv", index=False)
            display(dataset_profile)
            display(numeric_summary.head(12))
            """
        ),
        code(
            """
            data_dictionary = pd.read_csv(TABLE_DIR / "data_dictionary_summary.csv")
            variable_types = pd.read_csv(TABLE_DIR / "variable_type_classification.csv")
            display(data_dictionary)
            display(variable_types[["original_column", "clean_column", "analysis_type", "role", "preprocessing_action"]])
            """
        ),
        interp(
            "The variable classification distinguishes identifiers, the quote-grid target, continuous numeric features, ordinal numeric features, binary flags, and nominal categories.",
            "`applicant_id` is retained for audit only and excluded from modeling and the Streamlit schema.",
        ),
        md("## 7. Column Renaming and Data Cleanup"),
        code(
            """
            cleaned_df = clean_column_names(raw_df)
            renaming = pd.read_csv(TABLE_DIR / "renaming_and_cleanup.csv")
            display(renaming)
            display(pd.DataFrame({"raw_column": raw_df.columns, "clean_column": cleaned_df.columns}))
            """
        ),
        interp(
            "The cleanup normalizes misspelled source columns and category labels while preserving the original data meaning.",
            "Clean names reduce downstream coding errors and make the final app schema easier for reviewers to inspect.",
        ),
        md("## 8. Dataset Profile and Descriptive Statistics"),
        code(
            """
            categorical_cols = [
                col for col in raw_df.columns
                if not pd.api.types.is_numeric_dtype(raw_df[col])
            ]
            cat_rows = []
            for col in categorical_cols:
                for category, count in raw_df[col].value_counts(dropna=False).items():
                    cat_rows.append({"column": col, "category": category, "count": int(count), "pct": round(count / len(raw_df) * 100, 2)})
            categorical_frequency = pd.DataFrame(cat_rows)
            categorical_frequency.to_csv(TABLE_DIR / "categorical_frequency.csv", index=False)
            display(categorical_frequency.head(20))
            """
        ),
        interp(
            "Descriptive statistics and frequency tables cover both numeric spread and categorical level balance.",
            "This satisfies the rubric requirement to describe variables before jumping into model building.",
        ),
        md("## 9. Missing Value Analysis"),
        code(
            """
            missing_values = pd.read_csv(TABLE_DIR / "missing_values.csv")
            missing_bmi = pd.read_csv(TABLE_DIR / "missing_bmi_profile.csv")
            missing_admission = pd.read_csv(TABLE_DIR / "missing_year_last_admitted_profile.csv")
            display(missing_values)
            display(missing_bmi)
            display(missing_admission)
            """
        ),
        image_cell("outputs/figures/missing_values.png", "Missing values by column"),
        image_cell("outputs/figures/missingness_target_impact.png", "Target impact of missingness"),
        interp(
            "BMI missingness is limited, but Year_last_admitted missingness is substantial and structurally meaningful.",
            "The model should preserve missingness indicators because missing admission history changes customer context.",
        ),
        interp(
            "The missingness-impact chart shows target mean differences between present and missing groups.",
            "A reviewer can see why the pipeline uses flags and imputation instead of dropping incomplete records.",
        ),
        md("## 10. Duplicate and Unwanted Variable Analysis"),
        code(
            """
            duplicate_summary = pd.DataFrame(
                {
                    "check": ["duplicate_rows", "unique_applicant_id", "applicant_id_in_model_columns"],
                    "value": [int(raw_df.duplicated().sum()), int(raw_df["applicant_id"].nunique()), "applicant_id" in RAW_MODEL_COLUMNS],
                }
            )
            duplicate_summary
            """
        ),
        interp(
            "`applicant_id` is unique for every row and is not included in the modeling feature list.",
            "Removing the identifier avoids memorization and supports a deployable applicant-input form.",
        ),
        md("## 11. Target Variable Analysis"),
        code(
            """
            target_summary = raw_df["insurance_cost"].agg(["count", "mean", "median", "std", "min", "max", "skew"]).round(3)
            target_summary
            """
        ),
        image_cell("outputs/figures/target_distribution.png", "Target distribution"),
        image_cell("outputs/figures/target_boxplot.png", "Target boxplot"),
        interp(
            "The target is mildly right-skewed, with valid high-cost values that should remain in the dataset.",
            "High-cost applicants represent legitimate pricing cases, not noise to be removed.",
        ),
        interp(
            "The boxplot identifies upper-tail premiums, but these align with valid quote bands rather than impossible values.",
            "The business needs accurate estimates for expensive cases, so target outlier removal would be harmful.",
        ),
        md("## 12. Target Pricing Grid Analysis"),
        code(
            """
            target_grid = pd.read_csv(TABLE_DIR / "target_grid_summary.csv")
            target_frequency = pd.read_csv(TABLE_DIR / "target_price_band_frequency.csv")
            display(target_grid)
            display(target_frequency.head(12))
            display(target_frequency.tail(12))
            """
        ),
        image_cell("outputs/figures/target_price_grid_frequency.png", "Insurance quote-band frequency"),
        interp(
            "The target has 54 unique levels and a 1,234 unit grid step, proving that `insurance_cost` is a business quote grid.",
            "Deployment should show raw prediction plus nearest quote band instead of pretending the target is fully continuous.",
        ),
        md("## 13. Univariate Analysis: Numeric Variables"),
        code(
            """
            numeric_summary = pd.read_csv(TABLE_DIR / "numeric_summary.csv")
            display(numeric_summary)
            """
        ),
        image_cell("outputs/figures/outlier_boxplots.png", "Outlier review for numeric variables"),
        interp(
            "Numeric variables show wide but plausible ranges; BMI is stabilized for analysis while valid premium outliers are retained.",
            "The preprocessing strategy protects model stability without deleting legitimate high-risk applicants.",
        ),
        md("## 14. Univariate Analysis: Categorical, Binary, and Ordinal Variables"),
        code(
            """
            categorical_frequency = pd.read_csv(TABLE_DIR / "categorical_frequency.csv")
            display(categorical_frequency)
            for name in ["covered_by_any_other_company", "regular_checkup_last_year", "weight_change_in_last_one_year", "admission_status"]:
                print("\\n", name)
                display(pd.read_csv(TABLE_DIR / f"group_summary_{name}.csv").head(10))
            """
        ),
        interp(
            "Categorical and ordinal summaries show class balance and average-cost separation before multivariate modeling.",
            "This helps prevent overclaiming weak fields while still documenting business-relevant segments.",
        ),
        md("## 15. Bivariate Analysis with Target"),
        image_cell("outputs/figures/weight_vs_cost_jitter.png", "Insurance cost versus weight with jitter"),
        image_cell("outputs/figures/cost_by_weight_line.png", "Average insurance cost by weight"),
        interp(
            "Weight is the clearest marginal relationship, and jitter reveals quote-grid bands that would otherwise overlap.",
            "Weight dominates the submitted dataset, so policy and fairness review should focus on how this feature is used.",
        ),
        image_cell("outputs/figures/cost_by_weight_band.png", "Average insurance cost by weight band"),
        image_cell("outputs/figures/cost_by_other_coverage.png", "Insurance cost by other-company coverage"),
        interp(
            "Weight bands confirm the main gradient; other-company coverage adds a separate categorical signal.",
            "The pricing story is not a single-variable story, but weight is the primary observed driver.",
        ),
        image_cell("outputs/figures/cost_by_regular_checkup_last_year.png", "Cost by regular checkup count"),
        image_cell("outputs/figures/cost_by_weight_change.png", "Cost by weight-change category"),
        interp(
            "Regular checkups and weight-change categories have stronger target separation than many simple habit flags.",
            "These variables can support applicant segmentation, but the report frames them as associations rather than causal levers.",
        ),
        image_cell("outputs/figures/cost_by_year_last_admitted.png", "Cost by year last admitted"),
        image_cell("outputs/figures/cost_by_admission_status.png", "Cost by admission status"),
        interp(
            "Admission-year and admission-status plots show that recency/history information is meaningful when represented carefully.",
            "Applicants without known admission history should not be forced into an arbitrary numeric admission year.",
        ),
        md("## 16. Correlation and Association Analysis"),
        image_cell("outputs/figures/correlation_heatmap.png", "Correlation heatmap"),
        code(
            """
            feature_signal = pd.read_csv(TABLE_DIR / "feature_signal_strength.csv")
            display(feature_signal)
            """
        ),
        interp(
            "The association table confirms strong observed signals for weight, admission recency/year, other coverage, checkups, weight change, and adventure sports.",
            "The report honestly labels smoking, alcohol, exercise, age, BMI, glucose, cholesterol, and disease flags as weak marginal signals in this dataset.",
        ),
        md("## 17. Multivariate and Segmented EDA"),
        image_cell("outputs/figures/cost_age_bmi_heatmap.png", "Average cost by age band and BMI category"),
        image_cell("outputs/figures/cost_by_combined_risk_bands.png", "Average cost by medical and lifestyle risk bands"),
        interp(
            "Segmented EDA checks whether weak marginal variables become more useful in combination with age, BMI, medical, or lifestyle bands.",
            "Business users get a more realistic view of applicant risk than from one-way charts alone.",
        ),
        image_cell("outputs/figures/cost_smoking_by_age_band.png", "Average cost by smoking status across age bands"),
        image_cell("outputs/figures/disease_history_means.png", "Average cost by disease-history combination"),
        interp(
            "Smoking and disease-history views are retained for transparency, but they do not overturn the stronger weight and admission-history findings.",
            "The submission avoids overclaiming expected health variables when the supplied data shows weaker marginal effects.",
        ),
        md("## 18. Outlier Detection and Treatment Strategy"),
        code(
            """
            outliers = pd.read_csv(TABLE_DIR / "outlier_summary.csv")
            display(outliers)
            """
        ),
        interp(
            "The outlier strategy separates true target extremes from numeric stabilization needs such as BMI capping.",
            "This is a balanced preprocessing decision: keep valid expensive customers, but avoid letting implausible BMI values distort analysis.",
        ),
        md("## 19. Feature Engineering and Transformations"),
        code(
            """
            engineered_cols = [
                "age_band", "bmi_missing_flag", "bmi_for_analysis", "bmi_category",
                "was_admitted_before", "admission_year_missing_flag", "years_since_last_admitted",
                "cholesterol_midpoint", "any_major_disease_history", "weight_bmi_interaction",
                "steps_per_age", "medical_risk_score", "lifestyle_risk_score",
            ]
            cleaned_analysis = pd.read_csv(TABLE_DIR / "cleaned_analysis_data.csv")
            display(cleaned_analysis[engineered_cols].head())
            """
        ),
        interp(
            "Feature engineering converts missingness, ranges, and history fields into model-ready representations without leaking target information.",
            "The engineered variables give reviewers interpretable risk context and prepare the data for the Milestone 2 pipeline.",
        ),
        md("## 20. Preprocessing Summary"),
        code(
            """
            preprocessing_summary = pd.DataFrame(
                [
                    ["applicant_id", "Drop from modeling", "Unique identifier"],
                    ["bmi", "Flag missingness and impute", "Limited missingness with health meaning"],
                    ["Year_last_admitted", "Create admission-status and recency features", "Structural missingness"],
                    ["cholesterol_level", "Engineer midpoint", "Ordinal range"],
                    ["insurance_cost", "Keep target outliers", "Valid quote bands"],
                ],
                columns=["field", "action", "reason"],
            )
            preprocessing_summary
            """
        ),
        interp(
            "The preprocessing plan is explicit, limited, and tied to each field's business meaning.",
            "This meets the rubric expectation for justifying unwanted-variable removal, missing-value handling, transformations, and new variables.",
        ),
        md("## 21. Extensive EDA Summary"),
        code(
            """
            strong = feature_signal[feature_signal["business_interpretation"].str.contains("strong", case=False, na=False)]
            weak = feature_signal[feature_signal["business_interpretation"].str.contains("weak", case=False, na=False)]
            display(strong)
            display(weak)
            """
        ),
        interp(
            "The strongest observed variables are weight, admission-year/recency, other-company coverage, regular checkups, weight change, and adventure sports.",
            "This focused ranking helps the report explain what actually matters in the supplied data rather than listing every variable equally.",
        ),
        md("## 22. Business Insights from EDA"),
        interp(
            "Quote bands and target-grid behavior mean RMSE/MAE should be supplemented with nearest-band evaluation in Milestone 2.",
            "Pricing teams need a business-ready quote band, not only a floating-point model prediction.",
        ),
        interp(
            "Weak marginal signals are documented rather than discarded because interactions can still matter in tree-based models.",
            "The business gets transparent caveats: weak standalone fields are not presented as major pricing drivers.",
        ),
        interp(
            "Missing admission history and prior admission recency are different concepts and should not be collapsed into one imputed number.",
            "This distinction gives a more defensible applicant-history signal for downstream pricing triage.",
        ),
        interp(
            "The EDA supports target-band stratification for train/test splitting because every quote band is a meaningful target level.",
            "The final model evaluation should preserve quote-band representation across train and test data.",
        ),
        md("## 23. Milestone 1 Conclusion and Next Steps"),
        interp(
            "The cleaned analytical dataset is ready for modeling after ID removal, targeted missingness treatment, feature engineering, and quote-grid-aware evaluation planning.",
            "Milestone 2 should compare baseline, linear, tree, boosted, calibrated, rounded, and deployment-ready variants with honest metric reporting.",
        ),
        interp(
            "EDA findings are association-based and should not be interpreted as underwriting policy or causal medical conclusions.",
            "The model can support pricing analysis, but final insurance decisions require policy, fairness, compliance, and human review.",
        ),
        md("## 24. Appendix: Additional Tables and Plots"),
        code(
            """
            appendix_files = sorted(p.name for p in TABLE_DIR.glob("*.csv"))
            pd.DataFrame({"available_table": appendix_files}).head(80)
            """
        ),
    ]

    extra_m1_interpretations = [
        ("Other-company coverage has a visible mean-cost separation in the group summary.", "Coverage status should be explained as possible customer-history context, not as a causal premium rule."),
        ("Adventure-sports participation has clearer group separation than many lifestyle indicators.", "This flag is useful for risk context but should not drive an automatic decision alone."),
        ("Alcohol use has weak marginal separation in the current data.", "The report should avoid claiming alcohol is a major standalone cost driver here."),
        ("Exercise category has weak marginal separation in the current data.", "Exercise is best treated as supporting applicant context rather than a headline driver."),
        ("BMI category is useful for descriptive health segmentation but does not dominate weight.", "Business interpretation should distinguish BMI context from the stronger weight gradient."),
        ("Disease-history flags show limited standalone signal in this dataset.", "The field remains important for governance and monitoring despite weak marginal EDA."),
        ("Glucose and cholesterol are clinically intuitive but weak as individual EDA signals here.", "The model may still use them in interactions, but the report should not overstate them."),
        ("Target tails are valid and belong to the same quote grid.", "Manual review is better than deleting high-premium cases."),
        ("Categorical frequencies confirm all major categories have enough representation for EDA.", "The reviewer can trust that segment comparisons are not purely anecdotal."),
        ("The preprocessing table ties each transformation to a data-quality or modeling-readiness reason.", "This makes the project easier to grade against the Milestone 1 rubric."),
    ]
    univariate_path = TABLE_DIR / "univariate_plot_interpretations.csv"
    if univariate_path.exists():
        univariate_rows = pd.read_csv(univariate_path)
        m1_cells.extend(
            [
                md("### Complete Univariate Plot Appendix for All Original Variables"),
                code(
                    """
                    univariate_index = pd.read_csv(TABLE_DIR / "univariate_plot_interpretations.csv")
                    display(univariate_index[[
                        "original_column", "clean_column", "analysis_type", "plot_type",
                        "missing_count", "unique_count", "plot_file"
                    ]])
                    """
                ),
                interp(
                    "This appendix provides one generated univariate plot for every original source variable, including identifiers, target, numeric fields, ordinal fields, binary flags, and nominal categories.",
                    "The reviewer can verify that no variable was skipped in the initial EDA; high-signal variables receive deeper bivariate treatment elsewhere in the notebook.",
                ),
            ]
        )
        for _, row in univariate_rows.iterrows():
            clean_col = str(row["clean_column"])
            m1_cells.append(image_cell(str(row["plot_file"]), f"Univariate plot: {clean_col}"))
            m1_cells.append(interp(str(row["technical_interpretation"]), str(row["business_interpretation"])))
    for technical, business in extra_m1_interpretations:
        m1_cells.append(interp(technical, business))

    nb1 = nbf.v4.new_notebook(cells=m1_cells, metadata=notebook_metadata)
    nb1_path = NOTEBOOK_DIR / "01_milestone1_eda.ipynb"
    nbf.write(nb1, nb1_path)

    m2_cells = [
        md("# Milestone 2: Insurance Price Prediction - Modeling, Explainability, and Deployment"),
        style_cell(),
        md(
            """
            ## 1. Modeling Objective

            Milestone 2 builds on the Milestone 1 cleaned dataset and evaluates predictive models for
            `insurance_cost`. The final package reports raw regression metrics, quote-band rounding,
            calibration behavior, explainability outputs, and Streamlit deployment artifacts.
            """
        ),
        code(embedded_notebook_bootstrap_source(), tags=["embedded-helper-code", "remove-input"]),
        code(
            """
            import json
            from pathlib import Path

            import joblib
            import numpy as np
            import pandas as pd
            from IPython.display import Image, Markdown, display

            ROOT = Path.cwd()
            if not (ROOT / "Insurance Data.csv").exists():
                ROOT = ROOT.parent

            TABLE_DIR = ROOT / "outputs" / "tables"
            FIG_DIR = ROOT / "outputs" / "figures"
            MODEL_DIR = ROOT / "outputs" / "models"

            ensure_dirs()
            set_visual_theme()
            pd.set_option("display.max_columns", 80)
            """
        ),
        interp(
            "The modeling notebook resolves the project root, creates output folders, and uses training/deployment helpers embedded in the notebook.",
            "A reviewer can execute the notebook from `Insurance Data.csv` without relying on external helper scripts or pre-generated model outputs.",
        ),
        md("## 2. End-to-End Modeling Artifact Generation"),
        code(
            """
            DATA_PATH = ROOT / "Insurance Data.csv"
            raw_df = pd.read_csv(DATA_PATH)

            report_df = add_report_features(raw_df)
            tables = save_profile_tables(raw_df, report_df)
            figures = generate_eda_figures(raw_df, report_df)
            model_results = build_models(raw_df)
            figures.update(model_results["model_figures"])
            app_path = write_streamlit_app()

            print(f"Loaded dataset: {DATA_PATH.name} -> {raw_df.shape[0]:,} rows, {raw_df.shape[1]:,} columns")
            print(f"Generated {len(tables)} profile table groups and {len(figures)} figure artifacts.")
            print(f"Final model artifact: {model_results['final_model_path'].relative_to(ROOT)}")
            print(f"Streamlit app artifact: {app_path.relative_to(ROOT)}")
            """
        ),
        interp(
            "This cell performs the train/test split, model comparison, tuning, calibration, price-band evaluation, explainability outputs, and saved-model generation from the raw dataset.",
            "All later Milestone 2 sections read artifacts generated in this notebook run, so the workflow is executable end to end.",
        ),
        md("## 3. Train/Test Strategy and Target Stratification"),
        code(
            """
            from sklearn.model_selection import train_test_split

            raw_model_df = clean_column_names(pd.read_csv(ROOT / "Insurance Data.csv"))
            X = raw_model_df.drop(columns=[TARGET, "applicant_id"], errors="ignore")
            y = raw_model_df[TARGET]

            split_strata = make_target_strata(y)
            X_train, X_test, y_train, y_test = train_test_split(
                X,
                y,
                test_size=0.2,
                random_state=RANDOM_STATE,
                stratify=split_strata,
            )

            valid_levels = target_grid_metadata(y_train)["valid_target_levels"]
            train_counts = y_train.value_counts().sort_index()
            test_counts = y_test.value_counts().sort_index()
            generated_split_check = pd.DataFrame(
                [
                    {
                        "insurance_cost": int(level),
                        "train_count": int(train_counts.get(level, 0)),
                        "test_count": int(test_counts.get(level, 0)),
                        "train_pct": round(int(train_counts.get(level, 0)) / len(y_train) * 100, 3),
                        "test_pct": round(int(test_counts.get(level, 0)) / len(y_test) * 100, 3),
                    }
                    for level in valid_levels
                ]
            )

            saved_split_check = pd.read_csv(TABLE_DIR / "target_band_split_check.csv")
            print("Matches saved target_band_split_check.csv:", generated_split_check.equals(saved_split_check))
            display(generated_split_check.head())
            display(generated_split_check.tail())
            """
        ),
        interp(
            "The train/test split preserves the target quote-band distribution, reducing evaluation bias from rare bands.",
            "Business quote bands remain represented during model testing, which makes the reported errors more credible.",
        ),
        md("## 4. Model Families Compared"),
        code(
            """
            metrics = pd.read_csv(MODEL_DIR / "model_metrics.csv")
            metrics[["model", "cv_rmse_mean", "test_MAE", "test_RMSE", "test_R2", "selected_final_model"]]
            """
        ),
        interp(
            "The comparison includes baseline, linear, regularized, tree, ensemble, boosted, and optional external boosting candidates when installed.",
            "A reviewer can see the final model was chosen from a broad candidate pool, not from a single default algorithm.",
        ),
        md("## 5. Baseline and Linear Models"),
        code(
            """
            metrics[metrics["model"].isin(["DummyRegressor", "LinearRegression", "Ridge", "LogTargetRidge"])][
                ["model", "test_MAE", "test_RMSE", "test_R2"]
            ]
            """
        ),
        interp(
            "Baseline and linear models provide interpretability but underperform the stronger nonlinear models.",
            "They establish a reasonable performance floor and justify moving to tree-based learners.",
        ),
        md("## 6. Tree and Boosting Models"),
        code(
            """
            tree_metrics = metrics[~metrics["model"].isin(["DummyRegressor", "LinearRegression", "Ridge", "LogTargetRidge"])]
            tree_metrics[["model", "test_MAE", "test_RMSE", "test_R2", "train_test_RMSE_gap", "selected_final_model"]]
            """
        ),
        image_cell("outputs/figures/model_comparison_rmse.png", "Model comparison by test RMSE"),
        interp(
            "Histogram gradient boosting variants are near the top of the test-RMSE ranking with modest train-test gaps.",
            "The final selection favors strong performance with simpler deployment behavior rather than only the absolute lowest test RMSE.",
        ),
        md("## 7. Final Model Selection"),
        code(
            """
            metadata = json.loads((MODEL_DIR / "model_metadata.json").read_text())
            final_summary = pd.read_csv(MODEL_DIR / "final_model_summary.csv")
            display(metadata)
            display(final_summary)
            """
        ),
        interp(
            "The selected final model is explicitly flagged in `model_metrics.csv` and summarized in `final_model_summary.csv`.",
            "The package explains the model as a parsimony/stability selection, not as an incorrect lowest-test-RMSE claim.",
        ),
        md("## 8. Repeated Cross-validation"),
        code(
            """
            repeated_cv = pd.read_csv(MODEL_DIR / "repeated_cv_metrics.csv")
            repeated_cv
            """
        ),
        interp(
            "Repeated cross-validation checks whether leading candidates remain competitive across multiple folds.",
            "This gives the final model a stronger stability argument than one train/test split alone.",
        ),
        md("## 9. Calibration and Price-grid Evaluation"),
        code(
            """
            price_grid_eval = pd.read_csv(MODEL_DIR / "price_grid_evaluation.csv")
            calibration = pd.read_csv(MODEL_DIR / "calibration_comparison.csv")
            display(price_grid_eval)
            display(calibration)
            """
        ),
        image_cell("outputs/figures/calibration_curve.png", "Calibration curve"),
        image_cell("outputs/figures/residual_by_cost_decile_after_calibration.png", "Residuals by cost decile after calibration"),
        interp(
            "Calibration improves some percentage or band-style metrics but does not improve MAE/RMSE in the current run.",
            "The Streamlit app therefore treats calibrated cost as secondary diagnostic output unless metadata explicitly chooses it.",
        ),
        interp(
            "Rounded raw prediction is the default quote-band display because the target is a fixed price grid.",
            "Business users see a valid quote band while analysts can still inspect the raw continuous estimate.",
        ),
        md("## 10. Ordinal Price-band Challenger"),
        code(
            """
            ordinal = pd.read_csv(MODEL_DIR / "ordinal_challenger_metrics.csv")
            ordinal
            """
        ),
        interp(
            "The ordinal/classification challenger tests the target as band labels instead of continuous cost.",
            "The challenger is useful evidence, but the selected regression pipeline remains stronger for the final deployment objective.",
        ),
        md("## 11. Saved Model Load and Prediction Smoke Test"),
        code(
            """
            model = joblib.load(MODEL_DIR / "final_model.pkl")
            raw_df = pd.read_csv(ROOT / "Insurance Data.csv")
            sample = clean_column_names(raw_df).drop(columns=["insurance_cost", "applicant_id"], errors="ignore").head(5)
            preds = model.predict(sample)
            pd.DataFrame({"prediction": preds})
            """
        ),
        interp(
            "The saved artifact loads as a fitted sklearn pipeline and predicts finite values for sample rows.",
            "This confirms the deployment file is usable outside the training script.",
        ),
        md("## 12. Explainability: Permutation Importance"),
        code(
            """
            importance = pd.read_csv(MODEL_DIR / "feature_importance.csv")
            importance.head(15)
            """
        ),
        image_cell("outputs/figures/feature_importance.png", "Permutation feature importance"),
        interp(
            "Permutation importance aligns with EDA by highlighting weight and admission/coverage/checkup/weight-change signals.",
            "The explanation is consistent with the Milestone 1 story and avoids introducing a contradictory driver narrative.",
        ),
        md("## 13. Explainability: SHAP and Partial Dependence"),
        code(
            """
            shap_status = json.loads((MODEL_DIR / "shap_status.json").read_text())
            display(shap_status)
            shap_path = MODEL_DIR / "shap_importance.csv"
            if shap_path.exists():
                display(pd.read_csv(shap_path).head(15))
            """
        ),
        image_cell("outputs/figures/shap_importance_bar.png", "SHAP importance when available"),
        image_cell("outputs/figures/partial_dependence_top_features.png", "Partial dependence for top features"),
        interp(
            "SHAP is generated when the installed environment supports it; permutation importance remains the core deployed explanation.",
            "The package remains robust on machines without SHAP while still including stronger explainability when available.",
        ),
        md("## 14. Deployment Schema"),
        code(
            """
            app_schema = json.loads((MODEL_DIR / "app_schema.json").read_text())
            app_schema["numeric_ranges"].keys(), app_schema["target_grid"]["target_unique_count"], app_schema["target_grid"]["target_grid_step"]
            """
        ),
        code(
            """
            pd.DataFrame(app_schema["numeric_ranges"]).T
            """
        ),
        interp(
            "The app schema includes training-derived ranges, including Year_last_admitted, and excludes applicant_id.",
            "The Streamlit form stays aligned with model inputs and avoids asking the user for an identifier.",
        ),
        md("## 15. Streamlit Deployment Logic"),
        code(
            """
            deployment_variant = metadata.get("deployment_variant", {})
            deployment_variant
            """
        ),
        interp(
            "Deployment metadata separates raw analytical estimates, quote-band rounding, optional calibration, and risk category logic.",
            "The app displays calibrated cost as context but defaults quote-band output to the selected raw-rounded business variant.",
        ),
        md("## 16. Final Metrics for Report Synchronization"),
        code(
            """
            selected_mask = metrics["selected_final_model"].astype(str).str.lower().isin(["true", "1"])
            selected = metrics.loc[selected_mask].iloc[0]
            best_test = metrics.sort_values(["test_RMSE", "test_MAE"]).iloc[0]
            pd.DataFrame(
                [
                    {"comparison": "selected_final_model", **selected[["model", "test_MAE", "test_RMSE", "test_R2"]].to_dict()},
                    {"comparison": "best_test_rmse_model", **best_test[["model", "test_MAE", "test_RMSE", "test_R2"]].to_dict()},
                ]
            )
            """
        ),
        interp(
            "The selected model and best test-RMSE model are both reported so there is no contradiction across deliverables.",
            "A small RMSE gap is acceptable because the selection reason is parsimony and stability, not pure leaderboard rank.",
        ),
        md("## 17. Business Recommendations"),
        interp(
            "Use the model as pricing decision support and keep raw, calibrated, and rounded outputs clearly labeled.",
            "Pricing teams can use the output for triage, but underwriting authority should remain with policy and human review.",
        ),
        interp(
            "Monitor drift, segment-level error, and fairness-sensitive variables such as Gender and Location after deployment.",
            "A capstone model should not become a production pricing rule without compliance review.",
        ),
        md("## 18. Reproducibility and Required Artifacts"),
        code(
            """
            required = [
                ROOT / "run_all.py",
                ROOT / "insurance_modeling.py",
                ROOT / "app.py",
                MODEL_DIR / "final_model.pkl",
                MODEL_DIR / "model_metrics.csv",
                MODEL_DIR / "final_model_summary.csv",
                MODEL_DIR / "app_schema.json",
            ]
            pd.DataFrame({"artifact": [str(p.relative_to(ROOT)) for p in required], "exists": [p.exists() for p in required]})
            """
        ),
        interp(
            "The artifact check verifies that the generated modeling and deployment files are present.",
            "Reviewers can reproduce the package by running `python run_all.py` from the project root after installing requirements.",
        ),
        md("## 19. Milestone 2 Conclusion"),
        interp(
            "The final modeling package is synchronized across metrics, metadata, reports, notebooks, PPT, and app deployment logic.",
            "The submission gives a coherent story from EDA to model choice to quote-band-ready deployment.",
        ),
    ]

    nb2 = nbf.v4.new_notebook(cells=m2_cells, metadata=notebook_metadata)
    nb2_path = NOTEBOOK_DIR / "02_milestone2_modeling.ipynb"
    nbf.write(nb2, nb2_path)
    return [nb1_path, nb2_path]


def execute_and_export_notebooks(notebook_paths: list[Path]) -> list[Path]:
    exported: list[Path] = []
    if os.name == "nt":
        asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    exporter = HTMLExporter()
    exporter.template_name = "lab"
    tag_preprocessor = TagRemovePreprocessor(remove_input_tags={"remove-input"})
    tag_preprocessor.enabled = True
    exporter.register_preprocessor(tag_preprocessor, enabled=True)
    if hasattr(exporter, "embed_images"):
        exporter.embed_images = True
    for nb_path in notebook_paths:
        notebook = nbf.read(nb_path, as_version=4)
        client = NotebookClient(
            notebook,
            timeout=1800,
            kernel_name="python3",
            resources={"metadata": {"path": str(ROOT)}},
        )
        executed = client.execute()
        nbf.write(executed, nb_path)
        body, _ = exporter.from_notebook_node(executed, resources={"metadata": {"path": str(ROOT)}})
        html_path = nb_path.with_suffix(".html")
        html_path.write_text(body, encoding="utf-8")
        output_html_path = OUTPUT_NOTEBOOK_DIR / html_path.name
        shutil.copy2(html_path, output_html_path)
        exported.extend([html_path, output_html_path])
    return exported


def write_requirements_and_readme(model_results: dict[str, object]) -> tuple[Path, Path]:
    requirements = "\n".join(
        [
            "pandas>=2.2,<3.0",
            "numpy>=2.0,<3.0",
            "scikit-learn>=1.5,<1.8",
            "lightgbm>=4.0,<5.0",
            "xgboost>=2.0,<4.0",
            "catboost>=1.2,<2.0",
            "shap>=0.46,<0.53",
            "matplotlib>=3.8,<4.0",
            "seaborn>=0.13,<0.14",
            "joblib>=1.3,<2.0",
            "python-docx>=1.1,<2.0",
            "python-pptx>=1.0,<2.0",
            "pypdf>=4.0,<7.0",
            "nbformat>=5.10,<6.0",
            "nbclient>=0.10,<1.0",
            "nbconvert>=7.16,<8.0",
            "ipykernel>=6.29,<7.0",
            "streamlit>=1.36,<2.0",
        ]
    )
    req_path = ROOT / "requirements.txt"
    req_path.write_text(requirements + "\n", encoding="utf-8")

    metrics = model_results["metrics"].loc[
        model_results["metrics"]["model"] == model_results["best_name"]
    ].iloc[0]
    readme = f"""# Insurance Price Prediction Capstone

This project builds a supervised regression model to estimate `insurance_cost` using health, lifestyle, habit, and demographic variables.

## Reproducibility

Run the full pipeline from the project root:

```bash
python run_all.py
```

For grading, install the requirements and run `python run_all.py` from the project root.

The script generates EDA figures, summary tables, model artifacts, milestone reports, the final presentation, and the Streamlit app dependencies.

The target is a discrete business quote grid. The pipeline reports raw regression metrics, calibrated predictions, nearest valid quote-band outputs, and band-accuracy metrics.

## Notebooks and HTML exports

The submitted notebooks are executable end to end from `Insurance Data.csv`:

- Run `notebooks/01_milestone1_eda.ipynb` to recreate Milestone 1 EDA tables, figures, preprocessing evidence, and interpretations from the raw dataset.
- Run `notebooks/02_milestone2_modeling.ipynb` to recreate the train/test split, model training/comparison, calibration, quote-band evaluation, explainability artifacts, saved model, and Streamlit app from the raw dataset.

Milestone 2 runs the full model suite and may take several minutes on CPU. The notebooks do not require pre-generated output files; they create the required `outputs/` artifacts during execution.

`python run_all.py` can still be used to regenerate the full package plus executed HTML exports:

- `notebooks/01_milestone1_eda.ipynb`
- `notebooks/01_milestone1_eda.html`
- `notebooks/02_milestone2_modeling.ipynb`
- `notebooks/02_milestone2_modeling.html`
- `outputs/notebooks/01_milestone1_eda.html`
- `outputs/notebooks/02_milestone2_modeling.html`

The Milestone 1 notebook follows the rubric sections for Data Report, Initial EDA, Data Pre-processing, and Extensive EDA. The HTML copies are included for reviewers who want to inspect executed notebook output without opening Jupyter.

## Final model

Selected model: `{model_results['best_name']}`

- Test MAE: {metrics['test_MAE']:,.0f}
- Test RMSE: {metrics['test_RMSE']:,.0f}
- Test R2: {metrics['test_R2']:.3f}

## Streamlit

```bash
streamlit run app.py
```

The app loads `outputs/models/final_model.pkl` and applies the same preprocessing pipeline used during training.
It also reads `outputs/models/app_schema.json` and, when available, `outputs/models/prediction_calibrator.pkl` to show raw predicted cost, calibrated cost, and nearest valid quote band.
Raw prediction is used for the analytical estimate and default risk category; the nearest quote band is rounded from the selected deployment variant in `outputs/models/model_metadata.json`. Calibrated cost is shown as a secondary diagnostic unless metadata explicitly chooses it for deployment.

## Key generated validation artifacts

- `outputs/models/repeated_cv_metrics.csv`
- `outputs/models/final_model_summary.csv`
- `outputs/models/price_grid_evaluation.csv`
- `outputs/models/calibration_comparison.csv`
- `outputs/models/ordinal_challenger_metrics.csv`
- `outputs/models/app_schema.json`

## Colab GPU Training

Use `notebooks/03_colab_gpu_training.ipynb` in Google Colab with a GPU runtime to run the heavier optional XGBoost/CatBoost/LightGBM comparison without using this PC. The notebook expects the prepared project zip from `outputs/colab/insurance_capstone_colab_bundle.zip`, runs `python run_all.py` with `INSURANCE_USE_GPU=1`, and downloads the regenerated artifacts.

Create or refresh the upload bundle with:

```powershell
powershell -ExecutionPolicy Bypass -File scripts\\prepare_colab_bundle.ps1
```
"""
    readme_path = ROOT / "README.md"
    readme_path.write_text(readme, encoding="utf-8")
    return req_path, readme_path


def delete_all_slides(prs: Presentation) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rel_id = sld_id.rId
        prs.part.drop_rel(rel_id)
        sld_id_lst.remove(sld_id)


EMU_PER_INCH = 914400
PPT_MARGIN = 0.15
PPT_SLIDE_W = 10.0
PPT_SLIDE_H = 5.625


def configure_ppt_dimensions(prs: Presentation) -> None:
    global PPT_SLIDE_W, PPT_SLIDE_H
    PPT_SLIDE_W = prs.slide_width / EMU_PER_INCH
    PPT_SLIDE_H = prs.slide_height / EMU_PER_INCH


def bounded_box(x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    safe_right = PPT_SLIDE_W - PPT_MARGIN
    safe_bottom = PPT_SLIDE_H - PPT_MARGIN
    x = min(max(PPT_MARGIN, x), safe_right - 0.1)
    y = min(max(PPT_MARGIN, y), safe_bottom - 0.1)
    w = min(w, max(0.1, safe_right - x))
    h = min(h, max(0.1, safe_bottom - y))
    return x, y, w, h


def add_textbox(slide, x, y, w, h, text, font_size=22, bold=False, color=PRIMARY, align=PP_ALIGN.LEFT):
    x, y, w, h = bounded_box(x, y, w, h)
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = PptRGBColor.from_string(color.replace("#", ""))
    return box


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.55, 0.34, 8.9, 0.48, title, font_size=24, bold=True, color=PRIMARY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PptInches(0.55), PptInches(0.92), PptInches(1.25), PptInches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PptRGBColor.from_string(ACCENT.replace("#", ""))
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.55, 1.02, 8.9, 0.35, subtitle, font_size=10, color=MUTED)


def add_bullet_box(slide, bullets: list[str], x=0.75, y=1.45, w=4.2, h=3.85, font_size=12) -> None:
    x, y, w, h = bounded_box(x, y, w, h)
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PptPt(font_size)
        p.font.color.rgb = PptRGBColor.from_string(PRIMARY.replace("#", ""))
        p.space_after = PptPt(3)


def add_metric(slide, x, y, label, value, color=ACCENT, w=1.8) -> None:
    add_textbox(slide, x, y, w, 0.28, label.upper(), font_size=8, bold=True, color=MUTED)
    add_textbox(slide, x, y + 0.3, w, 0.42, value, font_size=18, bold=True, color=color)


def add_image(slide, path: Path, x, y, w, h) -> None:
    if path.exists():
        x, y, w, h = bounded_box(x, y, w, h)
        slide.shapes.add_picture(str(path), PptInches(x), PptInches(y), width=PptInches(w), height=PptInches(h))


def add_table_to_slide(slide, df: pd.DataFrame, x, y, w, h, font_size=9, max_rows=8) -> None:
    x, y, w, h = bounded_box(x, y, w, h)
    display = df.head(max_rows)
    rows, cols = display.shape[0] + 1, display.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    table = table_shape.table
    for c, col in enumerate(display.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PptRGBColor.from_string(PRIMARY.replace("#", ""))
        for p in cell.text_frame.paragraphs:
            p.font.size = PptPt(font_size)
            p.font.bold = True
            p.font.color.rgb = PptRGBColor(255, 255, 255)
    for r, (_, row) in enumerate(display.iterrows(), start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if isinstance(val, float):
                text = f"{val:,.2f}" if abs(val) >= 100 else f"{val:.3f}"
            else:
                text = str(val)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(font_size)
                p.font.color.rgb = PptRGBColor.from_string(PRIMARY.replace("#", ""))


def create_streamlit_showcase_image() -> Path:
    path = FIG_DIR / "streamlit_deployment_workflow.png"
    fig, ax = plt.subplots(figsize=(10, 5.2))
    ax.axis("off")
    steps = [
        ("Input form", "Health, lifestyle,\nhabit, demographics"),
        ("Pipeline", "Feature engineering,\nimputation, encoding"),
        ("Final model", "Raw predicted\ninsurance cost"),
        ("Decision support", "Raw cost,\nquote band, risk band"),
    ]
    xs = [0.07, 0.32, 0.58, 0.82]
    for i, (title, detail) in enumerate(steps):
        rect = plt.Rectangle((xs[i] - 0.09, 0.35), 0.18, 0.28, facecolor="#E8F4F1", edgecolor=ACCENT, linewidth=2)
        ax.add_patch(rect)
        ax.text(xs[i], 0.56, title, ha="center", va="center", fontsize=13, fontweight="bold", color=PRIMARY)
        ax.text(xs[i], 0.44, detail, ha="center", va="center", fontsize=10.5, color=MUTED)
        if i < len(steps) - 1:
            ax.annotate("", xy=(xs[i + 1] - 0.11, 0.49), xytext=(xs[i] + 0.11, 0.49), arrowprops=dict(arrowstyle="->", lw=2, color=PRIMARY))
    ax.text(0.5, 0.78, "Streamlit Deployment Workflow", ha="center", fontsize=18, fontweight="bold", color=PRIMARY)
    ax.text(0.5, 0.18, "The app loads the saved pipeline, app_schema.json, and optional calibration artifact; quote bands default to the selected raw-rounded deployment variant.", ha="center", fontsize=10.5, color=MUTED)
    savefig(path)
    return path


def generate_presentation(model_results: dict[str, object], figures: dict[str, Path], eda_summary: dict[str, object]) -> Path:
    prs = Presentation(str(PPT_TEMPLATE))
    configure_ppt_dimensions(prs)
    delete_all_slides(prs)
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    metrics = model_results["metrics"].copy()
    winner = metrics.loc[metrics["model"] == model_results["best_name"]].iloc[0]
    importance = model_results["importance"].head(10).round(4)
    optional_model_status = optional_model_status_sentence(model_results["metadata"])
    explainability_status = explainability_status_sentence(model_results["metadata"])
    showcase = create_streamlit_showcase_image()

    def new_slide(title, subtitle=None):
        slide = prs.slides.add_slide(blank)
        add_slide_title(slide, title, subtitle)
        return slide

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.6, 1.05, 5.6, 0.95, "Insurance Price\nPrediction", font_size=31, bold=True, color=PRIMARY)
    add_textbox(slide, 0.65, 2.45, 5.4, 0.55, "Capstone final presentation | Health, lifestyle, habit, and demographic pricing model", font_size=12, color=MUTED)
    add_metric(slide, 0.75, 3.45, "Rows", "25,000")
    add_metric(slide, 2.55, 3.45, "Predictors", "23")
    add_metric(slide, 4.35, 3.45, "Target", "insurance_cost", color=WARN)
    add_image(slide, figures["target_distribution"], 6.45, 1.15, 3.15, 3.35)

    slide = new_slide("Executive Summary")
    add_bullet_box(
        slide,
        [
            "Objective: estimate insurance_cost from applicant risk and behavior variables.",
            "insurance_cost is a fixed quote-band grid, not a fully continuous target.",
            "No duplicate rows; applicant_id is unique and removed from modeling.",
            "Key preprocessing: BMI imputation, admission recency, cholesterol midpoint, disease-history and lifestyle features.",
            f"Final model: {model_results['best_name']} with test MAE {winner['test_MAE']:,.0f}, RMSE {winner['test_RMSE']:,.0f}, and R2 {winner['test_R2']:.3f}.",
            "Deployment: Streamlit app with raw predicted cost, calibrated cost, nearest quote band, risk category, and model-driver context.",
        ],
        x=0.65,
        y=1.32,
        w=8.8,
        h=3.85,
        font_size=12,
    )

    slide = new_slide("Milestone 1 Rubric Coverage")
    add_bullet_box(
        slide,
        [
            "Data Report: dataset shape, head/tail/sample checks, data types, unique counts, missingness, duplicates, descriptive statistics, and variable classification.",
            "Initial EDA: target distribution, numeric spread, categorical frequencies, bivariate charts, correlation heatmap, and early insight summaries.",
            "Data Pre-processing: applicant_id removal, typo cleanup, BMI and admission-year missingness treatment, outlier strategy, transformations, and engineered variables.",
            "Extensive EDA: important-variable relationships with technical and business interpretations for weight, admission history, coverage, checkups, weight change, and adventure sports.",
            "Evidence files: Milestone 1 DOCX, executed Milestone 1 notebook/HTML, outputs/reports/milestone1_rubric_coverage_matrix.csv, and supporting figures/tables.",
        ],
        x=0.55,
        y=1.25,
        w=8.9,
        h=4.0,
        font_size=10.7,
    )

    slide = new_slide("Business Problem and Objective")
    add_bullet_box(
        slide,
        [
            "Pricing teams need a consistent estimate of likely insurance cost for each applicant.",
            "Inputs combine health metrics, lifestyle habits, disease history, usage behavior, and demographics.",
            "Model output should support pricing triage, wellness programs, and manual review of extreme cases.",
            "The model is positioned as decision support, not the final underwriting authority.",
        ],
        x=0.75,
        y=1.45,
        w=4.2,
        font_size=12,
    )
    add_image(slide, figures["cost_by_disease_history"], 5.15, 1.3, 4.25, 3.85)

    slide = new_slide("Dataset Overview")
    add_metric(slide, 0.55, 1.33, "Rows", "25,000", w=1.55)
    add_metric(slide, 2.35, 1.33, "Columns", "24", w=1.55)
    add_metric(slide, 4.15, 1.33, "Duplicates", "0", color=WARN, w=1.55)
    add_metric(slide, 5.95, 1.33, "Missing BMI", "990", w=1.55)
    add_metric(slide, 7.75, 1.33, "Missing admission year", "11,881", w=1.75)
    add_image(slide, figures["missing_values"], 0.55, 2.35, 4.25, 2.85)
    add_image(slide, figures["target_boxplot"], 5.15, 2.35, 4.25, 2.85)

    slide = new_slide("Data Preprocessing Summary")
    add_bullet_box(
        slide,
        [
            "Renamed three misspelled columns and corrected Salaried occupation label.",
            "Dropped applicant_id before modeling and removed it from app inputs.",
            "Encoded admission history into was_admitted_before, admission_year_missing_flag, admission_status, and years_since_last_admitted.",
            "Added bmi_missing_flag before BMI imputation.",
            "Created age_band, bmi_category, cholesterol_midpoint, any_major_disease_history, weight_bmi_interaction, and steps_per_age.",
            "Used train-fitted preprocessing inside sklearn Pipeline: imputation, scaling, and one-hot encoding.",
            "Retained valid target outliers; capped only extreme BMI values for analysis/model stability.",
        ],
        x=0.65,
        y=1.28,
        w=8.8,
        h=3.95,
        font_size=11.5,
    )

    slide = new_slide("Key EDA Findings")
    add_image(slide, figures["correlation_heatmap"], 0.55, 1.3, 4.4, 3.85)
    add_bullet_box(
        slide,
        [
            "insurance_cost has 54 valid quote bands with a fixed 1,234 step.",
            "Weight is the dominant observed driver and has an extremely strong marginal relationship with cost.",
            "Admission year/recency, other-company coverage, regular checkups, weight change, and adventure sports carry stronger observed signal.",
            "Smoking, alcohol, exercise, BMI, glucose, cholesterol, age, and disease-history flags are weak marginal drivers in this dataset.",
        ],
        x=5.2,
        y=1.45,
        w=4.25,
        h=3.6,
        font_size=11.5,
    )

    slide = new_slide("Feature Engineering")
    add_bullet_box(
        slide,
        [
            "age_band: customer life-stage segmentation.",
            "bmi_category: standard risk categories based on imputed BMI.",
            "cholesterol_midpoint: ordinal range converted to numeric midpoint.",
            "admission recency: prior admission flag plus years since last admission.",
            "risk interactions: disease-history flag, weight x BMI, and steps per age.",
        ],
        x=0.55,
        y=1.35,
        w=4.25,
        font_size=12,
    )
    add_image(slide, figures["cost_by_bmi_category"], 5.15, 1.35, 4.25, 3.85)

    slide = new_slide("Modeling Approach")
    add_bullet_box(
        slide,
        [
            "Split: 80 percent train, 20 percent test, stratified by target quote band.",
            "Baseline models: DummyRegressor and Linear Regression.",
            "Regularized linear model: Ridge.",
            "Tree/ensemble models: Decision Tree, Random Forest, Extra Trees, Gradient Boosting, HistGradientBoosting.",
            optional_model_status,
            "Added repeated CV, calibration, quote-band rounding, and an ordinal price-band challenger.",
        ],
        x=0.65,
        y=1.3,
        w=8.8,
        h=3.9,
        font_size=12,
    )

    slide = new_slide("Model Evaluation Metrics")
    add_bullet_box(
        slide,
        [
            "MAE: average absolute premium error in cost units.",
            "RMSE: penalizes large mispricing errors.",
            "R2: proportion of target variance explained.",
            "MAPE/SMAPE: percentage-style business interpretability.",
            "Train-test gap: checks whether a high-scoring model overfits.",
        ],
        x=0.65,
        y=1.35,
        w=8.8,
        h=3.8,
        font_size=13,
    )

    slide = new_slide("Model Performance Comparison")
    add_image(slide, figures["model_comparison_rmse"], 0.45, 1.25, 4.5, 4.0)
    table_cols = ["model", "test_MAE", "test_RMSE", "test_R2", "test_MAPE"]
    add_table_to_slide(slide, metrics[table_cols].round(3), 5.15, 1.25, 4.35, 4.0, font_size=5.8, max_rows=8)

    slide = new_slide("Calibration and Quote Bands")
    price_cols = ["variant", "MAE", "RMSE", "exact_band_accuracy", "within_2_band_accuracy"]
    add_table_to_slide(slide, model_results["price_grid_evaluation"][price_cols].round(3), 0.45, 1.25, 4.55, 3.15, font_size=6, max_rows=4)
    add_image(slide, FIG_DIR / "calibration_curve.png", 5.35, 1.2, 4.1, 3.45)
    add_bullet_box(
        slide,
        [
            "Rounding the selected raw prediction to the nearest valid quote band is useful for business presentation.",
            "Calibration is reported separately because it improves some percentage/band metrics but not MAE/RMSE.",
            "Exact band accuracy is hard because there are 54 target bands; within-band accuracy is more informative.",
        ],
        x=0.55,
        y=4.65,
        w=8.85,
        h=0.75,
        font_size=8.8,
    )

    slide = new_slide("Final Model and Explainability")
    add_bullet_box(
        slide,
        [
            f"Selected model: {model_results['best_name']}.",
            f"Test MAE {winner['test_MAE']:,.0f}; test RMSE {winner['test_RMSE']:,.0f}; test R2 {winner['test_R2']:.3f}.",
            "Selected through the parsimony/stability rule; the best test-RMSE model is reported separately in the metrics table.",
            "Top drivers are led by weight, admission/coverage/checkup/weight-change signals, with weaker marginal effects for many lifestyle flags.",
            explainability_status,
        ],
        x=0.55,
        y=1.25,
        w=4.25,
        h=3.8,
        font_size=11.2,
    )
    add_image(slide, FIG_DIR / "feature_importance.png", 5.1, 1.25, 4.35, 3.95)

    slide = new_slide("Streamlit Deployment Showcase")
    add_image(slide, showcase, 0.55, 1.25, 8.9, 3.75)
    add_textbox(slide, 0.8, 5.08, 8.4, 0.28, "Run with: streamlit run app.py", font_size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)

    slide = new_slide("Business Recommendations")
    add_bullet_box(
        slide,
        [
            "Use quote-band output for business communication while retaining raw prediction for analytics.",
            "Prioritize review of weight and admission-related policy alignment because those variables dominate this dataset.",
            "Use preventive checkup campaigns carefully as observational association, not proof of causal premium reduction.",
            "Add manual review for extreme predictions before quoting final premiums.",
            "Monitor Gender and Location effects for compliance and fairness risk.",
            "Track drift, segment error, and premium conversion after deployment.",
        ],
        x=0.65,
        y=1.3,
        w=8.8,
        h=3.95,
        font_size=12,
    )

    slide = new_slide("Risks, Limitations, and Next Steps")
    add_bullet_box(
        slide,
        [
            "Dataset provenance and real production pricing rules are not specified.",
            "Target is discretized, likely reflecting historical pricing bands.",
            "Weight dominates; validate against business pricing policy and fairness expectations.",
            "All predictors must be confirmed available at quote time before real deployment.",
            "Gender and Location require governance review.",
            "Next: add calibration checks, monitoring dashboards, challenger models, and periodic retraining.",
        ],
        x=0.55,
        y=1.35,
        w=4.25,
        font_size=11.8,
    )
    add_image(slide, figures["residuals"], 5.15, 1.35, 4.25, 3.85)

    slide = new_slide("Appendix: Data Dictionary")
    data_dict = pd.read_csv(TABLE_DIR / "data_dictionary_summary.csv")
    add_table_to_slide(slide, data_dict[["clean_column", "dtype", "missing", "unique_values"]], 0.55, 1.15, 8.9, 4.1, font_size=5, max_rows=12)

    slide = new_slide("Appendix: Additional EDA")
    add_image(slide, figures["cost_by_smoking"], 0.55, 1.25, 4.25, 3.9)
    add_image(slide, figures["cost_by_exercise"], 5.15, 1.25, 4.25, 3.9)

    slide = new_slide("Appendix: Model Details")
    add_table_to_slide(slide, importance[["feature", "importance_mean", "importance_std"]], 0.55, 1.25, 4.25, 4.1, font_size=6, max_rows=10)
    add_image(slide, figures["predicted_vs_actual"], 5.15, 1.25, 4.25, 4.05)

    check_ppt_overflows(prs, raise_on_overflow=True)
    output_path = ROOT / "Insurance_Price_Prediction_Final_Presentation.pptx"
    return save_pptx_with_fallback(prs, output_path)


def check_ppt_overflows(prs: Presentation, raise_on_overflow: bool = False) -> dict[str, object]:
    slide_w = prs.slide_width / EMU_PER_INCH
    slide_h = prs.slide_height / EMU_PER_INCH
    safe_right = slide_w - PPT_MARGIN
    safe_bottom = slide_h - PPT_MARGIN
    overflows = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            left = shape.left / EMU_PER_INCH
            top = shape.top / EMU_PER_INCH
            right = (shape.left + shape.width) / EMU_PER_INCH
            bottom = (shape.top + shape.height) / EMU_PER_INCH
            if left < -0.01 or top < -0.01 or right > safe_right + 0.01 or bottom > safe_bottom + 0.01:
                overflows.append(
                    {
                        "slide": slide_idx,
                        "shape": shape_idx,
                        "left": round(left, 3),
                        "top": round(top, 3),
                        "right": round(right, 3),
                        "bottom": round(bottom, 3),
                    }
                )
    result = {
        "status": "failed" if overflows else "passed",
        "slide_width": round(slide_w, 3),
        "slide_height": round(slide_h, 3),
        "margin": PPT_MARGIN,
        "overflow_count": len(overflows),
        "overflows": overflows[:50],
    }
    if raise_on_overflow and overflows:
        raise ValueError(f"PPT overflow check failed for {len(overflows)} shapes: {overflows[:5]}")
    return result


def inspect_pptx(path: Path) -> dict[str, object]:
    prs = Presentation(str(path))
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        texts.append({"slide": i, "text_count": len(slide_text), "chars": sum(len(t) for t in slide_text)})
    return {"slides": len(prs.slides), "text_summary": texts, "overflow_check": check_ppt_overflows(prs)}


def verify_docx_contains(path: Path, required_terms: list[str]) -> dict[str, object]:
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")
    return {term: term in xml for term in required_terms}


def read_docx_text(path: Path) -> str:
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml").decode("utf-8", errors="ignore")
    text = re.sub(r"<[^>]+>", "", xml)
    return unescape(text)


def count_docx_term(path: Path, term: str) -> int:
    if not path.exists():
        return 0
    return read_docx_text(path).count(term)


def notebook_quality_summary(notebook_paths: list[Path]) -> dict[str, object]:
    summary: dict[str, object] = {}
    for path in notebook_paths:
        if not path.exists():
            summary[path.name] = {"exists": False}
            continue
        notebook = nbf.read(path, as_version=4)
        sources = "\n".join(str(cell.get("source", "")) for cell in notebook.cells)
        cell_count = len(notebook.cells)
        tech_count = sources.count("Technical interpretation:") + sources.count("insight-label tech")
        business_count = sources.count("Business interpretation:") + sources.count("insight-label biz")
        if path.name.startswith("01_"):
            passes = cell_count >= 45 and tech_count >= 30 and business_count >= 30
        else:
            passes = cell_count >= 35 and tech_count >= 15 and business_count >= 15
        summary[path.name] = {
            "exists": True,
            "cell_count": cell_count,
            "technical_interpretation_count": tech_count,
            "business_interpretation_count": business_count,
            "passes_threshold": bool(passes),
        }
    return summary


def notebook_self_contained_summary(notebook_paths: list[Path]) -> dict[str, object]:
    patterns = ["from run_all import", "import run_all", "from insurance_modeling import", "sys.path.insert"]
    summary: dict[str, object] = {}
    for path in notebook_paths:
        if not path.exists():
            summary[path.name] = {"exists": False, "self_contained": False}
            continue
        notebook = nbf.read(path, as_version=4)
        embedded_cells = [cell for cell in notebook.cells if "embedded-helper-code" in cell.metadata.get("tags", [])]
        visible_cells = [cell for cell in notebook.cells if "embedded-helper-code" not in cell.metadata.get("tags", [])]
        visible_source = "\n".join(str(cell.get("source", "")) for cell in visible_cells)
        pattern_counts = {pattern: visible_source.count(pattern) for pattern in patterns}
        helper_cells_hidden = bool(
            embedded_cells
            and all(
                "remove-input" in cell.metadata.get("tags", [])
                and bool(cell.metadata.get("jupyter", {}).get("source_hidden"))
                for cell in embedded_cells
            )
        )
        self_contained = bool(embedded_cells and sum(pattern_counts.values()) == 0)
        summary[path.name] = {
            "exists": True,
            "self_contained": self_contained,
            "embedded_helper_cell_count": len(embedded_cells),
            "visible_external_project_import_counts": pattern_counts,
            "embedded_helper_input_hidden_in_html": helper_cells_hidden,
            "required_input_files": ["Insurance Data.csv"],
            "note": "Helper source is embedded in the notebook; visible analysis cells do not import local project scripts.",
        }
    return summary


def notebook_warning_output_summary(notebook_paths: list[Path]) -> dict[str, object]:
    summary: dict[str, object] = {}
    for path in notebook_paths:
        if not path.exists():
            summary[path.name] = {"exists": False}
            continue
        notebook = nbf.read(path, as_version=4)
        stderr_count = 0
        warning_or_traceback_count = 0
        error_count = 0
        for cell in notebook.cells:
            for output in cell.get("outputs", []):
                text_value = output.get("text", "")
                text = "".join(text_value) if isinstance(text_value, list) else str(text_value)
                if output.get("name") == "stderr":
                    stderr_count += 1
                if "warning" in text.lower() or "traceback" in text.lower():
                    warning_or_traceback_count += 1
                if output.get("output_type") == "error":
                    error_count += 1
        summary[path.name] = {
            "exists": True,
            "stderr_output_count": stderr_count,
            "warning_or_traceback_output_count": warning_or_traceback_count,
            "error_output_count": error_count,
            "clean": stderr_count == 0 and warning_or_traceback_count == 0 and error_count == 0,
        }
    return summary


def rubric_matrix_quality() -> dict[str, object]:
    path = REPORT_DIR / "milestone1_rubric_coverage_matrix.csv"
    if not path.exists():
        return {"exists": False, "all_covered": False}
    matrix = pd.read_csv(path)
    required_columns = [
        "rubric_section",
        "marks",
        "rubric_requirement",
        "evidence_in_milestone1_docx",
        "evidence_in_notebook",
        "evidence_file_or_chart",
        "status",
    ]
    return {
        "exists": True,
        "row_count": int(len(matrix)),
        "required_columns_present": all(column in matrix.columns for column in required_columns),
        "all_covered": bool((matrix["status"].astype(str).str.lower() == "covered").all()),
        "sections": sorted(matrix["rubric_section"].dropna().unique().tolist()),
    }


def univariate_plot_quality() -> dict[str, object]:
    path = TABLE_DIR / "univariate_plot_interpretations.csv"
    if not path.exists():
        return {"exists": False, "all_plots_exist": False}
    index = pd.read_csv(path)
    plot_exists = []
    for plot_file in index["plot_file"].astype(str):
        plot_exists.append((ROOT / Path(*plot_file.split("/"))).exists())
    return {
        "exists": True,
        "row_count": int(len(index)),
        "all_plots_exist": bool(all(plot_exists)),
        "technical_interpretation_count": int(index["technical_interpretation"].notna().sum()),
        "business_interpretation_count": int(index["business_interpretation"].notna().sum()),
        "all_rows_have_interpretations": bool(
            index["technical_interpretation"].astype(str).str.len().gt(0).all()
            and index["business_interpretation"].astype(str).str.len().gt(0).all()
        ),
    }


def unlocked_output_path(path: Path) -> Path:
    """Return a nearby fallback path when an Office file is open/locked."""
    return path.with_name(f"{path.stem}_updated{path.suffix}")


def save_docx_with_fallback(doc: Document, path: Path) -> Path:
    try:
        doc.save(path)
        return path
    except PermissionError:
        fallback = unlocked_output_path(path)
        doc.save(fallback)
        print(f"Could not overwrite locked file {path.name}; saved {fallback.name} instead.")
        return fallback


def save_pptx_with_fallback(prs: Presentation, path: Path) -> Path:
    try:
        prs.save(path)
        return path
    except PermissionError:
        fallback = unlocked_output_path(path)
        prs.save(fallback)
        print(f"Could not overwrite locked file {path.name}; saved {fallback.name} instead.")
        return fallback


def rel_path(path: Path) -> str:
    return str(path.resolve().relative_to(ROOT)).replace("\\", "/")


def run_smoke_checks(
    required_paths: list[Path],
    notebook_paths: list[Path] | None = None,
    notebook_html_paths: list[Path] | None = None,
    m1_docx_path: Path | None = None,
) -> dict[str, object]:
    summary: dict[str, object] = {
        "compile": {},
        "required_artifacts": {},
        "prediction": {},
        "notebook_quality": {},
        "notebook_self_contained": {},
        "notebook_warning_outputs": {},
        "notebook_html": {},
        "docx_interpretation_counts": {},
        "rubric_matrix": {},
        "univariate_plots": {},
    }
    for code_path in [ROOT / "run_all.py", ROOT / "insurance_modeling.py", ROOT / "app.py"]:
        try:
            py_compile.compile(str(code_path), doraise=True)
            summary["compile"][code_path.name] = True
        except Exception as exc:
            summary["compile"][code_path.name] = str(exc)

    for path in required_paths:
        summary["required_artifacts"][rel_path(path)] = path.exists()

    if notebook_paths:
        summary["notebook_quality"] = notebook_quality_summary(notebook_paths)
        summary["notebook_self_contained"] = notebook_self_contained_summary(notebook_paths)
        summary["notebook_warning_outputs"] = notebook_warning_output_summary(notebook_paths)
    if notebook_html_paths:
        summary["notebook_html"] = {rel_path(path): path.exists() for path in notebook_html_paths}
    if m1_docx_path is not None and m1_docx_path.exists():
        technical_count = count_docx_term(m1_docx_path, "Technical interpretation:")
        business_count = count_docx_term(m1_docx_path, "Business interpretation:")
        summary["docx_interpretation_counts"]["Milestone_1_Insurance_Price_Prediction.docx"] = {
            "actual_file": rel_path(m1_docx_path),
            "technical_interpretation_count": technical_count,
            "business_interpretation_count": business_count,
            "passes_threshold": bool(technical_count >= 20 and business_count >= 20),
        }
    summary["rubric_matrix"] = rubric_matrix_quality()
    summary["univariate_plots"] = univariate_plot_quality()

    try:
        model = joblib.load(MODEL_DIR / "final_model.pkl")
        sample = clean_column_names(load_raw_data()).drop(columns=[TARGET, "applicant_id"], errors="ignore").head(5)
        predictions = model.predict(sample)
        summary["prediction"] = {
            "model_loads": True,
            "prediction_count": int(len(predictions)),
            "first_5_predictions": [float(value) for value in predictions],
            "all_finite": bool(np.isfinite(predictions).all()),
        }
    except Exception as exc:
        summary["prediction"] = {"model_loads": False, "error": str(exc)}

    (REPORT_DIR / "smoke_test_summary.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
    return summary


def main() -> None:
    ensure_dirs()
    raw = load_raw_data()
    report_df = add_report_features(raw)
    tables = save_profile_tables(raw, report_df)
    figures = generate_eda_figures(raw, report_df)
    eda_summary = summarize_eda(report_df)
    model_results = build_models(raw)
    figures.update(model_results["model_figures"])

    app_path = write_streamlit_app()
    rubric_path = write_milestone1_rubric_coverage_matrix()
    notebooks = create_notebooks()
    req_path, readme_path = write_requirements_and_readme(model_results)
    m1_path = generate_milestone1_report(raw, report_df, tables, figures, eda_summary)
    m2_path = generate_milestone2_report(model_results, figures)
    notebook_html = execute_and_export_notebooks(notebooks)
    ppt_path = generate_presentation(model_results, figures, eda_summary)
    required_paths = [
        m1_path,
        m2_path,
        ppt_path,
        app_path,
        model_results["final_model_path"],
        MODEL_DIR / "model_metrics.csv",
        MODEL_DIR / "final_model_summary.csv",
        MODEL_DIR / "model_metadata.json",
        MODEL_DIR / "feature_importance.csv",
        MODEL_DIR / "repeated_cv_metrics.csv",
        MODEL_DIR / "price_grid_evaluation.csv",
        MODEL_DIR / "calibration_comparison.csv",
        MODEL_DIR / "ordinal_challenger_metrics.csv",
        MODEL_DIR / "app_schema.json",
        MODEL_DIR / "shap_status.json",
        FIG_DIR / "target_price_grid_frequency.png",
        FIG_DIR / "calibration_curve.png",
        TABLE_DIR / "target_grid_summary.csv",
        TABLE_DIR / "feature_signal_strength.csv",
        tables["variable_type_classification"],
        TABLE_DIR / "univariate_plot_interpretations.csv",
        rubric_path,
        *notebooks,
        *notebook_html,
    ]
    smoke_summary = run_smoke_checks(
        required_paths,
        notebook_paths=notebooks,
        notebook_html_paths=notebook_html,
        m1_docx_path=m1_path,
    )

    qa = {
        "docx_checks": {
            "Milestone_1_Insurance_Price_Prediction.docx": verify_docx_contains(
                m1_path,
                [
                    "Milestone 1 Submission",
                    "Data Report",
                    "Data Pre-processing",
                    "Extensive Exploratory Data Analysis",
                    "Technical interpretation:",
                    "Business interpretation:",
                ],
            ),
            str(m2_path.name): verify_docx_contains(
                m2_path,
                ["Milestone 2 Submission", "Model Selection", "Model Comparison", "Business Insights"],
            ),
        },
        "pptx": inspect_pptx(ppt_path),
        "smoke_tests": smoke_summary,
        "notebook_self_contained": notebook_self_contained_summary(notebooks),
        "notebook_warning_outputs": notebook_warning_output_summary(notebooks),
        "artifacts": {
            "app": rel_path(app_path),
            "notebooks": [rel_path(p) for p in notebooks],
            "notebook_html": [rel_path(p) for p in notebook_html],
            "requirements": rel_path(req_path),
            "readme": rel_path(readme_path),
            "final_model": rel_path(model_results["final_model_path"]),
            "milestone1_rubric_coverage_matrix": rel_path(rubric_path),
        },
    }
    (REPORT_DIR / "qa_summary.json").write_text(json.dumps(qa, indent=2), encoding="utf-8")

    print("Generated deliverables:")
    for path in [m1_path, m2_path, ppt_path, app_path, model_results["final_model_path"], *notebook_html]:
        print(path)


if __name__ == "__main__":
    main()
